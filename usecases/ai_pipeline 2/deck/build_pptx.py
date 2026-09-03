"""Generate the editable client presentation (.pptx) — neat layout + diagrams.

    python deck/build_pptx.py           # -> deck/AI_Pipeline_LLMOps.pptx

Design goals: no overlapping text/lines, real flow + architecture diagrams with
proper arrows, AFNI-style colours, no footer. Speaker notes (plain comments) on
every slide. Fully editable afterwards in PowerPoint / Teams / Google Slides.

Colours are defined once below — swap for AFNI's exact brand hex if needed.
"""
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ── palette (AFNI-style corporate blues) ─────────────────────────────────────
NAVY = RGBColor(0x0F, 0x2A, 0x47)
BLUE = RGBColor(0x2E, 0x6F, 0xB7)
BLUE_L = RGBColor(0xE7, 0xF0, 0xFA)
TEAL = RGBColor(0x1F, 0xA8, 0x9A)
TEAL_L = RGBColor(0xE6, 0xF6, 0xF3)
AMBER = RGBColor(0xD0, 0x8A, 0x2E)
INK = RGBColor(0x22, 0x30, 0x3F)
SOFT = RGBColor(0x5A, 0x6B, 0x7B)
LINE = RGBColor(0xD8, 0xE1, 0xEC)
LIGHT = RGBColor(0xF2, 0xF6, 0xFB)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
CLOUD = RGBColor(0xD7, 0xE3, 0xF1)

OUT = Path(__file__).resolve().parent / "AI_Pipeline_LLMOps.pptx"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
SW, SH = prs.slide_width, prs.slide_height
FONT = "Segoe UI"


def slide(bg=WHITE):
    s = prs.slides.add_slide(BLANK)
    r = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
    r.fill.solid(); r.fill.fore_color.rgb = bg; r.line.fill.background(); r.shadow.inherit = False
    return s


def _runs(tf, lines, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, wrap=True, space=6):
    tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    for m in ("margin_left", "margin_right"):
        setattr(tf, m, Pt(8))
    tf.margin_top = Pt(4); tf.margin_bottom = Pt(4)
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align; p.space_after = Pt(space); p.space_before = Pt(0)
        for (t, sz, b, c) in line:
            run = p.add_run(); run.text = t
            run.font.size = Pt(sz); run.font.bold = b; run.font.color.rgb = c; run.font.name = FONT


def textbox(s, x, y, w, h, lines, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, space=6):
    tb = s.shapes.add_textbox(x, y, w, h)
    _runs(tb.text_frame, lines, align, anchor, True, space)
    return tb


def box(s, x, y, w, h, fill, line=None, line_w=1.0, shape=MSO_SHAPE.ROUNDED_RECTANGLE,
        lines=None, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE):
    sh = s.shapes.add_shape(shape, x, y, w, h)
    sh.fill.solid(); sh.fill.fore_color.rgb = fill
    if line is not None:
        sh.line.color.rgb = line; sh.line.width = Pt(line_w)
    else:
        sh.line.fill.background()
    sh.shadow.inherit = False
    if lines:
        _runs(sh.text_frame, lines, align, anchor, True, 2)
    return sh


def header(s, kicker, title):
    box(s, 0, 0, SW, Inches(0.14), BLUE, shape=MSO_SHAPE.RECTANGLE)
    textbox(s, Inches(0.7), Inches(0.42), Inches(12), Inches(0.35),
            [[(kicker.upper(), 12, True, BLUE)]])
    textbox(s, Inches(0.7), Inches(0.78), Inches(12), Inches(0.7),
            [[(title, 28, True, NAVY)]])


def bullets(s, x, y, w, h, items, size=15, gap=10):
    lines = []
    for it in items:
        if isinstance(it, tuple):
            lead, rest = it
            lines.append([("—  ", size, True, TEAL), (lead, size, True, NAVY), (rest, size, False, INK)])
        else:
            lines.append([("—  ", size, True, TEAL), (it, size, False, INK)])
    textbox(s, x, y, w, h, lines, space=gap)


def notes(s, msg):
    s.notes_slide.notes_text_frame.text = msg


# ── 1 · TITLE ────────────────────────────────────────────────────────────────
s = slide(NAVY)
textbox(s, Inches(0.9), Inches(0.95), Inches(11.5), Inches(0.4),
        [[("AFNI  ·  OFFICE OF GenAI ARCHITECTURE", 13, True, TEAL)]])
textbox(s, Inches(0.9), Inches(1.55), Inches(11.6), Inches(1.9),
        [[("An Enterprise LLMOps Platform", 40, True, WHITE)],
         [("proven on the AI Pipeline (APIX)", 40, True, WHITE)]])
box(s, Inches(0.95), Inches(3.5), Inches(2.4), Inches(0.06), TEAL, shape=MSO_SHAPE.RECTANGLE)
textbox(s, Inches(0.9), Inches(3.8), Inches(11.2), Inches(1.1),
        [[("We didn't build one AI feature. We built the platform all our AI use cases run on —", 18, False, CLOUD)],
         [("and we already have a real application running on it.", 18, False, CLOUD)]])
box(s, Inches(0.9), Inches(5.35), Inches(11.5), Inches(0.75), RGBColor(0x14, 0x35, 0x59),
    lines=[[("14 reusable components      6 already used by APIX      one console to run & govern it",
             15, True, RGBColor(0xB6, 0xDF, 0xD7))]], align=PP_ALIGN.CENTER)
notes(s, "Open simple. Two things today: the platform we built, and proof it works because APIX already "
         "runs on it. The one idea to land is leverage — build once, reuse for every future use case.")

# ── 2 · THE IDEA ─────────────────────────────────────────────────────────────
s = slide()
header(s, "The idea", "Build the platform, not one-off features")
textbox(s, Inches(0.7), Inches(1.7), Inches(11.9), Inches(0.8),
        [[("If every AI use case is built from scratch, each one redoes the same plumbing — model calls, "
           "prompts, safety, cost tracking, evaluation — slowly and inconsistently.", 16, False, INK)]])
# two comparison cards + result
box(s, Inches(0.7), Inches(2.7), Inches(5.75), Inches(1.5), RGBColor(0xFB, 0xEE, 0xEE), line=RGBColor(0xE7, 0xC7, 0xC7),
    lines=[[("Without a platform", 15, True, RGBColor(0x9A, 0x3B, 0x3B))],
           [("Every use case re-invents model calls, safety, cost and evaluation. Slow, inconsistent, risky.",
             13.5, False, INK)]], align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE)
box(s, Inches(6.85), Inches(2.7), Inches(5.75), Inches(1.5), TEAL_L, line=RGBColor(0xB9, 0xE2, 0xDA),
    lines=[[("With the platform", 15, True, RGBColor(0x18, 0x6E, 0x63))],
           [("Every use case inherits guardrails, observability, evaluation and prompt versioning — the same way.",
             13.5, False, INK)]], align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE)
bullets(s, Inches(0.7), Inches(4.5), Inches(11.9), Inches(2),
        [("Build once. ", "The platform solves the hard, shared parts a single time."),
         ("Reuse everywhere. ", "A new use case plugs in with config, not a rebuild."),
         ("Governed by design. ", "Safety, cost and quality checks are built in, not bolted on."),
         ("Faster and cheaper. ", "Second and third use cases land in weeks, not quarters.")],
        size=16, gap=9)
notes(s, "The trap with GenAI is doing it feature by feature — each team rebuilds scaffolding and does safety "
         "and cost differently. We flipped it: one platform, use cases plug in. APIX is the first one on it.")


# ── component list slides ────────────────────────────────────────────────────
def component_slide(suffix, items):
    s = slide()
    header(s, "The platform", "What's in it — " + suffix)
    y = 1.72
    for num, name, desc in items:
        box(s, Inches(0.7), Inches(y), Inches(0.62), Inches(0.44), BLUE,
            lines=[[(num, 13, True, WHITE)]])
        textbox(s, Inches(1.5), Inches(y - 0.01), Inches(11.1), Inches(0.5),
                [[(name + "  —  ", 14.5, True, NAVY), (desc, 14.5, False, INK)]], anchor=MSO_ANCHOR.MIDDLE)
        y += 0.545
    return s


s = component_slide("part 1 of 2", [
    ("01", "Repo & foundation", "the shared project structure and cloud setup everything else builds on."),
    ("02", "Prompt management", "prompts are stored and versioned; you edit them without touching code."),
    ("03", "Model management", "pick the model by a simple name; swap models by config, per environment."),
    ("04", "Evaluation gate", "score prompts and models against known-good examples before anything ships."),
    ("05", "Observability", "every model call is traced — tokens, cost, time taken, and any safety flag."),
    ("06", "Guardrails", "checks each call for PII, secrets and prompt-injection, before and after the model."),
    ("07", "Data tools", "search over documents (RAG), speech to/from text, and API connectors."),
])
notes(s, "Walk these top to bottom, one line each — don't dwell. The ones that matter most for APIX are 02, "
         "03, 04, 05, 06. 01 and 07 are foundation and extras the other use cases will use.")

s = component_slide("part 2 of 2", [
    ("08", "Orchestration", "the engine that runs a pipeline's steps in order, with tracing built in."),
    ("09", "CI/CD", "automated checks on every change, including the evaluation gate as a required step."),
    ("10", "Serving / hosting", "expose a pipeline as a service or a container when we need to."),
    ("11", "Feedback", "capture human corrections and turn them into new test cases automatically."),
    ("12", "FinOps", "cost budgets and spend tracking so we stay in control of the bill."),
    ("13", "Governance & onboarding", "the step-by-step path to bring a new use case onto the platform."),
    ("14", "Security & compliance", "data-handling and privacy policy, mapped to industry standards."),
])
notes(s, "Second half, same pace. Reassure them these all exist and are reusable; APIX doesn't need every one "
         "yet, but the next use case picks up whatever it needs with no new platform work.")

# ── 5 · FLOW DIAGRAM: the pipeline ───────────────────────────────────────────
s = slide()
header(s, "The application", "What the AI Pipeline (APIX) does")
textbox(s, Inches(0.7), Inches(1.7), Inches(11.9), Inches(0.6),
        [[("It turns raw contact-centre call transcripts into per-agent coaching — in five steps:",
           16, False, INK)]])
flow = [("Denoise", "clean the\ntranscript"), ("Analysis", "score each\ncall"),
        ("Summary", "weekly\nreflection"), ("Metrics", "coaching\nmetrics"), ("KPI", "aggregate\nreport")]
n = len(flow); bw, gap = 2.05, 0.28
x = 0.7
for i, (t, d) in enumerate(flow):
    fill = BLUE if i < 4 else TEAL
    b = box(s, Inches(x), Inches(2.7), Inches(bw), Inches(1.15), fill,
            lines=[[(t, 16, True, WHITE)], [(d.replace("\n", " "), 12, False, RGBColor(0xE7, 0xF0, 0xFA))]])
    if i < n - 1:
        box(s, Inches(x + bw + 0.01), Inches(3.05), Inches(gap - 0.02), Inches(0.45), CLOUD,
            shape=MSO_SHAPE.RIGHT_ARROW)
    x += bw + gap
box(s, Inches(0.7), Inches(4.35), Inches(11.9), Inches(0.95), LIGHT, line=LINE,
    lines=[[("Every step makes its model calls through ", 15, False, INK),
            ("one shared function", 15, True, NAVY),
            (" — that single point is where the LLMOps platform plugs in (next slide).", 15, False, INK)]],
    align=PP_ALIGN.LEFT)
bullets(s, Inches(0.7), Inches(5.55), Inches(11.9), Inches(1.4),
        [("Only 'Analysis' and 'KPI' feed the scores; ", "the flow is linear and easy to reason about."),
         ("KPI is pure aggregation — ", "no model call, so it's fast and free.")], size=14.5, gap=8)
notes(s, "Keep it concrete: transcripts in, coaching out, five steps. The key fact for the next slide is that "
         "all five steps send their model calls through one function — that's our single integration point.")

# ── 6 · ARCHITECTURE DIAGRAM ─────────────────────────────────────────────────
s = slide()
header(s, "How it fits together", "The platform plugs in at one point")
# top: application container with the 5 steps
box(s, Inches(0.7), Inches(1.6), Inches(11.93), Inches(1.35), BLUE_L, line=RGBColor(0xC5, 0xDA, 0xF0))
textbox(s, Inches(0.9), Inches(1.66), Inches(6), Inches(0.3), [[("APPLICATION · AI PIPELINE (APIX)", 11, True, BLUE)]])
steps = ["Denoise", "Analysis", "Summary", "Metrics", "KPI"]
sbw, sgap = 1.9, 0.36; sx = 0.98
for i, t in enumerate(steps):
    box(s, Inches(sx), Inches(2.05), Inches(sbw), Inches(0.72), WHITE, line=RGBColor(0xC5, 0xDA, 0xF0),
        lines=[[(t, 13.5, True, NAVY)]])
    if i < 4:
        box(s, Inches(sx + sbw + 0.03), Inches(2.28), Inches(sgap - 0.06), Inches(0.28), RGBColor(0x9D, 0xBE, 0xE3),
            shape=MSO_SHAPE.RIGHT_ARROW)
    sx += sbw + sgap
# down arrow (each model call)
box(s, Inches(6.1), Inches(3.02), Inches(0.5), Inches(0.42), NAVY, shape=MSO_SHAPE.DOWN_ARROW)
textbox(s, Inches(6.75), Inches(3.05), Inches(4), Inches(0.35), [[("every model call goes down through…", 12, True, SOFT)]])
# seam band
box(s, Inches(0.7), Inches(3.55), Inches(11.93), Inches(0.95), NAVY,
    lines=[[("THE LLMOps SEAM  ", 13, True, TEAL),
            ("render prompt  ›  check guardrails  ›  pick model  ›  call  ›  check output  ›  record cost & trace",
             14, True, WHITE)]])
# up arrows from platform into the seam
for ax in (3.2, 6.35, 9.5):
    box(s, Inches(ax), Inches(4.55), Inches(0.5), Inches(0.4), TEAL, shape=MSO_SHAPE.UP_ARROW)
textbox(s, Inches(0.7), Inches(4.58), Inches(2.4), Inches(0.35), [[("…provided by ↑", 12, True, SOFT)]])
# platform container with component chips
box(s, Inches(0.7), Inches(5.05), Inches(11.93), Inches(1.75), TEAL_L, line=RGBColor(0xBF, 0xE0, 0xD9))
textbox(s, Inches(0.9), Inches(5.11), Inches(6), Inches(0.3), [[("LLMOps PLATFORM (reusable)", 11, True, RGBColor(0x18, 0x6E, 0x63))]])
chips = [("02", "Prompts"), ("03", "Models"), ("06", "Guardrails"), ("05", "Observability"), ("04", "Evaluation"), ("11", "Feedback")]
cbw, cgap = 1.78, 0.19; cx = 0.98
for num, name in chips:
    box(s, Inches(cx), Inches(5.5), Inches(cbw), Inches(1.05), WHITE, line=RGBColor(0xBF, 0xE0, 0xD9),
        lines=[[(num, 12, True, TEAL)], [(name, 13, True, NAVY)]])
    cx += cbw + cgap
notes(s, "This is the money slide. Top row is the application — the five pipeline steps. Every step sends its "
         "model calls down into the one 'seam' in the middle. That seam is where the platform does its work: "
         "the right prompt, guardrails, model choice, and cost/trace recording. The bottom row is the reusable "
         "platform providing those services. Point out we added the middle and bottom without changing the top.")

# ── 7 · IMPLEMENTATION APPROACH ──────────────────────────────────────────────
s = slide()
header(s, "How we built it", "A safe add-on, not a rewrite")
bullets(s, Inches(0.7), Inches(1.85), Inches(11.9), Inches(4),
        [("One integration point. ", "We wrapped the single shared function — the five steps were left untouched."),
         ("Low risk and reversible. ", "It's an add-on; we can switch any layer off and the pipeline still runs."),
         ("Shipped piece by piece. ", "Observability first, then models, guardrails, prompts, evaluation, feedback."),
         ("Config, not code. ", "APIX is registered with a few settings — model names, a guardrail policy, thresholds."),
         ("Fail-open. ", "If a cloud service is down, the pipeline keeps running; tracing just pauses.")],
        size=16, gap=13)
notes(s, "Plain version: the developer already funnelled every model call through one function; we wrapped that "
         "one spot. We added capabilities one at a time so nothing was big-bang. And it's reversible — that's why "
         "it was low risk.")

# ── 8 · WHAT'S WORKING ───────────────────────────────────────────────────────
s = slide()
header(s, "Where we are", "What's working today")
bullets(s, Inches(0.7), Inches(1.8), Inches(11.9), Inches(3.6),
        [("Done: ", "observability & cost, model management, guardrails, prompt versioning, evaluation gate, "
          "feedback — plus one console to run and operate it all."),
         ("Prompt change goes live by clicking Activate — ", "no redeployment."),
         ("PII flagged for audit, never dropped; secrets blocked. ", "Real customer data stays safe."),
         ("Runs on a locked-down machine — ", "the demo is plain Python, nothing to install."),
         ("Next: ", "send cost and traces to Azure Monitor once we have accounts and rates."),
         ("Ready when needed: ", "search/RAG, voice, and hosting for the next use case.")],
        size=15.5, gap=12)
notes(s, "The first four bullets are done and demonstrable. Stress two things: prompt changes go live without a "
         "deployment, and guardrails flag PII but never drop a call because these are real transcripts.")

# ── 9 · THE CONSOLE ──────────────────────────────────────────────────────────
s = slide()
header(s, "How we operate it", "One console for everything")
cards = [
    ("Application", "Run the pipeline; see KPIs, per-agent scores, reflections and the calls behind them."),
    ("Playground", "Edit a prompt, pick a model, score it against known-good examples; see the cost."),
    ("Evaluation", "Quality over time and per prompt version — the gate before anything ships."),
    ("Golden datasets", "View, edit, add or upload the known-good examples we score against."),
    ("Monitoring", "Calls, tokens, cost, time taken, and guardrail activity."),
    ("Feedback & Guardrails", "Coach feedback on calls, and an audit trail of every safety decision."),
]
cw, ch, cgap = 3.86, 1.5, 0.18
positions = [(0.7, 1.8), (0.7 + cw + cgap, 1.8), (0.7 + 2 * (cw + cgap), 1.8),
             (0.7, 1.8 + ch + cgap), (0.7 + cw + cgap, 1.8 + ch + cgap), (0.7 + 2 * (cw + cgap), 1.8 + ch + cgap)]
for (title, body), (px, py) in zip(cards, positions):
    b = box(s, Inches(px), Inches(py), Inches(cw), Inches(ch), LIGHT, line=LINE, anchor=MSO_ANCHOR.TOP,
            lines=[[(title, 15, True, NAVY)], [(body, 12.5, False, INK)]], align=PP_ALIGN.LEFT)
textbox(s, Inches(0.7), Inches(5.15), Inches(11.9), Inches(0.7),
        [[("Note: the demo uses a small sample of made-up call data so it runs anywhere. The same screens "
           "show real numbers once we connect live data.", 13, True, SOFT)]])
notes(s, "Describe the tabs simply — don't read every word. If we're live, click Application then Playground. Be "
         "upfront that the numbers are sample data so it runs on any machine; the same screens work with real data.")

# ── 10 · GOVERNANCE & COST ───────────────────────────────────────────────────
s = slide()
header(s, "Why it matters to leadership", "Safe and measured by default")
bullets(s, Inches(0.7), Inches(1.9), Inches(11.9), Inches(3.4),
        [("Safety on every call. ", "PII flagged, secrets blocked, injection checked. Change the policy in one "
          "place and it applies everywhere."),
         ("Nothing is a black box. ", "Every model call is traced with its cost and time."),
         ("Cost is a first-class number. ", "Track spend live; use a cheaper model for bulk steps to save."),
         ("Same rules for the next use case. ", "It inherits this safety and cost story automatically.")],
        size=16.5, gap=15)
notes(s, "This is the de-risking slide for a regulated contact-centre business. Nothing hits a model without "
         "guardrails, and everything is costed and traced. Use case two inherits the same safety — we don't redo "
         "governance every time.")

# ── 11 · DECISIONS ───────────────────────────────────────────────────────────
s = slide()
header(s, "What we need from you", "A few decisions to move forward")
bullets(s, Inches(0.7), Inches(1.85), Inches(11.9), Inches(3.6),
        [("Azure accounts and environments — ", "where we run, and dev/test/prod split."),
         ("Approved models and their real rates — ", "turns cost tracking into exact numbers."),
         ("PII policy — ", "flag for audit (today) or mask/redact; anything that must be hard-blocked."),
         ("Ground-truth examples — ", "who owns them, and a first small batch for APIX."),
         ("Human sign-off — ", "do coaches approve the AI's notes before agents see them?"),
         ("The next use case — ", "which one to onboard to prove reuse (voice, performance index, hiring).")],
        size=15.5, gap=12)
box(s, Inches(0.7), Inches(6.15), Inches(11.9), Inches(0.7), TEAL_L, line=RGBColor(0xB9, 0xE2, 0xDA),
    lines=[[("If we get the first four, we can take APIX live and start the next use case at the same time.",
             14.5, True, RGBColor(0x18, 0x6E, 0x63))]], align=PP_ALIGN.LEFT)
notes(s, "These aren't blockers to the demo — they turn it into production and set the roadmap. If I can get four "
         "things: approved models and rates, the PII stance, a small set of ground-truth examples, and which "
         "second use case you want.")

# ── 12 · CLOSE (reuse flow) ──────────────────────────────────────────────────
s = slide(NAVY)
textbox(s, Inches(0.9), Inches(0.85), Inches(11.5), Inches(0.9), [[("Build once. Onboard many.", 38, True, WHITE)]])
textbox(s, Inches(0.9), Inches(1.85), Inches(11.5), Inches(0.9),
        [[("The next use case reuses the same guardrails, cost controls and console — and lands in weeks.",
           17, False, CLOUD)]])
# onboarding flow with arrows
steps = ["Intake", "Configure", "Evaluate", "Operate", "Improve"]
bw2, gap2 = 2.05, 0.28; x = 0.9
for i, t in enumerate(steps):
    box(s, Inches(x), Inches(3.0), Inches(bw2), Inches(0.8), RGBColor(0x14, 0x35, 0x59), line=TEAL,
        lines=[[(t, 14.5, True, WHITE)]])
    if i < len(steps) - 1:
        box(s, Inches(x + bw2 + 0.01), Inches(3.24), Inches(gap2 - 0.02), Inches(0.32), TEAL,
            shape=MSO_SHAPE.RIGHT_ARROW)
    x += bw2 + gap2
textbox(s, Inches(0.9), Inches(4.1), Inches(11.5), Inches(0.4),
        [[("Only the config and the known-good examples are new each time — everything else is inherited.",
           13.5, False, RGBColor(0x9F, 0xB6, 0xCF))]])
# three next use cases
nxt = [("Voice Agent", "adds voice; reuses the rest."),
       ("Performance Index", "reuses the exact APIX analysis approach."),
       ("Hiring Intelligence", "reuses guardrails and human sign-off.")]
cw3 = 3.7; x = 0.9
for (t, d) in nxt:
    box(s, Inches(x), Inches(4.75), Inches(cw3), Inches(1.15), RGBColor(0x14, 0x35, 0x59), line=RGBColor(0x2B, 0x4B, 0x6E),
        lines=[[(t, 15, True, RGBColor(0x8F, 0xD6, 0xC9))], [(d, 13, False, CLOUD)]], align=PP_ALIGN.LEFT)
    x += cw3 + 0.28
textbox(s, Inches(0.9), Inches(6.2), Inches(11.5), Inches(0.6),
        [[("The ask: approve the decisions, and let us take APIX live and start use case #2.",
           16, True, RGBColor(0xB6, 0xDF, 0xD7))]])
notes(s, "Close on leverage: you funded a platform, not a feature. The next use case reuses all of this and lands "
         "fast with the same governance. Ask for the decisions and the green light to take APIX live and start #2.")

prs.save(str(OUT))
print(f"Wrote {OUT}  ({len(prs.slides._sldIdLst)} slides)")
