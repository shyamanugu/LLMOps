#!/usr/bin/env python3
"""v2 implementation deck (enhanced) — how we actually implement LLMOps (Azure + GitHub).
Fully editable (native shapes + real code panels, no images). Speaker notes. No timelines.
Costs are indicative — confirm at a sizing exercise.
Output: presentation/LLMOps-Implementation-v2.pptx
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR

NAVY=RGBColor(0x1F,0x3A,0x5F); BLUE=RGBColor(0x2F,0x5C,0x9E); TEAL=RGBColor(0x2A,0x9D,0x8F)
GRAYTX=RGBColor(0x3C,0x46,0x54); LIGHT=RGBColor(0xEE,0xF1,0xF5); MUTE=RGBColor(0x6B,0x74,0x80)
WHITE=RGBColor(0xFF,0xFF,0xFF); LINEC=RGBColor(0xC9,0xD2,0xDD); AMBER=RGBColor(0xC8,0x7B,0x1E)
SKY=RGBColor(0x3E,0x6F,0xB0); SLATE=RGBColor(0x51,0x6B,0x8A); DTEAL=RGBColor(0x2E,0x7D,0x7D)
PURP=RGBColor(0x6A,0x51,0x9E); GREEN=RGBColor(0x3E,0x8E,0x5A); ROSE=RGBColor(0xB0,0x4A,0x5A)
CODEBG=RGBColor(0x1E,0x28,0x36); CODEFG=RGBColor(0xDD,0xE4,0xEC); CODEKEY=RGBColor(0x7F,0xD0,0xC8)
FONT="Calibri"; MONO="Consolas"; SW,SH=Inches(13.333),Inches(7.5)
prs=Presentation(); prs.slide_width=SW; prs.slide_height=SH; BLANK=prs.slide_layouts[6]

def add():
    s=prs.slides.add_slide(BLANK)
    r=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,0,0,SW,SH); r.fill.solid(); r.fill.fore_color.rgb=WHITE; r.line.fill.background(); r.shadow.inherit=False
    return s
def _txt(s,x,y,w,h,runs,align=PP_ALIGN.LEFT,anchor=MSO_ANCHOR.TOP,sa=6,ls=1.05):
    tf=s.shapes.add_textbox(x,y,w,h).text_frame; tf.word_wrap=True; tf.vertical_anchor=anchor
    tf.margin_left=0; tf.margin_right=0; tf.margin_top=0; tf.margin_bottom=0
    for i,para in enumerate(runs):
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph()
        p.alignment=align; p.space_after=Pt(sa); p.space_before=Pt(0); p.line_spacing=ls
        for (t,sz,b,c,*rest) in para:
            it=rest[0] if rest else False
            r=p.add_run(); r.text=t; r.font.size=Pt(sz); r.font.bold=b; r.font.italic=it; r.font.name=FONT; r.font.color.rgb=c
    return tf
def fillrect(s,x,y,w,h,color):
    sp=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,x,y,w,h); sp.fill.solid(); sp.fill.fore_color.rgb=color
    sp.line.fill.background(); sp.shadow.inherit=False; return sp
def rrect(s,x,y,w,h,fill,line=LINEC,lw=1.0):
    sp=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,x,y,w,h); sp.fill.solid(); sp.fill.fore_color.rgb=fill
    if line is None: sp.line.fill.background()
    else: sp.line.color.rgb=line; sp.line.width=Pt(lw)
    sp.shadow.inherit=False
    try: sp.adjustments[0]=0.06
    except Exception: pass
    return sp
def settext(sp,lines,anchor=MSO_ANCHOR.MIDDLE,align=PP_ALIGN.CENTER):
    tf=sp.text_frame; tf.word_wrap=True; tf.vertical_anchor=anchor
    tf.margin_left=Inches(0.08); tf.margin_right=Inches(0.08); tf.margin_top=Inches(0.03); tf.margin_bottom=Inches(0.03)
    for i,(t,sz,b,c) in enumerate(lines):
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph()
        p.alignment=align; p.space_after=Pt(1); p.line_spacing=1.0
        r=p.add_run(); r.text=t; r.font.size=Pt(sz); r.font.bold=b; r.font.name=FONT; r.font.color.rgb=c
def box(s,x,y,w,h,head,sub=None,fill=LIGHT,headc=NAVY,hs=12.5,ss=10,subc=GRAYTX,align=PP_ALIGN.CENTER):
    sp=rrect(s,x,y,w,h,fill); lines=[(head,hs,True,headc)]
    if sub: lines.append((sub,ss,False,subc))
    settext(sp,lines,align=align); return sp
def arrow(s,x,y,w=Inches(0.34),h=Inches(0.3),color=TEAL,shape=MSO_SHAPE.RIGHT_ARROW):
    sp=s.shapes.add_shape(shape,x,y,w,h); sp.fill.solid(); sp.fill.fore_color.rgb=color; sp.line.fill.background(); sp.shadow.inherit=False; return sp
def connect(s,x1,y1,x2,y2,color=LINEC,w=1.4):
    ln=s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,x1,y1,x2,y2); ln.line.color.rgb=color; ln.line.width=Pt(w); return ln
def title(s,text,num,kicker=None):
    if kicker: _txt(s,Inches(0.7),Inches(0.32),Inches(11.9),Inches(0.3),[[(kicker.upper(),11,True,TEAL)]])
    _txt(s,Inches(0.7),Inches(0.56),Inches(12.0),Inches(0.7),[[(text,22,True,NAVY)]])
    fillrect(s,Inches(0.72),Inches(1.14),Inches(1.5),Pt(3),TEAL)
    _txt(s,SW-Inches(1.1),SH-Inches(0.42),Inches(0.7),Inches(0.3),[[(str(num),10,False,MUTE)]],align=PP_ALIGN.RIGHT)
    _txt(s,Inches(0.7),SH-Inches(0.42),Inches(9),Inches(0.3),[[("LLMOps implementation (v2) — Azure + GitHub",9,False,MUTE)]])
def bullets(s,items,x=Inches(0.75),y=Inches(1.5),w=Inches(11.8),h=Inches(5.2),size=14,gap=8,color=GRAYTX):
    tf=s.shapes.add_textbox(x,y,w,h).text_frame; tf.word_wrap=True; tf.margin_left=0; tf.margin_top=0
    for i,it in enumerate(items):
        text=it[0]; lvl=it[1] if len(it)>1 else 0; bold=it[2] if len(it)>2 else False; col=it[3] if len(it)>3 else color
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph()
        p.space_after=Pt(gap); p.space_before=Pt(0); p.line_spacing=1.04; p.level=lvl
        r=p.add_run(); r.text=("•  " if lvl==0 else "–  ")+text
        r.font.size=Pt(size-lvl); r.font.bold=bold; r.font.name=FONT; r.font.color.rgb=col
    return tf
def codepanel(s,x,y,w,h,lines,titlebar=None):
    if titlebar:
        rrect(s,x,y,w,Inches(0.32),RGBColor(0x2C,0x3A,0x4C),line=None)
        _txt(s,x+Inches(0.15),y,w-Inches(0.3),Inches(0.32),[[(titlebar,10,True,CODEKEY)]],anchor=MSO_ANCHOR.MIDDLE)
        y=y+Inches(0.32); h=h-Inches(0.32)
    rrect(s,x,y,w,h,CODEBG,line=None)
    tf=s.shapes.add_textbox(x+Inches(0.14),y+Inches(0.08),w-Inches(0.28),h-Inches(0.16)).text_frame
    tf.word_wrap=True; tf.margin_left=0; tf.margin_top=0
    for i,ln in enumerate(lines):
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph(); p.space_after=Pt(0); p.line_spacing=1.02
        r=p.add_run(); r.text=ln if ln else " "; r.font.name=MONO; r.font.size=Pt(9); r.font.color.rgb=CODEFG
def delta(s,y,today,ours,change,x=Inches(0.6),w=Inches(12.13),h=Inches(1.05)):
    cw=(w.inches-0.2*2)/3
    cols=[("TODAY",today,ROSE,RGBColor(0xFB,0xEE,0xEF)),("OUR SETUP",ours,TEAL,RGBColor(0xEA,0xF5,0xF2)),
          ("WHAT CHANGES",change,NAVY,RGBColor(0xEC,0xF0,0xF7))]
    xx=x
    for lbl,txt,hc,bg in cols:
        rrect(s,xx,y,Inches(cw),h,bg,line=None)
        _txt(s,xx+Inches(0.15),y+Inches(0.08),Inches(cw-0.3),Inches(0.3),[[(lbl,10,True,hc)]])
        _txt(s,xx+Inches(0.15),y+Inches(0.38),Inches(cw-0.3),h-Inches(0.46),[[(txt,10.5,False,GRAYTX)]])
        xx=xx+Inches(cw+0.2)
def table(s,x,y,w,headers,rows,cw,fs=10.5,hs=10.5,rh=Inches(0.5)):
    nr=len(rows)+1; t=s.shapes.add_table(nr,len(headers),x,y,w,rh*nr).table
    for i,cc in enumerate(cw): t.columns[i].width=cc
    for j,hh in enumerate(headers):
        c=t.cell(0,j); c.fill.solid(); c.fill.fore_color.rgb=NAVY; c.vertical_anchor=MSO_ANCHOR.MIDDLE
        c.margin_left=Inches(0.07); c.margin_top=Inches(0.02); c.margin_bottom=Inches(0.02)
        r=c.text_frame.paragraphs[0].add_run(); r.text=hh; r.font.size=Pt(hs); r.font.bold=True; r.font.color.rgb=WHITE; r.font.name=FONT
    for i,row in enumerate(rows):
        for j,v in enumerate(row):
            c=t.cell(i+1,j); c.fill.solid(); c.fill.fore_color.rgb=WHITE if i%2==0 else LIGHT; c.vertical_anchor=MSO_ANCHOR.MIDDLE
            c.margin_left=Inches(0.07); c.margin_top=Inches(0.02); c.margin_bottom=Inches(0.02)
            r=c.text_frame.paragraphs[0].add_run(); r.text=v; r.font.size=Pt(fs); r.font.name=FONT; r.font.color.rgb=GRAYTX
            if j==0: r.font.bold=True; r.font.color.rgb=NAVY
    return t
def notes(s,t): s.notes_slide.notes_text_frame.text=t
print("Building enhanced v2 deck...")

# 1 TITLE
s=add(); fillrect(s,0,0,Inches(0.28),SH,NAVY)
_txt(s,Inches(0.9),Inches(1.95),Inches(11.5),Inches(1.2),[[("LLMOps: how we implement it",30,True,NAVY)]])
fillrect(s,Inches(0.94),Inches(2.78),Inches(2.0),Pt(3),TEAL)
_txt(s,Inches(0.92),Inches(3.08),Inches(11.3),Inches(1.6),
     [[("Component by component — the actual setup on Azure + GitHub, what changes from how we work today, and what it costs.",16,False,GRAYTX)],
      [("Grounded in APIX and Hiring Intelligence; reusable for any use case. Costs are indicative (confirm at sizing). No timelines.",12.5,False,MUTE,True)]],sa=10)
_txt(s,Inches(0.92),Inches(6.35),Inches(11),Inches(0.5),[[("Working draft for review",12,True,NAVY)]])
notes(s,"This version adds the depth from the feedback: reusability made explicit, real repo structure, cost numbers, evaluation thresholds, a guardrails list, and clearer terminology. Every component slide keeps the Today / Our setup / What changes strip.")

# 2 COMPONENTS + REUSABILITY
s=add(); title(s,"The components — and which are shared vs per use case",2,"Reusability")
table(s,Inches(0.6),Inches(1.4),Inches(6.15),
      ["Shared platform (build once)",""],
      [["Source control & CI/CD","backbone"],["Prompt registry & management",""],["Model catalog & routing",""],
       ["Evaluation engine & gate",""],["Observability & tracing",""],["Guardrails engine",""],
       ["Data-access & RAG framework",""],["Reusable tool catalog (MCP)",""],["Orchestration / pipeline runtime",""],
       ["Serving & gateway",""],["Identity & secrets, FinOps",""],["Feedback capture & analytics",""]],
      cw=[Inches(4.4),Inches(1.75)],fs=10.5,rh=Inches(0.35))
table(s,Inches(7.0),Inches(1.4),Inches(5.73),
      ["Per use case (new each time)",""],
      [["Prompt content",""],["Agent / pipeline design",""],["Golden dataset + thresholds",""],
       ["Data sources & connectors",""],["Use-case-specific tools",""],["Guardrail policy tuning",""],
       ["Dashboards & alerts",""]],
      cw=[Inches(4.0),Inches(1.73)],fs=10.5,rh=Inches(0.42))
_txt(s,Inches(0.6),Inches(6.35),Inches(12.13),Inches(0.5),
     [[("Most of the machinery is shared — a new use case reuses all 12 left-hand components and only adds the right-hand items.",12.5,True,NAVY,True)]])
notes(s,"Point 1: the fuller list, with reusability explicit. Twelve shared components built once; a new use case only adds the right-hand column. This is why the platform pays off after the first use case.")

# 3 REPO STRUCTURE
s=add(); title(s,"Repository structure — shared vs per use case, scales to N",3,"Repository")
codepanel(s,Inches(0.6),Inches(1.45),Inches(7.3),Inches(4.9),
    ["llmops-platform/","├─ platform/                 # SHARED — built once, reused",
     "│  ├─ common/    prompt_loader  model_router","│  │              tracing  guardrails  data_access",
     "│  ├─ tools/     search_knowledge/ query_sql/","│  │              extract_document/ get_record/   # MCP",
     "│  ├─ evaluators/ ragas  deepeval  tool_selection","│  ├─ gateway/   apim-policies/","│  └─ infra/     bicep modules",
     "├─ usecases/","│  ├─ apix/","│  │   ├─ prompts/  *.prompt.yaml","│  │   ├─ agents/   pipeline.agent.yaml",
     "│  │   ├─ evals/    golden.*.jsonl  evaluators.yaml","│  │   ├─ tools/    (use-case-specific only)","│  │   └─ config/   datasources.yaml",
     "│  └─ hiring/       (same shape)","├─ models.yaml               # shared alias -> model","└─ .github/workflows/        # shared pipelines"],titlebar="one monorepo")
_txt(s,Inches(8.05),Inches(1.5),Inches(4.6),Inches(0.5),[[("How it scales:",13,True,NAVY)]])
bullets(s,[("platform/ = the shared, reusable code every use case runs on.",0),
    ("usecases/<name>/ = one folder per use case, always the same shape.",0),
    ("the Nth use case just adds a folder under usecases/ and reuses everything in platform/.",0,True,NAVY),
    ("prompts, agents, evals, tools, config all have a clear, fixed home.",0)],
    x=Inches(8.05),y=Inches(2.0),w=Inches(4.7),h=Inches(4.0),size=12,gap=10)
notes(s,"Point 2: exactly where everything sits. platform/ is shared; usecases/<name>/ is per use case with a fixed shape. Adding a use case = adding a folder. Shared components are clearly visible at the top.")

# 4 CI/CD FLOW
s=add(); title(s,"CI/CD: the stages a change goes through",4,"Delivery")
stages=[("Author change","edit a prompt / agent / model config"),("Pull request","peer + CODEOWNERS review"),
        ("Automated checks","lint, unit, contract tests"),("Evaluation gate","golden-set metrics vs baseline — blocks on regression"),
        ("Merge","into main"),("Promotion gates","dev→test→prod: each needs approver + eval-full pass"),
        ("Canary release","new version to a small % of traffic"),("Full rollout / rollback","ramp to 100%, or auto-revert")]
table(s,Inches(0.6),Inches(1.5),Inches(12.13),["Stage","What happens"],
      [[a,b] for a,b in stages],cw=[Inches(2.7),Inches(9.43)],fs=11,rh=Inches(0.52))
_txt(s,Inches(0.6),Inches(6.3),Inches(12.13),Inches(0.5),
     [[("The Evaluation gate is what makes it LLMOps, not just DevOps: ",12.5,True,NAVY),("no change ships unless it passes the golden-set thresholds.",12.5,False,GRAYTX)]])
notes(s,"Point 3: clearer terminology and a description per stage. Walk down the table. The evaluation gate (row 4) is the difference from ordinary DevOps.")

# 5 PROMPT MGMT DIFFERENTIATOR
s=add(); title(s,"Prompt management — what's actually different",5,"Component: Prompts")
rrect(s,Inches(0.6),Inches(1.42),Inches(12.13),Inches(0.68),RGBColor(0xFB,0xEE,0xEF),line=None)
_txt(s,Inches(0.8),Inches(1.5),Inches(11.7),Inches(0.55),
     [[("The prompts are already in Git, inside the code files. So \"store prompts in Git\" is NOT the change. ",12.5,True,ROSE),
       ("Three things are:",12.5,True,NAVY)]],anchor=MSO_ANCHOR.MIDDLE)
bullets(s,[("One YAML file PER prompt — id, version, template, variables, eval_refs, changelog. (No YAML-per-prompt today.)",0,True,NAVY),
    ("Every prompt change runs the pipeline and must PASS the golden-dataset thresholds before it deploys.",0,True,NAVY),
    ("A registry that holds prompts so we can roll back / swap / compare versions on their evaluation scores.",0,True,NAVY)],
    x=Inches(0.75),y=Inches(2.3),w=Inches(6.0),h=Inches(2.6),size=12,gap=10)
codepanel(s,Inches(6.95),Inches(2.25),Inches(5.75),Inches(4.0),
    ["# prompts/apix/coaching-report.prompt.yaml","id: apix.coaching_report","version: 3","labels: [prod]",
     "model_alias: reason        # via models.yaml","temperature: 0.2","inputs: [agent_name, program, scores, evidence]",
     "template: |","  Using ONLY the evidence below, write a","  coaching note. Cite evidence. Do not invent.","eval_refs: [evals/apix/golden.telesales.jsonl]",
     "changelog:","  - v3: require evidence citation"],titlebar="one YAML per prompt")
notes(s,"The key slide from last time. Say it plainly: prompts already in Git is not new; the YAML-per-prompt artifact, the evaluation gate, and the registry for rollback/compare are.")

# 6 REGISTRY — HOW EACH WORKS
s=add(); title(s,"Where prompts are held: how the three registries work",6,"Component: Prompts")
rrect(s,Inches(0.6),Inches(1.5),Inches(3.9),Inches(4.5),LIGHT)
_txt(s,Inches(0.6),Inches(1.65),Inches(3.9),Inches(0.4),[[("Git + in-app cache",13.5,True,NAVY)]],align=PP_ALIGN.CENTER)
bullets(s,[("prompts live in usecases/<uc>/prompts/*.yaml",0),("app reads at startup, caches in memory",0),
    ("version = git + the YAML version field",0),("rollback = git revert / move the label",0),("cost: $0 (in the repo)",0,True,TEAL)],
    x=Inches(0.82),y=Inches(2.3),w=Inches(3.5),h=Inches(3.5),size=11,gap=9)
rrect(s,Inches(4.7),Inches(1.5),Inches(3.9),Inches(4.5),RGBColor(0xEA,0xF5,0xF2),line=RGBColor(0xBF,0xD6,0xD2))
_txt(s,Inches(4.7),Inches(1.65),Inches(3.9),Inches(0.4),[[("Langfuse (self-hosted)",13.5,True,TEAL)]],align=PP_ALIGN.CENTER)
bullets(s,[("open-source product; we host it in our network",0),("UI to edit, compare, label (prod/staging), roll back",0),
    ("app fetches by id + label at runtime",0),("also gives token/cost dashboards + tracing",0),("CI syncs the Git YAML into it on merge",0),
    ("cost: MIT-free + infra ≈ $50–150/mo",0,True,TEAL)],
    x=Inches(4.92),y=Inches(2.3),w=Inches(3.5),h=Inches(3.5),size=11,gap=8)
rrect(s,Inches(8.8),Inches(1.5),Inches(3.9),Inches(4.5),LIGHT)
_txt(s,Inches(8.8),Inches(1.65),Inches(3.9),Inches(0.4),[[("Foundry prompt assets",13.5,True,NAVY)]],align=PP_ALIGN.CENTER)
bullets(s,[("versioned assets inside an Azure AI Foundry project",0),("accessed via the SDK",0),
    ("integrates with Foundry evaluations + tracing",0),("fully managed by Azure",0),("cost: folded into Azure usage (minor)",0,True,TEAL)],
    x=Inches(9.02),y=Inches(2.3),w=Inches(3.5),h=Inches(3.5),size=11,gap=9)
notes(s,"Point 4: how each works and where prompts live. Langfuse is a self-hosted open-source product with a UI and dashboards; Foundry keeps everything managed in Azure; Git is simplest. Next slide compares and recommends.")

# 7 REGISTRY — COMPARISON + DECISION
s=add(); title(s,"Registry comparison, and what we recommend",7,"Component: Prompts")
table(s,Inches(0.6),Inches(1.5),Inches(12.13),
      ["Dimension","Git + in-app cache","Langfuse (self-host)","Foundry prompt assets"],
      [["Compare / rollback UI","git only","yes, visual","yes"],
       ["Extra features","none","observability + cost dashboards","evals + tracing (Azure)"],
       ["Ops burden","none","we run a container + DB","none (managed)"],
       ["Data residency","our repo","our network","Azure tenant"],
       ["Cost / month","$0","≈ $50–150 (infra)","minor (Azure usage)"]],
      cw=[Inches(2.4),Inches(2.9),Inches(3.6),Inches(3.23)],fs=10,rh=Inches(0.5))
rrect(s,Inches(0.6),Inches(4.95),Inches(12.13),Inches(1.35),RGBColor(0xEC,0xF0,0xF7),line=None)
_txt(s,Inches(0.8),Inches(5.05),Inches(11.7),Inches(1.2),
     [[("Recommendation:  ",13,True,NAVY),("start with Git + in-app cache (zero cost, simplest). Add Langfuse when we want the "
       "visual compare/rollback and built-in cost dashboards. Use Foundry if we prefer everything fully managed in Azure. "
       "All three can coexist — Git is the source of truth, the registry is the runtime copy.",12,False,GRAYTX)]])
notes(s,"Point 4 continued: a real comparison to decide. Decisive recommendation, with cost. Git first; Langfuse when the UI/observability is worth ~$50-150/mo; Foundry if fully managed.")

# 8 MODEL MGMT (code vs pipeline)
s=add(); title(s,"Model management — config-as-code (code and pipeline)",8,"Component: Models")
_txt(s,Inches(0.6),Inches(1.4),Inches(12.13),Inches(0.5),
     [[("Is it code-level or pipeline-level? Both — and neither is a manual portal setting. It is config-as-code:",12.5,True,NAVY,True)]])
box(s,Inches(0.7),Inches(2.0),Inches(3.5),Inches(1.15),"1. Repo config","models.yaml changed via pull request (code-level)",fill=LIGHT,hs=12.5,ss=10)
arrow(s,Inches(4.25),Inches(2.5))
box(s,Inches(4.7),Inches(2.0),Inches(3.5),Inches(1.15),"2. CI eval gate","pipeline validates + compares vs current (devops-level)",fill=RGBColor(0xF6,0xE9,0xE0),headc=AMBER,hs=12.5,ss=10)
arrow(s,Inches(8.25),Inches(2.5))
box(s,Inches(8.7),Inches(2.0),Inches(4.0),Inches(1.15),"3. Runtime resolver","app resolves alias→deployment per env (APP_ENV)",fill=LIGHT,hs=12.5,ss=10)
codepanel(s,Inches(0.6),Inches(3.45),Inches(7.4),Inches(2.5),
    ["# models.yaml (in the repo, reviewed via PR)","environments:","  prod: { aliases: { reason: gpt-5.2, bulk: gpt-5-mini } }","  dev:  { aliases: { reason: gpt-5-mini, bulk: gpt-5-mini } }","",
     "# app code:  resolve('reason')  -> deployment for APP_ENV","# never 'gpt-5.2' hard-coded; a swap = a reviewed, gated PR"],titlebar="models.yaml + resolver")
bullets(s,[("app asks for a task alias (reason / bulk / voice); the model is resolved from config.",0),
    ("a model swap is a config change that must pass the evaluation gate.",0),
    ("one shared config reused by every agent and use case under one hub.",0)],
    x=Inches(8.2),y=Inches(3.55),w=Inches(4.5),h=Inches(2.4),size=11.5,gap=10)
notes(s,"Point 5: resolve the confusion. It's config-as-code — a file in the repo (code), changed by PR and gated by the pipeline (devops), read by the app at runtime per environment. Not hard-coded, not a portal click.")

# 9 GOLDEN DATASET
s=add(); title(s,"Evaluation: the golden dataset is the gate",9,"Component: Evaluation")
bullets(s,[("A golden dataset is the ground-truth set of test cases for a use case — the first thing we build.",0),
    ("Same idea as normal ground truth, but in LLMOps it runs as a GATE at every change / pipeline run.",0,True,NAVY),
    ("Sources (three-step): SME-authored first → real traffic over time → reviewed again by SMEs & business.",0),
    ("Per use case AND per program (APIX: Telesales and WCC score differently). Start ~50–200 cases.",0)],
    x=Inches(0.7),y=Inches(1.5),w=Inches(6.0),h=Inches(3.2),size=12.5,gap=11)
codepanel(s,Inches(6.95),Inches(1.5),Inches(5.75),Inches(2.9),
    ["// evals/apix/golden.telesales.jsonl  (one line)","{","  \"id\": \"apix-telesales-014\",","  \"input\": {\"transcript_id\":\"c-88421\",",
     "            \"program\":\"telesales\"},","  \"grading\": {\"must_cite_evidence\": true,","    \"expected_score_band\": [70,85],","    \"must_flag\": [\"missed_upsell\"]},",
     "  \"meta\": {\"source\": \"sme_authored\"}","}"],titlebar="golden dataset record")
delta(s,Inches(4.9),"quality checked manually / by spot-checking a few calls","a versioned golden dataset run automatically on every change","evaluation becomes a release gate, not a manual look",h=Inches(1.05))
notes(s,"Kiran asked how golden data differs from normal ground truth — it IS ground truth, but run as a gate on every change, and it grows from real traffic.")

# 10 METRIC MEANINGS + THRESHOLDS
s=add(); title(s,"Evaluation: what each metric means, and thresholds",10,"Component: Evaluation")
table(s,Inches(0.6),Inches(1.42),Inches(7.4),
      ["Metric","What it measures"],
      [["Groundedness","every claim is backed by the retrieved source"],
       ["Context relevance","the retrieved chunks are on-point"],
       ["Answer relevance","the answer addresses the question"],
       ["Coherence / fluency","it reads clearly"],
       ["Correctness","matches the reference answer"],
       ["Tool-selection acc.","% of times the right tool was chosen"],
       ["Task success","% end-to-end correct"]],
      cw=[Inches(2.4),Inches(5.0)],fs=10.5,rh=Inches(0.48))
_txt(s,Inches(8.2),Inches(1.42),Inches(4.5),Inches(0.4),[[("How we set thresholds:",13,True,NAVY)]])
bullets(s,[("run a baseline on current production.",0),
    ("gate rule: no metric may drop more than X% below baseline.",0),
    ("absolute floors for safety: PII leak = 0, unsafe = 0.",0),
    ("minimums for critical metrics: e.g. groundedness ≥ 0.90.",0),
    ("thresholds live in evaluators.yaml, enforced by the gate.",0,True,TEAL)],
    x=Inches(8.2),y=Inches(1.9),w=Inches(4.5),h=Inches(3.2),size=11.5,gap=10)
codepanel(s,Inches(8.2),Inches(4.5),Inches(4.5),Inches(1.75),
    ["# evaluators.yaml","groundedness: {min: 0.90, regression: 0.02}","tool_selection: {min: 0.95}","pii_leak: {max: 0.0}"],titlebar="thresholds")
notes(s,"Point 6a/6b: what each metric means in plain words, and how thresholds are set — baseline + regression tolerance + absolute floors. Show evaluators.yaml.")

# 11 EVAL TECHNIQUE COST
s=add(); title(s,"Evaluation: the techniques, and what they cost",11,"Component: Evaluation")
table(s,Inches(0.6),Inches(1.5),Inches(12.13),
      ["Technique","Covers","How it's charged","Indicative cost"],
      [["Custom Python","tool selection, exact match, task path","compute only (CI minutes)","~free"],
       ["Ragas","RAG metrics (groundedness, relevance)","free code; LLM-metrics call a judge","~free + judge tokens"],
       ["DeepEval","writing quality, custom (G-Eval), CI","free code; LLM-metrics call a judge","~free + judge tokens"],
       ["LLM-as-judge","subjective quality with a rubric","judge model tokens per case","cents–low $ per 200-case run"],
       ["Azure Foundry evals","built-in + custom, trace-linked","Azure judge tokens","usage only, no license"],
       ["LangSmith","eval + observability platform","per-seat + per-trace license","≈ $1,500–2,800 / mo at scale"]],
      cw=[Inches(2.3),Inches(3.6),Inches(3.6),Inches(2.63)],fs=10,rh=Inches(0.52))
_txt(s,Inches(0.6),Inches(5.5),Inches(12.13),Inches(0.7),
     [[("Driver: ",12.5,True,NAVY),("judge tokens × dataset size × runs. Keep it cheap — use a small judge model (e.g. GPT-5-mini), run a "
       "subset on each pull request and the full set nightly. Open-source tools + a small judge cover most needs; LangSmith is the pricey option.",12,False,GRAYTX)]])
notes(s,"Point 6c: cost of each evaluation technique. Most is near-free open-source plus a small judge-model token cost. LangSmith is the expensive licensed one. Control cost with a small judge and subset-on-PR.")

# 12 TOOL SELECTION HARNESS
s=add(); title(s,"Evaluation: checking the agent picked the right tool",12,"Component: Evaluation")
_txt(s,Inches(0.6),Inches(1.4),Inches(12.13),Inches(0.4),
     [[("When an MCP tool server exposes several tools, a wrong tool that still answers is unreliable. We test it in code:",12.5,False,GRAYTX,True)]])
codepanel(s,Inches(0.6),Inches(1.95),Inches(7.6),Inches(3.4),
    ["# platform/evaluators/tool_selection.py","def evaluate_tool_selection(cases, run_agent):","  for c in cases:            # input + expected_tool (+ args)",
     "    trace  = run_agent(c['input'])","    chosen = trace.tool_calls[0].name if trace.tool_calls else None",
     "    args_ok = compare_args(trace.tool_calls[0].args, c['expected_args'])","    record(correct = chosen == c['expected_tool'], args_ok = args_ok)",
     "  return dict(accuracy=..., wrong_tool_rate=...,","              missing_tool_rate=..., arg_correctness=...)"],titlebar="custom Python (not Ragas/DeepEval)")
bullets(s,[("we know the right tool for each test case",0),("run the agent; read the tool it chose from the trace",0),
    ("score: accuracy, per-tool precision/recall, wrong-tool, missing-tool, argument correctness",0),
    ("this is why observability records the tool call",0,True,NAVY)],
    x=Inches(8.4),y=Inches(2.05),w=Inches(4.3),h=Inches(3.3),size=11.5,gap=10)
notes(s,"Kiran's MCP point. The harness reads the chosen tool from the trace and compares to the expected tool. Ragas/DeepEval don't do this — it's a small custom check that lives in the shared evaluators folder.")

# 13 OBSERVABILITY TRACE TREE
s=add(); title(s,"Observability: what gets tracked on every request",13,"Component: Observability")
spinex=Inches(1.0)
box(s,Inches(0.7),Inches(1.5),Inches(6.5),Inches(0.58),"Request (trace)","one call analysed / one candidate screened",fill=NAVY,headc=WHITE,hs=12,ss=9,subc=RGBColor(0xAE,0xC6,0xDE))
connect(s,spinex,Inches(2.08),spinex,Inches(4.95),color=LINEC,w=1.6)
connect(s,spinex,Inches(2.63),Inches(1.3),Inches(2.63),color=LINEC,w=1.6)
box(s,Inches(1.3),Inches(2.36),Inches(5.9),Inches(0.55),"Agent step (span)","dimension analysis / résumé rank",fill=BLUE,headc=WHITE,hs=11.5,ss=9,subc=RGBColor(0xCF,0xDC,0xF0))
connect(s,Inches(2.85),Inches(2.91),Inches(2.85),Inches(3.14),color=LINEC,w=1.4)
connect(s,Inches(5.55),Inches(2.91),Inches(5.55),Inches(3.14),color=LINEC,w=1.4)
box(s,Inches(1.55),Inches(3.14),Inches(2.6),Inches(0.5),"Model call (span)",None,fill=TEAL,headc=WHITE,hs=10.5)
box(s,Inches(4.3),Inches(3.14),Inches(2.6),Inches(0.5),"Tool call (span)",None,fill=DTEAL,headc=WHITE,hs=10.5)
connect(s,spinex,Inches(3.93),Inches(1.3),Inches(3.93),color=LINEC,w=1.6)
box(s,Inches(1.3),Inches(3.66),Inches(5.9),Inches(0.55),"Agent session","links multi-turn / the whole pipeline run",fill=SLATE,headc=WHITE,hs=11.5,ss=9,subc=RGBColor(0xD5,0xDD,0xE8))
connect(s,spinex,Inches(4.83),Inches(1.3),Inches(4.83),color=LINEC,w=1.6)
box(s,Inches(1.3),Inches(4.56),Inches(5.9),Inches(0.5),"Feedback event","thumbs / edit / override — same trace id",fill=LIGHT,hs=11.5,ss=9)
rrect(s,Inches(7.55),Inches(2.36),Inches(5.15),Inches(2.7),RGBColor(0xEA,0xF5,0xF2),line=RGBColor(0xBF,0xD6,0xD2))
_txt(s,Inches(7.75),Inches(2.48),Inches(4.8),Inches(0.35),[[("Answers the three questions directly:",12.5,True,TEAL)]])
bullets(s,[("model calls → tracked (model+version, prompt+version, tokens, cost, latency)",0),
    ("tool calls → tracked (tool, args, result, and was-it-the-correct-tool)",0),
    ("agent sessions → tracked (one trace id links every step and turn)",0)],
    x=Inches(7.75),y=Inches(2.98),w=Inches(4.75),h=Inches(2.0),size=11.5,gap=9)
notes(s,"Kiran's exact three questions: model calls, tool calls, agent sessions. Every box is a span tied by one trace id. The 'correct tool' field powers tool-selection evaluation.")

# 14 OBSERVABILITY — COST TRACKING
s=add(); title(s,"Observability: how cost is tracked (App Insights + Langfuse)",14,"Component: Observability")
box(s,Inches(0.7),Inches(1.7),Inches(3.2),Inches(1.3),"Each model call","span sets app.cost_usd = tokens × unit price",fill=LIGHT,hs=12.5,ss=10)
arrow(s,Inches(3.95),Inches(2.3))
box(s,Inches(4.4),Inches(1.6),Inches(3.6),Inches(1.5),"App Insights","the record source: query cost by use case / day / model (KQL, Workbook)",fill=RGBColor(0xEA,0xF5,0xF2),headc=TEAL,hs=12.5,ss=9.5)
box(s,Inches(4.4),Inches(3.3),Inches(3.6),Inches(1.3),"Langfuse","ready-made cost dashboards per model / prompt / user",fill=LIGHT,hs=12.5,ss=9.5)
connect(s,Inches(3.9),Inches(2.35),Inches(4.4),Inches(2.35)); connect(s,Inches(2.3),Inches(3.0),Inches(2.3),Inches(3.95)); connect(s,Inches(2.3),Inches(3.95),Inches(4.4),Inches(3.95))
box(s,Inches(8.4),Inches(2.35),Inches(4.3),Inches(1.3),"Azure Cost Management","reconcile monthly vs the actual invoice",fill=LIGHT,hs=12.5,ss=9.5)
arrow(s,Inches(8.05),Inches(2.85))
_txt(s,Inches(0.7),Inches(5.0),Inches(12.13),Inches(1.2),
     [[("No double-counting: ",12.5,True,NAVY),("cost is computed ONCE on the span; both tools read the same attribute — App Insights is the "
       "queryable record of truth, Langfuse is the day-to-day LLM cost view. We reconcile monthly against Azure Cost Management (the invoice).",12,False,GRAYTX)]])
notes(s,"Point 7: how cost is tracked across both tools without double counting. Cost is on the span once; App Insights aggregates it, Langfuse visualises it, Azure Cost Management is the invoice reconciliation.")

# 15 GUARDRAILS LIST
s=add(); title(s,"Guardrails: the list, and how we implement each",15,"Component: Guardrails")
table(s,Inches(0.6),Inches(1.42),Inches(12.13),
      ["Guardrail","What it stops","How we implement it"],
      [["Prompt injection / jailbreak","hijacking the instructions","Azure AI Content Safety — Prompt Shields"],
       ["Unsafe content","hate / violence / sexual / self-harm","Azure AI Content Safety"],
       ["PII leakage","exposing personal data","Azure AI Language PII, or Presidio (open source)"],
       ["Hallucination","ungrounded claims","Content Safety groundedness + our eval"],
       ["Off-topic / out-of-scope","answering the wrong thing","system-prompt limits + NeMo Guardrails / guardrails-ai"],
       ["Bad output format","broken JSON / schema","JSON schema / Pydantic / guardrails-ai"],
       ["Secrets / exfiltration","leaking secrets or data","output scanning + Purview DLP + regex"],
       ["Rate / cost abuse","runaway spend","API Management policies + budget alerts"],
       ["Protected material","copyrighted text","Content Safety protected-material detection"]],
      cw=[Inches(3.0),Inches(3.4),Inches(5.73)],fs=9.5,rh=Inches(0.47))
_txt(s,Inches(0.6),Inches(6.65),Inches(12),Inches(0.35),[[("Run as input checks (before the model) and output checks (before returning/storing).",11.5,False,MUTE,True)]])
notes(s,"Point 8: the actual guardrail list with the tool/package for each. Most are Azure Content Safety; PII can use Presidio; format uses schema validation; abuse uses APIM. We wrap input and output.")

# 16 DATA ACCESS LAYER
s=add(); title(s,"Data: not just RAG — structured, unstructured, documents",16,"Component: Data")
table(s,Inches(0.6),Inches(1.45),Inches(12.13),
      ["Data type","How we handle it","Azure service","Reusable tool"],
      [["Unstructured text","retrieval-augmented generation (RAG)","Azure AI Search","search_knowledge"],
       ["Structured data (SQL)","a SQL / NL2SQL agent — NOT RAG; read-only, allow-listed tables","Azure SQL / databases","query_sql"],
       ["Documents / files (PDF, scans, forms)","extract text/fields, then RAG or structured","AI Document Intelligence","extract_document"],
       ["Systems of record","direct read via a tool","CRM / ATS / APIs","get_record"]],
      cw=[Inches(2.9),Inches(4.5),Inches(2.5),Inches(2.23)],fs=10,rh=Inches(0.62))
rrect(s,Inches(0.6),Inches(4.6),Inches(12.13),Inches(1.5),RGBColor(0xEA,0xF5,0xF2),line=RGBColor(0xBF,0xD6,0xD2))
_txt(s,Inches(0.8),Inches(4.72),Inches(11.7),Inches(1.35),
     [[("Reusable tool catalog (platform/tools/): ",12.5,True,TEAL),
       ("we build these tools once — search_knowledge, query_sql, extract_document, get_record — and every use case "
        "composes them. New tools are added to the catalog and reused. Structured data goes through query_sql with "
        "guardrails (read-only, parameterised), not through RAG.",12,False,GRAYTX)]])
notes(s,"Point 9: data is more than RAG. Structured data uses a SQL/NL2SQL tool, not RAG; documents use Document Intelligence; and we predefine a reusable tool catalog that use cases compose.")

# 17 SERVING / HOSTING (why + canary + gate)
s=add(); title(s,"Serving & hosting — what runs where, and why",17,"Component: Serving")
table(s,Inches(0.6),Inches(1.42),Inches(12.13),
      ["Piece","Why we need it","What it does"],
      [["Container Apps","host the pipeline services","each step a container; autoscale; scale-to-zero"],
       ["Azure Functions","event triggers & schedules","e.g. APIX runs when a new transcript lands, or nightly — serverless, cheap for bursts"],
       ["Foundry Agent Service","managed agent hosting","runs agents with state/memory so we don't run our own server (adopt as it matures)"],
       ["API Management","the gateway","one entry point, quotas, token metering, caching"]],
      cw=[Inches(2.6),Inches(3.6),Inches(5.93)],fs=10,rh=Inches(0.58))
box(s,Inches(0.6),Inches(4.35),Inches(6.0),Inches(1.9),"Canary release (what it means)",
    "Deploy the new version, send ~10% of traffic to it, watch health + errors + eval signals for a short window, then ramp to 100% if healthy — or auto-rollback if not.",
    fill=RGBColor(0xEA,0xF5,0xF2),headc=TEAL,hs=12.5,ss=11)
box(s,Inches(6.75),Inches(4.35),Inches(5.98),Inches(1.9),"Promotion gate (what it means)",
    "The condition to move dev → test → prod: a human approver signs off AND the full evaluation set passes. No promotion without both.",
    fill=LIGHT,hs=12.5,ss=11)
notes(s,"Point 10: WHY each piece exists. Functions = event/schedule triggers (a new transcript, nightly batch). Foundry Agent Service = managed agent hosting. Define canary (10%->100% with rollback) and the promotion gate (approver + eval pass) plainly.")

# 18 FEEDBACK LOOP (each step described)
s=add(); title(s,"Feedback & improvement — what each step means",18,"Component: Feedback")
steps=[("1 Capture feedback","thumbs + reason, coach edits, overrides — tied to the trace id"),
       ("2 Land it","stored as scores / events in App Insights + Langfuse"),
       ("3 Triage negatives","review low-rated / failed answers; sort by cause (bad retrieval, wrong tool, weak prompt, missing data) and prioritise"),
       ("4 Add to golden set","turn confirmed bad cases into new test cases with the correct expected answer"),
       ("5 Fix & re-evaluate","change the prompt / retrieval / agent; run the evaluation gate"),
       ("6 Ship","release through CI/CD")]
y=Inches(1.55)
for h,sub in steps:
    rrect(s,Inches(0.7),y,Inches(3.4),Inches(0.72),LIGHT); _txt(s,Inches(0.85),y,Inches(3.2),Inches(0.72),[[(h,12.5,True,NAVY)]],anchor=MSO_ANCHOR.MIDDLE)
    _txt(s,Inches(4.35),y,Inches(8.3),Inches(0.72),[[(sub,11.5,False,GRAYTX)]],anchor=MSO_ANCHOR.MIDDLE)
    y=y+Inches(0.82)
notes(s,"Point 11: each step in plain words. 'Triage negatives' = review the bad responses, group them by root cause, and prioritise which to fix. That is the heart of the improvement loop.")

# 19 ONBOARDING A USE CASE (honest)
s=add(); title(s,"Onboarding a new use case — inherited vs defined",19,"Reusability")
rrect(s,Inches(0.6),Inches(1.5),Inches(6.0),Inches(4.6),RGBColor(0xEA,0xF5,0xF2),line=RGBColor(0xBF,0xD6,0xD2))
_txt(s,Inches(0.6),Inches(1.65),Inches(6.0),Inches(0.4),[[("Inherited from the platform (free)",13.5,True,TEAL)]],align=PP_ALIGN.CENTER)
bullets(s,[("CI/CD pipeline + evaluation gate",0),("observability & tracing",0),("model catalog & routing",0),
    ("guardrails engine",0),("reusable tool catalog (RAG, SQL, docs)",0),("gateway, identity, secrets",0),("dashboards framework",0)],
    x=Inches(0.85),y=Inches(2.3),w=Inches(5.5),h=Inches(3.7),size=12,gap=10)
rrect(s,Inches(6.75),Inches(1.5),Inches(5.98),Inches(4.6),LIGHT)
_txt(s,Inches(6.75),Inches(1.65),Inches(5.98),Inches(0.4),[[("Defined per use case (varies)",13.5,True,NAVY)]],align=PP_ALIGN.CENTER)
bullets(s,[("prompts + agent / pipeline design",0),("data sources + connectors",0),("retrieval / index setup",0),
    ("tools: reuse from catalog OR build new",0),("guardrail policy for this use case",0),("golden dataset + thresholds",0),
    ("evaluation config, dashboards",0),("often: use-case-specific integration / UI",0)],
    x=Inches(7.0),y=Inches(2.3),w=Inches(5.5),h=Inches(3.7),size=12,gap=9)
_txt(s,Inches(0.6),Inches(6.35),Inches(12.13),Inches(0.5),
     [[("Honest note: ",12.5,True,ROSE),("each use case genuinely differs — it is not \"just add four files.\" But it reuses the whole left column, so it is far faster than starting over.",12,False,GRAYTX)]])
notes(s,"Point 12: correct the misleading 'just add prompts/agents/tools/golden data'. A use case inherits the shared platform but defines a real list of its own things and differs case to case. Still much faster because the left column is reused.")

# 20 HOSTING COST + CAPABILITIES
s=add(); title(s,"Hosting: capability and cost per service (indicative)",20,"Cost")
table(s,Inches(0.6),Inches(1.42),Inches(12.13),
      ["Service","Capability","Pricing model","Indicative / month"],
      [["Azure OpenAI","the models","per token (GPT-5.5 ≈ $5 in / $30 out per 1M; cached in ≈ $0.50; mini/nano cheaper)","usage-driven (biggest)"],
       ["  — or PTU","reserved throughput","provisioned units","≈ $2,448 / unit"],
       ["Azure AI Search","RAG index","per search unit","Basic ≈ $74 · S1 ≈ $245"],
       ["Container Apps","run services","consumption, scale-to-zero","tens of $ (small)"],
       ["Functions","event triggers","per execution","~negligible (low volume)"],
       ["Cosmos DB / SQL","state & scores","serverless / provisioned","tens of $ (small)"],
       ["App Insights","observability","per GB ingested","tens of $ (volume-based)"],
       ["Langfuse (self-host)","LLM obs + prompt mgmt","MIT-free + infra","≈ $50–150"],
       ["Content Safety / APIM","guardrails / gateway","per 1k records / tiered","minor / tiered"]],
      cw=[Inches(2.2),Inches(2.5),Inches(4.9),Inches(2.53)],fs=9.5,rh=Inches(0.44))
_txt(s,Inches(0.6),Inches(6.5),Inches(12.13),Inches(0.4),
     [[("All indicative — confirm at a sizing exercise. ",11.5,True,NAVY),("Model tokens dominate the bill; everything else is modest fixed cost. Prompt caching and right-sizing cut the model cost.",11,False,GRAYTX)]])
notes(s,"Point 12: capability + cost, since that's where the client's interest is. Model tokens are the main variable; the rest is modest fixed cost. All figures indicative, to confirm at sizing.")

# 21 HOW IT ALL FITS
s=add(); title(s,"How it all fits together (end to end)",21,"Overview")
box(s,Inches(0.5),Inches(1.55),Inches(1.85),Inches(1.5),"Channels & triggers","web, chat,\nnew transcript,\nschedule",fill=LIGHT,hs=11.5,ss=9)
arrow(s,Inches(2.4),Inches(2.2),w=Inches(0.25))
box(s,Inches(2.7),Inches(1.55),Inches(1.7),Inches(1.5),"API Management","gateway,\nquotas,\nmetering",fill=LIGHT,hs=11.5,ss=9)
arrow(s,Inches(4.45),Inches(2.2),w=Inches(0.25))
box(s,Inches(4.75),Inches(1.55),Inches(4.2),Inches(1.5),"Orchestration / pipeline (Container Apps)","agents → Model Router→models · data-access tools (RAG / SQL / Docs) · guardrails",fill=RGBColor(0xEA,0xF5,0xF2),headc=TEAL,hs=11.5,ss=9)
arrow(s,Inches(9.0),Inches(2.2),w=Inches(0.25))
box(s,Inches(9.3),Inches(1.55),Inches(3.4),Inches(1.5),"Systems of record","CRM / ATS /\ndatabases /\ndocuments",fill=LIGHT,hs=11.5,ss=9)
# cross-cutting band
rrect(s,Inches(0.5),Inches(3.35),Inches(12.2),Inches(1.5),LIGHT)
_txt(s,Inches(0.7),Inches(3.45),Inches(12),Inches(0.35),[[("Cross-cutting (every request & change):",12,True,NAVY)]])
for t_,c,x in [("Observability: App Insights + Langfuse",TEAL,0.7),("Evaluation gate: GitHub CI/CD",AMBER,5.2),("Feedback → golden datasets",PURP,9.2)]:
    rrect(s,Inches(x),Inches(3.9),Inches(3.9 if x<9 else 3.3),Inches(0.7),WHITE,line=LINEC); _txt(s,Inches(x+0.15),Inches(3.9),Inches(3.7),Inches(0.7),[[(t_,11.5,True,c)]],anchor=MSO_ANCHOR.MIDDLE)
_txt(s,Inches(0.6),Inches(5.15),Inches(12.13),Inches(1.2),
     [[("Four flows: ",12.5,True,NAVY),("(1) a change is gated by evaluation before it ships; (2) a live request runs channels → gateway → "
       "pipeline → systems; (3) every step is traced with cost & quality; (4) feedback flows back into the golden datasets. "
       "That closed loop is the LLMOps setup.",12,False,GRAYTX)]])
notes(s,"Point 13: richer end-to-end with more blocks and the four flows described. Renamed from 'Architecture'. Walk the top row left to right, then the cross-cutting band, then the four flows.")

# 22 CONSOLIDATED TODAY -> OURS
s=add(); title(s,"What changes, in one table",22,"Summary")
table(s,Inches(0.6),Inches(1.45),Inches(12.13),
      ["Component","Today (to confirm)","What we add"],
      [["Source control / CI-CD","repo, manual releases","evaluation gate; gated, reversible deploys"],
       ["Prompt management","prompts inside code files","YAML per prompt + eval gate + registry for rollback/compare"],
       ["Model management","model names in code","task-aliases in config; swap must pass the gate"],
       ["Evaluation","manual spot-checks","golden datasets + Ragas/DeepEval/custom Python, run as a gate"],
       ["Observability","app logs","full trace tree; model/tool/agent + cost tracked"],
       ["Guardrails","minimal","Content Safety + PII + schema + APIM, input & output"],
       ["Data access","ad-hoc, RAG only","reusable tools: RAG + SQL + documents"],
       ["Serving / deploy","manual","Container Apps + APIM; canary + rollback"]],
      cw=[Inches(2.4),Inches(3.5),Inches(6.23)],fs=10,rh=Inches(0.5))
notes(s,"The one table Kiran wanted: for each component, what exists today and exactly what we add.")

# 23 SUMMARY BEFORE/AFTER
s=add(); title(s,"Summary: before and after",23,"Wrap-up")
table(s,Inches(0.6),Inches(1.45),Inches(12.13),
      ["Question","Today","With this LLMOps setup"],
      [["Can a bad change reach users?","yes — nothing blocks it","no — the evaluation gate blocks it"],
       ["How do we change a prompt safely?","edit code, hope","edit a YAML, gate + rollback"],
       ["Do we know answer quality?","by spot-checking","measured every change + in production"],
       ["Do we know the cost?","the monthly invoice","per request / model / use case, live"],
       ["Can we debug a bad answer?","dig through logs","follow one trace to the exact step"],
       ["How fast to add a use case?","start over","reuse the platform; add the use-case parts"],
       ["How do we swap a model?","code change","config change through the gate"]],
      cw=[Inches(3.6),Inches(3.6),Inches(4.93)],fs=10,rh=Inches(0.5))
_txt(s,Inches(0.6),Inches(6.35),Inches(12.13),Inches(0.5),
     [[("Next: ",12.5,True,NAVY),("confirm the current state with the team, stand up the backbone (source control + CI/CD + the gate), then onboard components one by one.",12,False,GRAYTX)]])
notes(s,"Point 14: a real before/after so the value is concrete. Close on the next step — confirm current state, build the backbone, onboard components one by one. No dates.")

# 24 GLOSSARY
s=add(); title(s,"Terms in plain English",24,"Appendix")
gloss=[("LLMOps","running LLM apps reliably — versioned, tested, released, monitored"),
       ("Golden dataset","saved test cases with expected answers/rules; the evaluation gate"),
       ("Evaluation gate","an automatic check that blocks a change if quality drops"),
       ("Threshold","the pass mark for a metric (baseline-relative or absolute floor)"),
       ("Registry","where prompts are held so we can roll back / swap versions"),
       ("Langfuse","open-source, self-hosted prompt mgmt + observability (MIT)"),
       ("LangSmith","a licensed platform for evaluation + observability"),
       ("Ragas / DeepEval","open-source frameworks that score RAG & writing quality"),
       ("LLM-as-judge","using a model to grade another model's answer"),
       ("Trace / span","a trace is one request; a span is one step inside it"),
       ("MCP","Model Context Protocol — the standard way to give agents tools"),
       ("RAG","Retrieval-Augmented Generation — answer from our own unstructured data"),
       ("NL2SQL","turning a question into a database (SQL) query"),
       ("Canary","release to a small % of traffic first, then all"),
       ("Promotion gate","approver + eval pass needed to move to the next environment"),
       ("PTU","Provisioned Throughput Unit — reserved model capacity, flat monthly"),
       ("Prompt Shields","Azure Content Safety feature that blocks prompt injection"),
       ("Triage","sort the failures by cause and prioritise what to fix")]
colw=Inches(5.95); x0=Inches(0.6); y0=Inches(1.45); rowh=Inches(0.56)
for i,(term,defn) in enumerate(gloss):
    r=i//2; c=i%2; x=x0+c*(colw+Inches(0.25)); y=y0+r*rowh
    _txt(s,x,y,colw,rowh,[[(term+"  ",11.5,True,TEAL),(defn,10,False,GRAYTX)]],ls=1.0)
notes(s,"Appendix — leave up for questions. New terms added: threshold, NL2SQL, canary, promotion gate, PTU, Prompt Shields, triage.")

out=os.path.join(os.path.dirname(os.path.dirname(os.path.abspath(__file__))),"presentation","LLMOps-Implementation-v2.pptx")
prs.save(out)
print(f"Saved: {out}  ({len(prs.slides._sldIdLst)} slides)")
