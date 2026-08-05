#!/usr/bin/env python3
"""Walkthrough deck: LLMOps approach for APIX & Hiring Intelligence.
Fully editable (native shapes only, no images). Diagram-led. No timelines. Speaker notes.
Requires: python-pptx.  Output: presentation/LLMOps-Approach-Walkthrough.pptx
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR

NAVY=RGBColor(0x1F,0x3A,0x5F); BLUE=RGBColor(0x2F,0x5C,0x9E); TEAL=RGBColor(0x2A,0x9D,0x8F)
GRAYTX=RGBColor(0x3C,0x46,0x54); LIGHT=RGBColor(0xEE,0xF1,0xF5); MUTE=RGBColor(0x6B,0x74,0x80)
WHITE=RGBColor(0xFF,0xFF,0xFF); LINEC=RGBColor(0xC9,0xD2,0xDD); AMBER=RGBColor(0xC8,0x7B,0x1E)
SKY=RGBColor(0x3E,0x6F,0xB0); SLATE=RGBColor(0x51,0x6B,0x8A); DTEAL=RGBColor(0x2E,0x7D,0x7D)
PURP=RGBColor(0x6A,0x51,0x9E); GREEN=RGBColor(0x3E,0x8E,0x5A); ROSE=RGBColor(0xB0,0x4A,0x5A)
FONT="Calibri"; SW,SH=Inches(13.333),Inches(7.5)
prs=Presentation(); prs.slide_width=SW; prs.slide_height=SH; BLANK=prs.slide_layouts[6]

def add():
    s=prs.slides.add_slide(BLANK)
    r=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,0,0,SW,SH); r.fill.solid(); r.fill.fore_color.rgb=WHITE; r.line.fill.background(); r.shadow.inherit=False
    return s
def _txt(s,x,y,w,h,runs,align=PP_ALIGN.LEFT,anchor=MSO_ANCHOR.TOP,sa=6,ls=1.05):
    tf=s.shapes.add_textbox(x,y,w,h).text_frame; tf.word_wrap=True; tf.vertical_anchor=anchor
    tf.margin_left=0; tf.margin_right=0; tf.margin_top=0; tf.margin_bottom=0
    for i,para in enumerate(runs):
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph()
        p.alignment=align; p.space_after=Pt(sa); p.space_before=Pt(0); p.line_spacing=ls
        for (t,sz,b,c,*rest) in para:
            it=rest[0] if rest else False
            r=p.add_run(); r.text=t; r.font.size=Pt(sz); r.font.bold=b; r.font.italic=it; r.font.name=FONT; r.font.color.rgb=c
    return tf
def fillrect(s,x,y,w,h,color):
    sp=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,x,y,w,h); sp.fill.solid(); sp.fill.fore_color.rgb=color
    sp.line.fill.background(); sp.shadow.inherit=False; return sp
def rrect(s,x,y,w,h,fill,line=LINEC,lw=1.0):
    sp=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,x,y,w,h); sp.fill.solid(); sp.fill.fore_color.rgb=fill
    if line is None: sp.line.fill.background()
    else: sp.line.color.rgb=line; sp.line.width=Pt(lw)
    sp.shadow.inherit=False
    try: sp.adjustments[0]=0.08
    except Exception: pass
    return sp
def settext(sp,lines,anchor=MSO_ANCHOR.MIDDLE,align=PP_ALIGN.CENTER):
    tf=sp.text_frame; tf.word_wrap=True; tf.vertical_anchor=anchor
    tf.margin_left=Inches(0.08); tf.margin_right=Inches(0.08); tf.margin_top=Inches(0.03); tf.margin_bottom=Inches(0.03)
    for i,(t,sz,b,c) in enumerate(lines):
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph()
        p.alignment=align; p.space_after=Pt(1); p.line_spacing=1.0
        r=p.add_run(); r.text=t; r.font.size=Pt(sz); r.font.bold=b; r.font.name=FONT; r.font.color.rgb=c
def box(s,x,y,w,h,head,sub=None,fill=LIGHT,headc=NAVY,hs=12.5,ss=10,subc=GRAYTX,align=PP_ALIGN.CENTER):
    sp=rrect(s,x,y,w,h,fill); lines=[(head,hs,True,headc)]
    if sub: lines.append((sub,ss,False,subc))
    settext(sp,lines,align=align); return sp
def arrow(s,x,y,w=Inches(0.34),h=Inches(0.3),color=TEAL,shape=MSO_SHAPE.RIGHT_ARROW):
    sp=s.shapes.add_shape(shape,x,y,w,h); sp.fill.solid(); sp.fill.fore_color.rgb=color; sp.line.fill.background(); sp.shadow.inherit=False; return sp
def connect(s,x1,y1,x2,y2,color=LINEC,w=1.4):
    ln=s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,x1,y1,x2,y2); ln.line.color.rgb=color; ln.line.width=Pt(w); return ln
def title(s,text,num,kicker=None):
    if kicker: _txt(s,Inches(0.7),Inches(0.32),Inches(11.9),Inches(0.3),[[(kicker.upper(),11,True,TEAL)]])
    _txt(s,Inches(0.7),Inches(0.56),Inches(12.0),Inches(0.7),[[(text,23,True,NAVY)]])
    fillrect(s,Inches(0.72),Inches(1.18),Inches(1.5),Pt(3),TEAL)
    _txt(s,SW-Inches(1.1),SH-Inches(0.42),Inches(0.7),Inches(0.3),[[(str(num),10,False,MUTE)]],align=PP_ALIGN.RIGHT)
    _txt(s,Inches(0.7),SH-Inches(0.42),Inches(9),Inches(0.3),[[("LLMOps approach — APIX & Hiring Intelligence",9,False,MUTE)]])
def bullets(s,items,x=Inches(0.75),y=Inches(1.5),w=Inches(11.8),h=Inches(5.2),size=15,gap=9,color=GRAYTX):
    tf=s.shapes.add_textbox(x,y,w,h).text_frame; tf.word_wrap=True; tf.margin_left=0; tf.margin_top=0
    for i,it in enumerate(items):
        text=it[0]; lvl=it[1] if len(it)>1 else 0; bold=it[2] if len(it)>2 else False; col=it[3] if len(it)>3 else color
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph()
        p.space_after=Pt(gap); p.space_before=Pt(0); p.line_spacing=1.05; p.level=lvl
        r=p.add_run(); r.text=("•  " if lvl==0 else "–  ")+text
        r.font.size=Pt(size-lvl); r.font.bold=bold; r.font.name=FONT; r.font.color.rgb=col
    return tf
def table(s,x,y,w,headers,rows,cw,fs=11,hs=11,rh=Inches(0.5)):
    nr=len(rows)+1; t=s.shapes.add_table(nr,len(headers),x,y,w,rh*nr).table
    for i,cc in enumerate(cw): t.columns[i].width=cc
    for j,hh in enumerate(headers):
        c=t.cell(0,j); c.fill.solid(); c.fill.fore_color.rgb=NAVY; c.vertical_anchor=MSO_ANCHOR.MIDDLE
        c.margin_left=Inches(0.08); c.margin_top=Inches(0.02); c.margin_bottom=Inches(0.02)
        r=c.text_frame.paragraphs[0].add_run(); r.text=hh; r.font.size=Pt(hs); r.font.bold=True; r.font.color.rgb=WHITE; r.font.name=FONT
    for i,row in enumerate(rows):
        for j,v in enumerate(row):
            c=t.cell(i+1,j); c.fill.solid(); c.fill.fore_color.rgb=WHITE if i%2==0 else LIGHT; c.vertical_anchor=MSO_ANCHOR.MIDDLE
            c.margin_left=Inches(0.08); c.margin_top=Inches(0.02); c.margin_bottom=Inches(0.02)
            r=c.text_frame.paragraphs[0].add_run(); r.text=v; r.font.size=Pt(fs); r.font.name=FONT; r.font.color.rgb=GRAYTX
            if j==0: r.font.bold=True; r.font.color.rgb=NAVY
    return t
def notes(s,text): s.notes_slide.notes_text_frame.text=text
def flow(s,steps,y,bh=Inches(1.15),x0=Inches(0.6),total=Inches(12.13),fills=None,hs=12,ss=9.5):
    n=len(steps); gap=Inches(0.2); bw=Inches((total.inches-gap.inches*(n-1))/n); x=x0
    for i,(h,sub) in enumerate(steps):
        f=fills[i] if fills else LIGHT
        box(s,x,y,bw,bh,h,sub,fill=f,hs=hs,ss=ss)
        if i<n-1: arrow(s,x+bw-Inches(0.02),y+bh/2-Inches(0.13),w=Inches(0.22),h=Inches(0.26))
        x=x+bw+gap
    return bw

print("Building 5thAug walkthrough deck...")

# 1 TITLE
s=add()
fillrect(s,0,0,Inches(0.28),SH,NAVY)
_txt(s,Inches(0.9),Inches(1.95),Inches(11.5),Inches(1.3),[[("LLMOps: our approach",32,True,NAVY)]])
fillrect(s,Inches(0.94),Inches(2.8),Inches(2.0),Pt(3),TEAL)
_txt(s,Inches(0.92),Inches(3.1),Inches(11.3),Inches(1.5),
     [[("Built around two live use cases — APIX and Hiring Intelligence — and reusable for any project that follows.",17,False,GRAYTX)],
      [("LLMOps = Large Language Model Operations: how we run these apps reliably (tracking, testing, releasing, improving).",13,False,MUTE,True)]],sa=10)
_txt(s,Inches(0.92),Inches(6.35),Inches(11),Inches(0.5),[[("Working draft for review  ·  approach & activities (no dates yet)",12,True,NAVY)]])
notes(s,"Frame it: this is the approach for discussion, not a finished plan, and deliberately no timelines yet — as agreed. We ground it in APIX and Hiring so it's concrete, but the approach applies to any future use case.")

# 2 WHAT WE'LL COVER (Kiran's asks)
s=add(); title(s,"What this covers",2,"Agenda")
items=[("The approach","how we run LLMOps",TEAL),("The activities","what work is involved",BLUE),
       ("As-is vs to-be","what exists, what changes",SLATE),("Observability","what we track on every request",DTEAL),
       ("Evaluation","how we measure quality",PURP),("Infrastructure","Azure services & hosting",NAVY)]
cw=Inches(3.85); ch=Inches(1.5); gx=Inches(0.25); gy=Inches(0.25); x0=Inches(0.7); y0=Inches(1.7)
for i,(h,sub,c) in enumerate(items):
    r=i//3; cidx=i%3; x=x0+cidx*(cw+gx); y=y0+r*(ch+gy)
    box(s,x,y,cw,ch,h,sub,fill=LIGHT); fillrect(s,x,y,Inches(0.1),ch,c)
_txt(s,Inches(0.7),Inches(5.4),Inches(11.9),Inches(0.6),[[("Evaluation and observability get the most detail — that is where the real work and the real questions are.",13,True,NAVY,True)]])
notes(s,"These are exactly the things asked for. Signal early that we'll spend the most time on evaluation and observability.")

# 3 THE APPROACH IN ONE PICTURE
s=add(); title(s,"The approach: wrap the pipelines with an operational layer",3,"Approach")
box(s,Inches(0.7),Inches(1.65),Inches(5.6),Inches(1.5),"What already exists","APIX and Hiring run as agent pipelines (built by the product team)",fill=LIGHT,hs=14)
arrow(s,Inches(6.4),Inches(2.3),w=Inches(0.45))
box(s,Inches(6.95),Inches(1.65),Inches(5.7),Inches(1.5),"What we add: the LLMOps layer","tracking, evaluation, prompt & model control, safe releases, feedback",fill=RGBColor(0xE9,0xF3,0xF0),headc=TEAL,hs=14)
_txt(s,Inches(0.7),Inches(3.5),Inches(12),Inches(0.4),[[("We are not rebuilding the use cases. We add the operational layer around them, standardise it, and make it reusable.",13,True,NAVY,True)]])
comps=[("Observability","see every step"),("Evaluation","measure quality"),("Prompt mgmt","versioned, tested"),
       ("Model mgmt","catalog & aliases"),("CI/CD","safe releases"),("Guardrails","safety & PII"),
       ("Data / RAG","feed from our data"),("Feedback","improve over time")]
cw=Inches(2.9); ch=Inches(0.85); x0=Inches(0.7); y0=Inches(4.1)
for i,(h,sub) in enumerate(comps):
    r=i//4; cidx=i%4; x=x0+cidx*(cw+Inches(0.16)); y=y0+r*(ch+Inches(0.16))
    box(s,x,y,cw,ch,h,sub,fill=LIGHT,hs=12,ss=9.5)
notes(s,"The core message. The pipelines exist. LLMOps is the operational layer we wrap around them. The eight tiles are the pieces of that layer — same layer serves both use cases and any future one.")

# 4 APIX PIPELINE
s=add(); title(s,"APIX as an agent pipeline",4,"Use case 1")
_txt(s,Inches(0.7),Inches(1.3),Inches(12),Inches(0.4),
     [[("APIX turns AI-analysed call transcripts into a weekly performance report (score /100, KPIs, trend, coaching, risk flags).",12.5,False,GRAYTX,True)]])
flow(s,[("Transcript","+ metadata"),("Dimension agents","sales, CX, retention, compliance"),("Extraction","escalations, sentiment, sales"),
        ("Scoring","composite /100 per program"),("Coaching report","steps + risk flags"),("Dashboard","managers & coaches")],
     y=Inches(2.0),bh=Inches(1.4),fills=[LIGHT,TEAL,BLUE,SLATE,PURP,DTEAL])
_txt(s,Inches(0.7),Inches(3.7),Inches(12),Inches(0.5),
     [[("Sequential pipeline (not agent-to-agent). Multi-program: Telesales and WCC use different scoring criteria.",12.5,False,GRAYTX,True)]])
_txt(s,Inches(0.7),Inches(4.35),Inches(12),Inches(0.4),[[("Evaluation focus for APIX:",13,True,NAVY)]])
bullets(s,[("Groundedness — coaching must cite what was actually said, not invent moments.",0),
    ("Scoring agreement — do the /100 dimension scores match a human QA reviewer?",0),
    ("Consistency & fairness — same rubric applied the same way across all agents and programs.",0)],
    x=Inches(0.9),y=Inches(4.75),w=Inches(11.5),h=Inches(1.8),size=13,gap=8)
notes(s,"Walk the pipeline left to right. WCC = the care program; Telesales = the sales program; they score differently, which matters for evaluation (separate golden sets per program). Then the three evaluation priorities — groundedness is the big one for a coaching tool.")

# 5 HIRING PIPELINE
s=add(); title(s,"Hiring Intelligence as an agent pipeline",5,"Use case 2")
flow(s,[("Intake / router","what's needed"),("Resume rank","parse + rank vs role (RAG)"),("Screening Q&A","answer candidate (RAG)"),
        ("Scoring & summary","fit score for recruiter"),("Recruiter","human decides")],
     y=Inches(1.75),bh=Inches(1.4),fills=[LIGHT,TEAL,BLUE,SLATE,GREEN])
_txt(s,Inches(0.7),Inches(3.45),Inches(12),Inches(0.5),
     [[("Uses tools through MCP (Model Context Protocol): applicant tracking system (ATS), requisition database, scheduling.",12.5,False,GRAYTX,True)]])
_txt(s,Inches(0.7),Inches(4.1),Inches(12),Inches(0.4),[[("Evaluation focus for Hiring:",13,True,NAVY)]])
bullets(s,[("Tool selection — did the agent call the RIGHT tool, with the right inputs? (Kiran's point.)",0,True,NAVY),
    ("RAG quality — is the ranking grounded in the job description and rubric?",0),
    ("Fairness — no bias in ranking or screening.",0)],
    x=Inches(0.9),y=Inches(4.5),w=Inches(11.5),h=Inches(1.8),size=13,gap=8)
_txt(s,Inches(0.7),Inches(6.35),Inches(12),Inches(0.35),
     [[("Together, APIX + Hiring exercise every evaluation metric group — so the approach generalises.",12.5,True,TEAL,True)]])
notes(s,"Hiring is the one with tools, so it's where tool-selection evaluation lives. MCP is just the standard way an agent is given tools. Note the pairing point at the bottom: the two use cases together cover everything, which is why the framework generalises.")

# 6 ACTIVITIES / WORKSTREAMS
s=add(); title(s,"The activities involved (no dates — sequencing only)",6,"Activities")
found=[("A. Discovery","assess current pipelines"),("B. Foundation","repo, Azure landing, gateway"),("C. Observability","tracing on both pipelines")]
prio=[("D. Evaluation","golden data, evaluators, CI gate")]
later=[("E. Prompt & model mgmt",""),("F. CI/CD & releases",""),("G. Data / RAG pipelines",""),("H. Guardrails & governance",""),("I. Feedback loop","")]
_txt(s,Inches(0.7),Inches(1.5),Inches(12),Inches(0.3),[[("Foundational (do first)",12,True,SLATE)]])
x=Inches(0.7)
for h,sub in found:
    box(s,x,Inches(1.85),Inches(3.55),Inches(0.95),h,sub,fill=LIGHT,hs=12,ss=9.5); x+=Inches(3.75)
_txt(s,Inches(0.7),Inches(3.0),Inches(12),Inches(0.3),[[("Priority — runs early and continuously",12,True,TEAL)]])
box(s,Inches(0.7),Inches(3.35),Inches(11.9),Inches(0.9),"D. Evaluation framework","golden datasets (per use case & program) · evaluators (Ragas, DeepEval, custom Python) · CI gate · online sampling · human review",fill=RGBColor(0xE9,0xF3,0xF0),headc=TEAL,hs=13,ss=11)
_txt(s,Inches(0.7),Inches(4.45),Inches(12),Inches(0.3),[[("Layer in as we go",12,True,PURP)]])
x=Inches(0.7)
for i,(h,sub) in enumerate(later):
    box(s,x,Inches(4.8),Inches(2.28),Inches(0.9),h,sub,fill=LIGHT,hs=11,ss=9); x+=Inches(2.4)
_txt(s,Inches(0.7),Inches(5.95),Inches(12),Inches(0.4),[[("No timelines yet — this is the order of work and its dependencies, not a schedule.",12.5,False,MUTE,True)]])
notes(s,"Nine workstreams. Foundational first (discovery, foundation, observability). Evaluation is called out as the priority that runs early and never really stops. The rest layer in. Stress: no dates — sequencing and dependencies only, as agreed.")

# 7 AS-IS / TO-BE
s=add(); title(s,"What exists vs what changes (as-is / to-be)",7,"Gap")
table(s,Inches(0.6),Inches(1.45),Inches(12.13),
      ["Area","As-is (to confirm in discovery)","Target (to-be)"],
      [["Prompts","edited in code or portal, untracked","versioned in Git + registry, reviewed, tested"],
       ["Models","model names in code","task-aliases in config; swap = reviewed change"],
       ["Tracing","app logs; no per-step detail","full trace tree: agent, model, tool spans"],
       ["Evaluation","manual / spot-check; no gate","golden datasets + automated scoring + CI gate"],
       ["Data / RAG","ad-hoc","managed ingestion + scheduled refresh"],
       ["Guardrails","minimal","safety + personal-data checks + human review"],
       ["Deploy","manual","automated, gated, gradual, reversible"]],
      cw=[Inches(2.2),Inches(5.0),Inches(4.93)],fs=11,rh=Inches(0.55))
_txt(s,Inches(0.6),Inches(6.2),Inches(12.13),Inches(0.5),
     [[("As-is is our assumption until we audit the pipelines with the team — every row is a discovery question, not a claim.",12,False,MUTE,True)]])
notes(s,"Be careful and honest here: we haven't audited the pipelines yet, so as-is is written as 'to confirm'. The value is the target column and the gap it implies. We're adding an operational layer, not rebuilding.")

# 8 OBSERVABILITY: TRACE TREE
s=add(); title(s,"Observability: every request is a trace of nested steps",8,"Observability")
# left tree with a single spine at x=1.0
spinex=Inches(1.0)
box(s,Inches(0.7),Inches(1.55),Inches(6.6),Inches(0.6),"Request (trace)","one call analysed / one candidate screened",fill=NAVY,headc=WHITE,hs=12.5,ss=9.5,subc=RGBColor(0xAE,0xC6,0xDE))
# spine down the left
connect(s,spinex,Inches(2.15),spinex,Inches(5.02),color=LINEC,w=1.6)
# Agent 1
connect(s,spinex,Inches(2.72),Inches(1.3),Inches(2.72),color=LINEC,w=1.6)
box(s,Inches(1.3),Inches(2.44),Inches(6.0),Inches(0.56),"Agent 1 (span)","dimension analysis / resume rank",fill=BLUE,headc=WHITE,hs=12,ss=9,subc=RGBColor(0xCF,0xDC,0xF0))
# children of agent 1: two short vertical drops from agent-1 bottom, no crossing lines
connect(s,Inches(2.9),Inches(3.0),Inches(2.9),Inches(3.26),color=LINEC,w=1.4)
connect(s,Inches(5.7),Inches(3.0),Inches(5.7),Inches(3.26),color=LINEC,w=1.4)
box(s,Inches(1.6),Inches(3.26),Inches(2.6),Inches(0.5),"Model call (span)",None,fill=TEAL,headc=WHITE,hs=10.5)
box(s,Inches(4.4),Inches(3.26),Inches(2.6),Inches(0.5),"Tool call (span)",None,fill=DTEAL,headc=WHITE,hs=10.5)
# Agent 2
connect(s,spinex,Inches(4.2),Inches(1.3),Inches(4.2),color=LINEC,w=1.6)
box(s,Inches(1.3),Inches(3.92),Inches(6.0),Inches(0.56),"Agent 2 (span)","scoring / screening",fill=BLUE,headc=WHITE,hs=12,ss=9,subc=RGBColor(0xCF,0xDC,0xF0))
# Final output
connect(s,spinex,Inches(4.92),Inches(1.3),Inches(4.92),color=LINEC,w=1.6)
box(s,Inches(1.3),Inches(4.64),Inches(6.0),Inches(0.54),"Final output","report / candidate summary",fill=LIGHT,hs=12)
# right side note (plain panel; header + bullets, no overlap)
rrect(s,Inches(7.75),Inches(2.45),Inches(4.95),Inches(2.75),RGBColor(0xE9,0xF3,0xF0),line=RGBColor(0xBF,0xD6,0xD2))
_txt(s,Inches(8.0),Inches(2.62),Inches(4.5),Inches(0.4),[[("Every box is a span",14,True,TEAL)]])
bullets(s,[("It records its own inputs, outputs, timing, and cost.",0),
    ("Child spans roll up to the parent — total cost and time for the whole request add up automatically.",0),
    ("One trace id ties everything together, so a bad answer can be followed back to the exact step.",0)],
    x=Inches(8.0),y=Inches(3.15),w=Inches(4.45),h=Inches(2.0),size=12,gap=9)
notes(s,"This is the mental model for the whole observability section. A request is a trace; each step is a span nested under its parent. Point at the model-call and tool-call spans under Agent 1. The trace id is the thread that lets us follow any answer back to its source.")

# 9 OBSERVABILITY: WHAT IS TRACKED
s=add(); title(s,"What we track at each level",9,"Observability")
table(s,Inches(0.6),Inches(1.4),Inches(12.13),
      ["Level","What we record"],
      [["Request","use case, program, input id, final output, status, total time, total tokens, total cost"],
       ["Agent step","agent name + version, its input/output, which model & tools it used, time, tokens, cost"],
       ["Model call","model name + version, prompt id + version, prompt, answer, tokens in/out, cost, latency"],
       ["Tool call","tool name, inputs, result, success/failure, AND whether it was the correct tool (for evaluation)"],
       ["Session","for multi-turn chats: links the turns, conversation id, user id (hashed)"],
       ["Feedback","thumbs up/down + reason, coach edits, recruiter overrides — linked by trace id"]],
      cw=[Inches(1.9),Inches(10.23)],fs=11.5,rh=Inches(0.62))
_txt(s,Inches(0.6),Inches(6.3),Inches(12.13),Inches(0.5),
     [[("Built on OpenTelemetry -> Azure Application Insights (system of record) + self-hosted Langfuse (the LLM-specific view).",12,True,NAVY,True)]])
notes(s,"This directly answers the four questions asked: per request, per model call, per tool call, per agent session. Highlight the tool-call row — recording the correct-tool flag is what makes tool-selection evaluation possible. Tools: Azure-native App Insights for the record of truth, Langfuse for the LLM-specific lens.")

# 10 EVALUATION: METRIC GROUPS
s=add(); title(s,"Evaluation: the metric groups",10,"Evaluation")
table(s,Inches(0.6),Inches(1.4),Inches(12.13),
      ["Group","Example metrics","Judges..."],
      [["Retrieval / RAG","groundedness, context relevance, precision/recall","is it backed by our data?"],
       ["Writing quality","coherence, tone, completeness","how well it reads"],
       ["Task execution","tool selection, correct action, steps taken","whether it did the right thing"],
       ["Safety / fairness","unsafe content, PII leak, bias","is it safe and fair?"],
       ["Operational","latency, cost, tokens","is it fast and affordable?"]],
      cw=[Inches(2.7),Inches(5.2),Inches(4.23)],fs=11.5,rh=Inches(0.6))
box(s,Inches(0.6),Inches(5.15),Inches(12.13),Inches(1.35),"Why writing quality and task execution are separate",
    "A fluent, well-written answer can still be wrong — e.g. it called the wrong tool. \"How it reads\" and \"did it do the right thing\" are different questions, so we score them separately. Some metrics overlap (coherence can sit under RAG or writing); we assign each to one group to avoid double-counting.",
    fill=RGBColor(0xE9,0xF3,0xF0),headc=TEAL,hs=13,ss=11.5)
notes(s,"Directly answers Kiran's question about the grouping. The point: writing quality and task execution are independent. Acknowledge the overlap he raised and say how we handle it — assign each metric to one group.")

# 11 EVALUATION: TOOL SELECTION
s=add(); title(s,"Evaluation: did the agent pick the right tool?",11,"Evaluation")
_txt(s,Inches(0.7),Inches(1.35),Inches(12),Inches(0.4),
     [[("When an MCP tool server exposes several tools, a wrong tool that still returns an answer is unreliable. We test this directly.",12.5,False,GRAYTX,True)]])
box(s,Inches(0.7),Inches(2.0),Inches(3.3),Inches(1.5),"Test case","input + the KNOWN correct tool & inputs",fill=LIGHT,hs=13)
arrow(s,Inches(4.05),Inches(2.65))
box(s,Inches(4.55),Inches(2.0),Inches(3.3),Inches(1.5),"Run the agent","read the tool it actually chose (from the trace)",fill=LIGHT,hs=13)
arrow(s,Inches(7.9),Inches(2.65))
box(s,Inches(8.4),Inches(2.0),Inches(4.25),Inches(1.5),"Compare & score","chosen vs expected, and were the inputs right?",fill=RGBColor(0xE9,0xF3,0xF0),headc=TEAL,hs=13)
_txt(s,Inches(0.7),Inches(3.75),Inches(12),Inches(0.4),[[("We measure:",13,True,NAVY)]])
bullets(s,[("tool-selection accuracy, and precision/recall per tool",0),
    ("wrong-tool rate · called a tool when none was needed · missed a tool it should have used",0),
    ("argument correctness — right tool, but were the inputs right?",0)],
    x=Inches(0.9),y=Inches(4.15),w=Inches(11.5),h=Inches(1.6),size=13,gap=8)
_txt(s,Inches(0.7),Inches(6.15),Inches(12),Inches(0.4),
     [[("Ragas / DeepEval don't cover this — it's custom Python that reads the tool call from the trace. Hiring's ATS tools are the example.",12,False,MUTE,True)]])
notes(s,"This is the slide Kiran will care about most. Simple idea: we know the right tool for each test case, we run the agent, we read what it actually chose from the trace, and we compare. The metrics catch every failure mode. Point out this needs custom Python — the RAG frameworks don't do it.")

# 12 EVALUATION: TOOLING OPTIONS
s=add(); title(s,"Evaluation: the tools, and what each is for",12,"Evaluation")
table(s,Inches(0.6),Inches(1.4),Inches(12.13),
      ["Tool","Covers","Open source?","Use it for"],
      [["Ragas","RAG metrics (groundedness, relevance)","Yes","APIX groundedness, Hiring RAG"],
       ["DeepEval","broad LLM eval, custom metrics, CI-friendly","Yes","the CI gate, writing quality"],
       ["Custom Python","tool selection, scoring vs labels, extraction","Yes (own code)","agent & tool behaviour"],
       ["LLM-as-judge","subjective quality with a rubric","depends","coaching usefulness, summaries"],
       ["Azure AI Foundry evals","built-in + custom, links to traces","No (Azure)","staying inside Azure"],
       ["LangSmith","eval + observability platform","No — licensed","only if we standardise on it"]],
      cw=[Inches(2.5),Inches(4.4),Inches(2.0),Inches(3.23)],fs=10.5,rh=Inches(0.56))
_txt(s,Inches(0.6),Inches(5.5),Inches(12.13),Inches(0.5),
     [[("Recommendation: a mix — ",13,True,NAVY),("Ragas + DeepEval + custom Python for the CI gate, Foundry for trace-linked runs. Not one single tool.",13,False,GRAYTX)]])
notes(s,"Honest options, not a single vendor. The recommendation is a mix. Flag LangSmith as licensed — it's good but not open source, so it's a cost decision, and we can do most of it with open tools.")

# 13 EVALUATION: HOW IT RUNS
s=add(); title(s,"Evaluation: how and when it runs",13,"Evaluation")
box(s,Inches(0.7),Inches(1.7),Inches(3.85),Inches(2.6),"Offline (before release)","",fill=LIGHT,hs=14)
bullets(s,[("golden datasets per use case AND per program",0),("runs on every change (CI gate)",0),("blocks release if quality drops",0)],x=Inches(0.95),y=Inches(2.3),w=Inches(3.4),h=Inches(1.9),size=11.5,gap=7)
box(s,Inches(4.75),Inches(1.7),Inches(3.85),Inches(2.6),"Online (in production)","",fill=LIGHT,hs=14)
bullets(s,[("sample a share of live traffic",0),("score it automatically in the background",0),("alert if quality drifts",0)],x=Inches(5.0),y=Inches(2.3),w=Inches(3.4),h=Inches(1.9),size=11.5,gap=7)
box(s,Inches(8.8),Inches(1.7),Inches(3.85),Inches(2.6),"Human review","",fill=LIGHT,hs=14)
bullets(s,[("coaches & recruiters give feedback",0),("experts review a sample",0),("findings become new test cases",0)],x=Inches(9.05),y=Inches(2.3),w=Inches(3.4),h=Inches(1.9),size=11.5,gap=7)
box(s,Inches(0.7),Inches(4.6),Inches(11.95),Inches(1.4),"Per-agent AND end-to-end",
    "We score each pipeline step (each APIX dimension agent, each Hiring agent) and the final output. A pipeline can look fine end-to-end while one step quietly gets worse — checking both catches that.",
    fill=RGBColor(0xE9,0xF3,0xF0),headc=TEAL,hs=13,ss=11.5)
notes(s,"Three modes: offline gate, online sampling, human review — they reinforce each other. The bottom box is an important nuance for pipelines: test each step and the whole thing, because one weak step can hide behind a decent final answer.")

# 14 INFRASTRUCTURE
s=add(); title(s,"Infrastructure & hosting on Azure (proposed)",14,"Infrastructure")
_txt(s,Inches(0.7),Inches(1.35),Inches(12),Inches(0.4),[[("Hosting the pipelines — the main choice:",13,True,NAVY)]])
table(s,Inches(0.6),Inches(1.75),Inches(12.13),
      ["Option","Fit","Note"],
      [["Azure Container Apps (recommended)","each pipeline step as a service; scales to zero","best general fit"],
       ["Azure Functions","event-driven triggers (new transcript, new candidate)","good for APIX batch jobs"],
       ["Foundry Agent Service","managed hosted agents, less to run","consider as it matures"]],
      cw=[Inches(4.0),Inches(5.13),Inches(3.0)],fs=11,rh=Inches(0.55))
_txt(s,Inches(0.6),Inches(4.0),Inches(12.13),Inches(0.35),[[("Around it (shared platform):",13,True,NAVY)]])
tiles=[("Azure OpenAI","models"),("AI Search","RAG"),("Cosmos DB / SQL","state & scores"),("API Management","gateway"),
       ("App Insights + Langfuse","observability"),("Entra ID / Key Vault","identity & secrets"),("Content Safety","guardrails"),("GitHub Actions","CI/CD")]
cw=Inches(2.9); ch=Inches(0.8); x0=Inches(0.6); y0=Inches(4.4)
for i,(h,sub) in enumerate(tiles):
    r=i//4; cidx=i%4; x=x0+cidx*(cw+Inches(0.14)); y=y0+r*(ch+Inches(0.14))
    box(s,x,y,cw,ch,h,sub,fill=LIGHT,hs=11.5,ss=9.5)
_txt(s,Inches(0.6),Inches(6.35),Inches(12),Inches(0.35),[[("Target setup shown — sequencing lives in the activities section, still no dates.",11.5,False,MUTE,True)]])
notes(s,"Recommend Container Apps for the pipeline services, plus Functions for triggers. The tiles are the shared platform every use case sits on. Reinforce: this is the target picture, and we deliberately keep dates out.")

# 15 SUMMARY & NEXT STEPS
s=add(); title(s,"Summary and next step",15,"Wrap-up")
bullets(s,[
    ("We wrap the two live pipelines (APIX, Hiring) with an operational layer, and make it reusable for any project.",0,True,NAVY),
    ("Both are sequential pipelines, not agent-to-agent — the same Ops setup serves both.",0),
    ("Observability: every request is a trace; we record model calls, tool calls, and agent steps.",0),
    ("Evaluation (the priority): metric groups, tool-selection testing, a mix of tools, offline + online + human.",0),
    ("Infrastructure: Azure Container Apps + a shared platform of Azure services.",0),
    ("Next: confirm the current state of the pipelines with the team, agree the approach, then plan the work.",0,True,NAVY),
],y=Inches(1.6),gap=13)
notes(s,"Recap the five threads. The ask is not a sign-off on a schedule — it's to align on the approach and set up the discovery of the current pipelines, per the agreed next step.")

# 16 GLOSSARY
s=add(); title(s,"Terms in plain English",16,"Appendix")
gloss=[("LLMOps","Large Language Model Operations — running LLM apps reliably."),
       ("APIX","Afni Performance Intelligence Index — the coaching dashboard."),
       ("Pipeline","agents run one after another; each step feeds the next."),
       ("Agent","an LLM that can take steps and use tools, not just chat."),
       ("Trace / span","a trace is one request; a span is one step inside it."),
       ("RAG","Retrieval-Augmented Generation — answering from our own documents."),
       ("MCP","Model Context Protocol — a standard way to give agents tools."),
       ("Golden dataset","saved test cases with expected answers, used to check quality."),
       ("Groundedness","whether an answer is backed by the actual source (transcript)."),
       ("Tool selection","whether the agent chose the correct tool for the task."),
       ("CI gate","an automatic check that blocks a release if quality drops."),
       ("OpenTelemetry","an open standard for collecting traces and metrics."),
       ("ATS","Applicant Tracking System — the recruiting system of record."),
       ("PII","Personally Identifiable Information — personal data to protect."),
       ("Container Apps / Functions","Azure ways to run our services."),
       ("WCC / Telesales","two Afni programs with different scoring criteria.")]
colw=Inches(5.95); x0=Inches(0.6); y0=Inches(1.45); rowh=Inches(0.63)
for i,(term,defn) in enumerate(gloss):
    r=i//2; c=i%2; x=x0+c*(colw+Inches(0.25)); y=y0+r*rowh
    _txt(s,x,y,colw,rowh,[[(term+"  ",12,True,TEAL),(defn,11,False,GRAYTX)]],ls=1.0)
notes(s,"Appendix — leave up for questions.")

out=os.path.join(os.path.dirname(os.path.dirname(os.path.abspath(__file__))),"presentation","LLMOps-Approach-Walkthrough.pptx")
prs.save(out)
print(f"Saved: {out}  ({len(prs.slides._sldIdLst)} slides)")
