#!/usr/bin/env python3
"""Practical LLMOps research deck — fully editable (native shapes only, no images).
Generic, research-phase tone. Diagram-led; text kept minimal.
Requires: python-pptx.  Output: presentation/LLMOps-Research-Deep-Dive.pptx
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR

NAVY=RGBColor(0x1F,0x3A,0x5F); BLUE=RGBColor(0x2F,0x5C,0x9E); TEAL=RGBColor(0x2A,0x9D,0x8F)
GRAYTX=RGBColor(0x3C,0x46,0x54); LIGHT=RGBColor(0xEE,0xF1,0xF5); MUTE=RGBColor(0x6B,0x74,0x80)
WHITE=RGBColor(0xFF,0xFF,0xFF); LINEC=RGBColor(0xC9,0xD2,0xDD); AMBER=RGBColor(0xC8,0x7B,0x1E)
SKY=RGBColor(0x3E,0x6F,0xB0); SLATE=RGBColor(0x51,0x6B,0x8A); DTEAL=RGBColor(0x2E,0x7D,0x7D)
PURP=RGBColor(0x6A,0x51,0x9E); GREEN=RGBColor(0x3E,0x8E,0x5A); ROSE=RGBColor(0xB0,0x4A,0x5A)
FONT="Calibri"; SW,SH=Inches(13.333),Inches(7.5)

prs=Presentation(); prs.slide_width=SW; prs.slide_height=SH; BLANK=prs.slide_layouts[6]

def add():
    s=prs.slides.add_slide(BLANK)
    r=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,0,0,SW,SH); r.fill.solid(); r.fill.fore_color.rgb=WHITE; r.line.fill.background(); r.shadow.inherit=False
    return s
def _txt(s,x,y,w,h,runs,align=PP_ALIGN.LEFT,anchor=MSO_ANCHOR.TOP,sa=6,ls=1.05):
    tf=s.shapes.add_textbox(x,y,w,h).text_frame; tf.word_wrap=True; tf.vertical_anchor=anchor
    tf.margin_left=0; tf.margin_right=0; tf.margin_top=0; tf.margin_bottom=0
    for i,para in enumerate(runs):
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph()
        p.alignment=align; p.space_after=Pt(sa); p.space_before=Pt(0); p.line_spacing=ls
        for (t,sz,b,c,*rest) in para:
            it=rest[0] if rest else False
            r=p.add_run(); r.text=t; r.font.size=Pt(sz); r.font.bold=b; r.font.italic=it; r.font.name=FONT; r.font.color.rgb=c
    return tf
def fillrect(s,x,y,w,h,color):
    sp=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,x,y,w,h); sp.fill.solid(); sp.fill.fore_color.rgb=color
    sp.line.fill.background(); sp.shadow.inherit=False; return sp
def rrect(s,x,y,w,h,fill,line=LINEC,lw=1.0,radius=0.08):
    sp=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,x,y,w,h); sp.fill.solid(); sp.fill.fore_color.rgb=fill
    if line is None: sp.line.fill.background()
    else: sp.line.color.rgb=line; sp.line.width=Pt(lw)
    sp.shadow.inherit=False
    try: sp.adjustments[0]=radius
    except Exception: pass
    return sp
def settext(sp,lines,anchor=MSO_ANCHOR.MIDDLE,align=PP_ALIGN.CENTER):
    tf=sp.text_frame; tf.word_wrap=True; tf.vertical_anchor=anchor
    tf.margin_left=Inches(0.08); tf.margin_right=Inches(0.08); tf.margin_top=Inches(0.03); tf.margin_bottom=Inches(0.03)
    for i,(t,sz,b,c) in enumerate(lines):
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph()
        p.alignment=align; p.space_after=Pt(1); p.line_spacing=1.0
        r=p.add_run(); r.text=t; r.font.size=Pt(sz); r.font.bold=b; r.font.name=FONT; r.font.color.rgb=c
def box(s,x,y,w,h,head,sub=None,fill=LIGHT,headc=NAVY,hs=12.5,ss=10,subc=GRAYTX,align=PP_ALIGN.CENTER):
    sp=rrect(s,x,y,w,h,fill); lines=[(head,hs,True,headc)]
    if sub: lines.append((sub,ss,False,subc))
    settext(sp,lines,align=align); return sp
def arrow(s,x,y,w=Inches(0.34),h=Inches(0.3),color=TEAL,shape=MSO_SHAPE.RIGHT_ARROW):
    sp=s.shapes.add_shape(shape,x,y,w,h); sp.fill.solid(); sp.fill.fore_color.rgb=color; sp.line.fill.background(); sp.shadow.inherit=False; return sp
def connect(s,x1,y1,x2,y2,color=LINEC,w=1.4):
    ln=s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,x1,y1,x2,y2); ln.line.color.rgb=color; ln.line.width=Pt(w); return ln
def title(s,text,num,kicker=None):
    if kicker: _txt(s,Inches(0.7),Inches(0.32),Inches(11.9),Inches(0.3),[[(kicker.upper(),11,True,TEAL)]])
    _txt(s,Inches(0.7),Inches(0.56),Inches(11.9),Inches(0.7),[[(text,24,True,NAVY)]])
    fillrect(s,Inches(0.72),Inches(1.16),Inches(1.5),Pt(3),TEAL)
    _txt(s,SW-Inches(1.1),SH-Inches(0.42),Inches(0.7),Inches(0.3),[[(str(num),10,False,MUTE)]],align=PP_ALIGN.RIGHT)
    _txt(s,Inches(0.7),SH-Inches(0.42),Inches(9),Inches(0.3),[[("Practical LLMOps — research phase",9,False,MUTE)]])
def bullets(s,items,x=Inches(0.75),y=Inches(1.45),w=Inches(11.8),h=Inches(5.2),size=15,gap=9,color=GRAYTX):
    tf=s.shapes.add_textbox(x,y,w,h).text_frame; tf.word_wrap=True; tf.margin_left=0; tf.margin_top=0
    for i,it in enumerate(items):
        text=it[0]; lvl=it[1] if len(it)>1 else 0; bold=it[2] if len(it)>2 else False; col=it[3] if len(it)>3 else color
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph()
        p.space_after=Pt(gap); p.space_before=Pt(0); p.line_spacing=1.05; p.level=lvl
        r=p.add_run(); r.text=("•  " if lvl==0 else "–  ")+text
        r.font.size=Pt(size-lvl); r.font.bold=bold; r.font.name=FONT; r.font.color.rgb=col
    return tf
def table(s,x,y,w,headers,rows,cw,fs=11,hs=11,rh=Inches(0.5),wrap_col0_bold=True):
    nr=len(rows)+1; t=s.shapes.add_table(nr,len(headers),x,y,w,rh*nr).table
    for i,cc in enumerate(cw): t.columns[i].width=cc
    for j,hh in enumerate(headers):
        c=t.cell(0,j); c.fill.solid(); c.fill.fore_color.rgb=NAVY; c.vertical_anchor=MSO_ANCHOR.MIDDLE
        c.margin_left=Inches(0.08); c.margin_top=Inches(0.02); c.margin_bottom=Inches(0.02)
        r=c.text_frame.paragraphs[0].add_run(); r.text=hh; r.font.size=Pt(hs); r.font.bold=True; r.font.color.rgb=WHITE; r.font.name=FONT
    for i,row in enumerate(rows):
        for j,v in enumerate(row):
            c=t.cell(i+1,j); c.fill.solid(); c.fill.fore_color.rgb=WHITE if i%2==0 else LIGHT; c.vertical_anchor=MSO_ANCHOR.MIDDLE
            c.margin_left=Inches(0.08); c.margin_top=Inches(0.02); c.margin_bottom=Inches(0.02)
            r=c.text_frame.paragraphs[0].add_run(); r.text=v; r.font.size=Pt(fs); r.font.name=FONT; r.font.color.rgb=GRAYTX
            if j==0 and wrap_col0_bold: r.font.bold=True; r.font.color.rgb=NAVY
    return t
def notes(s,text): s.notes_slide.notes_text_frame.text=text
def levelbar(s,x,y,w,levels,activeidx=None):
    """levels: list of (label). Draws 4 connected segments."""
    n=len(levels); seg=Inches((w.inches-0.3*(n-1))/n) if hasattr(w,'inches') else w
    colors=[SLATE,TEAL,BLUE,NAVY]
    xx=x
    for i,lb in enumerate(levels):
        c=colors[i%len(colors)]
        b=rrect(s,xx,y,seg,Inches(0.55),c,line=None); settext(b,[(lb,12,True,WHITE)])
        if i<n-1: connect(s,xx+seg,y+Inches(0.275),xx+seg+Inches(0.3),y+Inches(0.275),color=LINEC,w=2)
        xx=xx+seg+Inches(0.3)

print("Building practical LLMOps research deck...")

# ============================================================ 1 TITLE
s=add()
fillrect(s,0,0,Inches(0.28),SH,NAVY)
_txt(s,Inches(0.9),Inches(1.9),Inches(11.5),Inches(1.3),[[("Practical LLMOps",34,True,NAVY)]])
fillrect(s,Inches(0.94),Inches(2.75),Inches(2.0),Pt(3),TEAL)
_txt(s,Inches(0.92),Inches(3.05),Inches(11.3),Inches(1.5),
     [[("What it means, what it covers, and exactly how to build it — component by component, on Azure.",17,False,GRAYTX)],
      [("LLM = Large Language Model. LLMOps = Large Language Model Operations. Research phase — implementation depth, no fluff.",13,False,MUTE,True)]],sa=10)
_txt(s,Inches(0.92),Inches(6.35),Inches(11),Inches(0.5),[[("Working draft  ·  editable  ·  research phase",12,True,NAVY)]])
notes(s,"Set the tone: this is a working, practical session — not a sales pitch. We'll go through each component of LLMOps, compare the real tool options, and show the exact setup, from a small start up to a full production platform.")

# ============================================================ 2 WHAT LLMOPS COVERS
s=add(); title(s,"What LLMOps covers",2,"Fundamentals")
_txt(s,Inches(0.7),Inches(1.3),Inches(11.9),Inches(0.4),
     [[("LLMOps = DevOps (build/test/release discipline) adapted for language-model apps. It is a set of components, not one tool.",13,False,GRAYTX,True)]])
comps=[("Source control & CI/CD","GitHub"),("Prompt management","versioned, tested"),("Model management","catalog, aliases, router"),
       ("Evaluation","golden datasets, scoring"),("Observability","traces, cost, quality"),("Feedback & analytics","dashboards, improve loop"),
       ("Data pipelines","ingestion, RAG, warehouse"),("Guardrails & safety","input/output checks"),
       ("Serving & gateway","routing, quotas, cache"),("Multi-agent orchestration","frameworks, workflows"),
       ("Security & identity","access, secrets, network"),("FinOps","cost tracking & budgets")]
cw=Inches(2.9); ch=Inches(1.05); gx=Inches(0.16); gy=Inches(0.16); x0=Inches(0.7); y0=Inches(1.85)
cols=[TEAL,BLUE,NAVY,SLATE,DTEAL,PURP,SKY,AMBER,GREEN,ROSE,GRAYTX,MUTE]
for i,((h,sub),c) in enumerate(zip(comps,cols)):
    r=i//4; cidx=i%4; x=x0+cidx*(cw+gx); y=y0+r*(ch+gy)
    box(s,x,y,cw,ch,h,sub,fill=LIGHT)
    fillrect(s,x,y,Inches(0.09),ch,c)
notes(s,"Twelve components. Not all need building on day one — that's the whole point of the maturity-level plan later. But this is the complete list so nothing gets missed.")

# ============================================================ 3 REUSABILITY
s=add(); title(s,"What's reusable, and what's per use case",3,"Fundamentals")
table(s,Inches(0.6),Inches(1.45),Inches(12.13),
      ["Component","Built once (shared)","Done again per use case"],
      [["Source control & CI/CD pipeline","Yes — one pipeline design","New folder + config"],
       ["Model gateway & catalog","Yes","Pick an alias"],
       ["Evaluation engine & CI gate","Yes","New golden dataset"],
       ["Observability & dashboards","Yes","New tags/filters"],
       ["Data pipeline framework","Yes","New source connectors"],
       ["Guardrail engine","Yes","Policy tuning"],
       ["Agent orchestration runtime","Yes","New agents & prompts"],
       ["Prompts, agents, golden data","No — always new","Always new"]],
      cw=[Inches(4.0),Inches(3.9),Inches(4.23)],fs=12,rh=Inches(0.55))
_txt(s,Inches(0.6),Inches(6.15),Inches(12.13),Inches(0.7),
     [[("This is why we build the platform first: ",13,True,NAVY),("most rows say 'Yes' — build once, reuse for every future use case.",13,False,GRAYTX)]])
notes(s,"The core argument. Most of the machinery is shared. Only the content — prompts, agents, and the golden test data — is written fresh each time. That's what makes a platform pay off after the first use case.")

# ============================================================ 4 MATURITY LEVELS
s=add(); title(s,"Start small, grow: four levels",4,"Fundamentals")
levels=[("Level 0","Baseline","Weeks 1-2","Repo, Azure landing, basic tracing"),
        ("Level 1","Managed","Weeks 3-6","CI evals, prompt registry, RAG, guardrails"),
        ("Level 2","Production-grade","Months 2-4","Full evals, canary, feedback loop, cost control"),
        ("Level 3","Scaled","Months 4+","Self-service onboarding, fine-tuning, multi-team agents")]
cw=Inches(2.95); gap=Inches(0.15); x=Inches(0.6); y=Inches(1.7)
cols=[SLATE,TEAL,BLUE,NAVY]
for i,(lvl,name,dur,desc) in enumerate(levels):
    c=cols[i]
    rrect(s,x,y,cw,Inches(0.6),c,line=None); settext(rrect(s,x,y,0,0,c) if False else s.shapes[-1],[(lvl,14,True,WHITE),(name,11,False,WHITE)]) if False else None
    hd=s.shapes[-1]; settext(hd,[(lvl+" · "+name,13,True,WHITE)])
    rrect(s,x,y+Inches(0.66),cw,Inches(0.4),LIGHT); settext(s.shapes[-1],[(dur,11.5,True,NAVY)])
    box(s,x,y+Inches(1.12),cw,Inches(1.6),"",desc,fill=WHITE)
    fillrect(s,x,y+Inches(1.12),cw,Inches(1.6),LIGHT) if False else None
    if i<3: arrow(s,x+cw+Inches(0.01),y+Inches(0.3),w=Inches(0.13),h=Inches(0.22))
    x+=cw+gap
_txt(s,Inches(0.6),Inches(4.6),Inches(12.13),Inches(0.6),
     [[("Each level adds capability without rework — nothing built at Level 0 is thrown away later.",14,True,NAVY,True)]])
notes(s,"Four levels. Level 0-1 gets one real use case live in about six weeks on a genuine (if minimal) LLMOps loop. Level 2 hardens it for production scale. Level 3 is when the platform serves many teams. This structure repeats through the deck.")

# ============================================================ 5 GITHUB REPO LAYOUT
s=add(); title(s,"The Ops backbone: repository layout on GitHub",5,"CI/CD")
items=[("/prompts","versioned prompt templates",TEAL),("/agents","agent & workflow definitions",BLUE),
       ("/evals","golden datasets + evaluator configs",AMBER),("/src","application & orchestration code",SLATE),
       ("/pipelines","GitHub Actions workflows",NAVY),("/infra","Bicep/Terraform (infrastructure as code)",PURP),
       ("/dashboards","dashboard & alert definitions",DTEAL)]
y=Inches(1.55); h=Inches(0.66)
for name,desc,c in items:
    rrect(s,Inches(1.1),y,Inches(2.6),h-Inches(0.08),c,line=None); settext(s.shapes[-1],[(name,13,True,WHITE)])
    rrect(s,Inches(3.85),y,Inches(8.4),h-Inches(0.08),LIGHT)
    _txt(s,Inches(4.05),y,Inches(8.0),h-Inches(0.08),[[(desc,12.5,False,GRAYTX)]],anchor=MSO_ANCHOR.MIDDLE)
    y+=h
_txt(s,Inches(1.1),y+Inches(0.1),Inches(11.15),Inches(0.7),
     [[("One monorepo for the platform; one subfolder per use case inside /prompts, /agents, /evals. ",13,True,NAVY),
       ("Everything here is reviewed by pull request — nothing changes untracked.",13,False,GRAYTX)]])
notes(s,"This is the literal folder structure in the Git repository. Point out: prompts and agents are code, not portal edits. Golden datasets live next to the code they test.")

# ============================================================ 6 CI/CD PIPELINE WITH GATES
s=add(); title(s,"GitHub Actions: from change to production",6,"CI/CD")
steps=[("Pull request","edit prompt/agent"),("pr-checks.yml","lint + quick eval subset"),
       ("Merge to main","eval-full.yml runs nightly"),("deploy.yml","build & push to dev"),
       ("Canary in prod","small traffic slice"),("Full rollout","or auto-rollback")]
x=Inches(0.6); y=Inches(1.6); bw=Inches(1.95); bh=Inches(1.3); gap=Inches(0.15)
for i,(h,sub) in enumerate(steps):
    fill=RGBColor(0xF6,0xE9,0xE0) if i in (1,2) else LIGHT
    box(s,x,y,bw,bh,h,sub,fill=fill,hs=12)
    if i<len(steps)-1: arrow(s,x+bw-Inches(0.01),y+bh/2-Inches(0.15),w=Inches(0.22))
    x+=bw+gap
_txt(s,Inches(0.6),Inches(3.2),Inches(12),Inches(0.4),[[("Required GitHub Environments: ",13,True,NAVY),("dev → test → prod, each with required reviewers before promotion.",13,False,GRAYTX)]])
table(s,Inches(0.6),Inches(3.75),Inches(12.13),
      ["Safeguard","How it works"],
      [["No stored cloud keys","GitHub OIDC (federated login) signs in to Azure per run"],
       ["Quality gate","pr-checks.yml fails the PR if eval score drops past a threshold"],
       ["Gradual release","canary slice first, health + eval checks, then full rollout"],
       ["Auto-rollback","deployment reverts automatically if alarms fire"]],
      cw=[Inches(3.0),Inches(9.13)],fs=12,rh=Inches(0.5))
notes(s,"Walk left to right. The amber boxes are where evaluation happens automatically. OIDC = the app authenticates to Azure without any stored password or key, using a short-lived federated token. This whole slide is 'what makes it LLMOps and not just DevOps' — the evaluation gate.")

# ============================================================ 7 PROMPT MANAGEMENT
s=add(); title(s,"Prompt management: where prompts actually live",7,"Prompts")
box(s,Inches(0.6),Inches(1.5),Inches(5.7),Inches(2.2),"Git (source of truth)","one YAML file per prompt: id, version, template, variables, eval refs, changelog",fill=LIGHT,hs=14)
arrow(s,Inches(6.35),Inches(2.5),w=Inches(0.4))
box(s,Inches(6.85),Inches(1.5),Inches(5.7),Inches(2.2),"Runtime registry","labelled prod / staging; app asks for \"prompt X, label prod\"; supports hot-swap and A/B",fill=RGBColor(0xE9,0xF3,0xF0),headc=TEAL,hs=14)
table(s,Inches(0.6),Inches(4.05),Inches(12.13),
      ["Registry option","What it is","Good for"],
      [["Langfuse Prompt Management","Open-source, self-hosted","Versioned prompts + built-in observability together"],
       ["Foundry prompt assets","Native to the Azure AI platform","Staying inside one Microsoft portal"],
       ["Git + in-app cache","Simplest option","Small setups; no separate registry service"]],
      cw=[Inches(3.3),Inches(4.0),Inches(4.83)],fs=11.5,rh=Inches(0.55))
_txt(s,Inches(0.6),Inches(6.15),Inches(12.13),Inches(0.6),
     [[("Anti-pattern: ",13,True,ROSE),("prompts typed directly into a chatbot portal or hard-coded in application code — untracked, untested, unreviewed.",13,False,GRAYTX)]])
notes(s,"Prompts are treated exactly like code: they live in Git, go through pull request review, and are tested by the evaluation gate before anyone sees them in production. The registry is just a fast-access, labelled copy for the running app.")

# ============================================================ 8 MODEL MANAGEMENT
s=add(); title(s,"Model management: how model choice works",8,"Models")
box(s,Inches(0.6),Inches(1.5),Inches(3.7),Inches(1.5),"App code","asks for a task alias, e.g. \"summarize\"",fill=LIGHT,hs=13)
arrow(s,Inches(4.45),Inches(2.15))
box(s,Inches(4.9),Inches(1.5),Inches(3.7),Inches(1.5),"models.yaml (in Git)","alias -> deployment, per environment",fill=RGBColor(0xE9,0xF3,0xF0),headc=TEAL,hs=13)
arrow(s,Inches(8.75),Inches(2.15))
box(s,Inches(9.2),Inches(1.5),Inches(3.35),Inches(1.5),"Model deployment","the actual model version running",fill=LIGHT,hs=13)
_txt(s,Inches(0.6),Inches(3.3),Inches(12.13),Inches(0.4),
     [[("Swapping a model = a config file change reviewed like any other code, and it must pass the evaluation gate.",13,True,NAVY,True)]])
table(s,Inches(0.6),Inches(3.85),Inches(12.13),
      ["Task type","Typical choice","Why"],
      [["Complex reasoning","Larger flagship model","Higher accuracy on hard multi-step tasks"],
       ["Bulk / simple tasks","Smaller, cheaper model","Good enough quality, much lower cost"],
       ["Real-time voice","Realtime speech-to-speech model","Built for low-delay spoken conversation"],
       ["Search / retrieval","Embedding model","Turns text into searchable vectors, not chat"]],
      cw=[Inches(2.8),Inches(3.8),Inches(5.53)],fs=11.5,rh=Inches(0.55))
notes(s,"No model name is hard-coded in the application. The app asks for a task alias; a config file (in Git, reviewed) decides which model that maps to today. A Model Router can also pick automatically by cost and quality, but explicit aliases give more control.")

# ============================================================ 9 GOLDEN DATASETS
s=add(); title(s,"Evaluation: what a golden dataset actually is",9,"Evaluation")
_txt(s,Inches(0.6),Inches(1.35),Inches(12.13),Inches(0.4),
     [[("A golden dataset is a saved, versioned set of test questions with the expected answer or grading rule — used to check every change before it ships.",13,False,GRAYTX,True)]])
srcs=[("SME-authored","experts write realistic cases",TEAL),("Mined from real traffic","anonymised past requests",BLUE),("Synthetic + reviewed","generated, then human-checked",AMBER)]
x=Inches(0.6); y=Inches(2.0)
for h,sub,c in srcs:
    box(s,x,y,Inches(3.9),Inches(1.15),h,sub,fill=LIGHT); fillrect(s,x,y,Inches(0.09),Inches(1.15),c); x+=Inches(4.08)
table(s,Inches(0.6),Inches(3.4),Inches(12.13),
      ["Metric group","Example metric","How it's scored"],
      [["Knowledge answers (RAG)","Groundedness — is it backed by real data","Automated grading model + rules"],
       ["Writing quality","Coherence, tone","Automated grading model"],
       ["Exact tasks","Exact match / F1 score","Rule-based comparison"],
       ["Agent behaviour","Task success rate, correct tool used","Rule-based + automated grading"],
       ["Running cost & speed","Latency, cost per request","Measured directly"],
       ["Safety","Unsafe content rate","Automated safety scanner + red-team tests"]],
      cw=[Inches(2.9),Inches(4.2),Inches(5.03)],fs=11,rh=Inches(0.5))
_txt(s,Inches(0.6),Inches(6.6),Inches(12.13),Inches(0.4),[[("Start with 50–200 test cases per use case; grow it from real feedback over time.",12.5,False,MUTE,True)]])
notes(s,"Golden dataset = the exam the system must keep passing. Three ways to build it. Then the metrics table — different use cases need different measures, and most are graded automatically, with a human check on a sample.")

# ============================================================ 10 OBSERVABILITY PROVIDERS
s=add(); title(s,"Observability: comparing the providers",10,"Observability")
table(s,Inches(0.6),Inches(1.45),Inches(12.13),
      ["Provider","Type","What it tracks","Best fit"],
      [["Azure Monitor + App Insights","Cloud-native (Microsoft)","Traces, tokens, latency, errors, custom events","Default; data stays in the tenant"],
       ["Langfuse","Open source (self-hosted)","Traces, cost per model, prompt versions, user feedback, datasets","Best LLM-specific view; also does prompt management"],
       ["LangSmith","SaaS (LangChain)","Traces, feedback, datasets, prompt hub","Teams already using LangChain tools"],
       ["Arize Phoenix","Open source","Traces, evaluation scores, embedding drift","Watching for data/answer drift over time"],
       ["Datadog LLM Observability","SaaS","Traces, evaluation, cost","Already a Datadog shop"]],
      cw=[Inches(2.6),Inches(2.1),Inches(4.13),Inches(3.3)],fs=10.5,rh=Inches(0.62))
box(s,Inches(0.6),Inches(5.35),Inches(12.13),Inches(1.2),"Recommended pairing",
    "Azure Monitor + Application Insights as the system of record  +  self-hosted Langfuse as the LLM-specific lens (cost per model, prompt versions, feedback)",
    fill=RGBColor(0xE9,0xF3,0xF0),headc=TEAL,hs=13,ss=12)
notes(s,"Five real options, honestly compared. The recommendation is a pairing, not a single tool: Azure-native for the system of record, plus Langfuse for the LLM-specific view that Azure Monitor doesn't give out of the box.")

# ============================================================ 11 WHAT GETS TRACKED
s=add(); title(s,"What exactly gets tracked on every request",11,"Observability")
items=[("Model calls","prompt, answer, tokens, cost, latency, model + prompt version",TEAL),
       ("Tool calls","which tool, inputs, outputs, success/failure",BLUE),
       ("Agent hops","which agent ran, in what order, handoffs",SLATE),
       ("Sessions & users","conversation id, user id (no raw personal data)",AMBER),
       ("Feedback events","thumbs up/down, reason, edits, escalations",DTEAL)]
y=Inches(1.6)
for h,sub,c in items:
    rrect(s,Inches(0.8),y,Inches(3.0),Inches(0.85),c,line=None); settext(s.shapes[-1],[(h,13,True,WHITE)])
    rrect(s,Inches(4.0),y,Inches(8.7),Inches(0.85),LIGHT); _txt(s,Inches(4.2),y,Inches(8.3),Inches(0.85),[[(sub,12.5,False,GRAYTX)]],anchor=MSO_ANCHOR.MIDDLE)
    y+=Inches(1.0)
_txt(s,Inches(0.8),y+Inches(0.05),Inches(11.5),Inches(0.5),[[("Every one of these carries a trace ID, so a bad answer can be followed back to its exact request.",13,True,NAVY,True)]])
notes(s,"This is the checklist for 'did we instrument this properly'. If any of these five rows is missing, we can't fully diagnose an issue or measure quality in production.")

# ============================================================ 12 FEEDBACK & ANALYTICS LOOP
s=add(); title(s,"Feedback, analytics, and the improvement loop",12,"Feedback")
steps=["Capture\nfeedback","Land in\ndashboards","Triage\nnegatives","Add to\ngolden set","Fix &\nre-evaluate","Ship"]
x=Inches(0.6); y=Inches(1.6); bw=Inches(1.9); bh=Inches(1.3); gap=Inches(0.15)
for i,st in enumerate(steps):
    head=st.split("\n")[0]; sub=st.split("\n")[1] if "\n" in st else ""
    box(s,x,y,bw,bh,head,sub,fill=LIGHT)
    if i<len(steps)-1: arrow(s,x+bw-Inches(0.01),y+bh/2-Inches(0.15),w=Inches(0.2))
    x+=bw+gap
arrow(s,Inches(0.8),Inches(3.25),w=Inches(10.9),h=Inches(0.2),color=RGBColor(0xBF,0xD6,0xD2),shape=MSO_SHAPE.LEFT_ARROW)
table(s,Inches(0.6),Inches(3.7),Inches(12.13),
      ["Dashboard tile","Shows"],
      [["Volume & containment","how many requests, how many resolved without a human"],
       ["Speed & cost","p95 latency, cost per use case per day"],
       ["Quality trend","evaluation scores over time"],
       ["Feedback rate","thumbs up/down rate and top negative reasons"]],
      cw=[Inches(3.5),Inches(8.63)],fs=12,rh=Inches(0.5))
_txt(s,Inches(0.6),Inches(6.1),Inches(12.13),Inches(0.5),
     [[("Later stage: ",13,True,NAVY),("once prompts and search stop improving quality further, well-reviewed accepted answers can become fine-tuning data for a smaller, cheaper model.",13,False,GRAYTX)]])
notes(s,"The loop that makes this an operation, not a one-off build. Point at the light-teal loop-back arrow — that is the whole idea of LLMOps in one picture. Fine-tuning is mentioned as a later, optional step, not a starting point.")

# ============================================================ 13 DATA PIPELINES / RAG
s=add(); title(s,"Data pipelines: feeding answers from our own data",13,"Data")
_txt(s,Inches(0.6),Inches(1.3),Inches(12.13),Inches(0.4),
     [[("RAG (Retrieval-Augmented Generation): search our documents first, then let the model answer using only what was found.",13,False,GRAYTX,True)]])
steps=[("Sources","files, records,\nsystems"),("Ingest","read & clean,\nremove personal data"),("Chunk & embed","split text,\nturn into vectors"),
       ("Search index","store for\nmeaning search"),("Retrieve","find matching\npieces at answer time")]
x=Inches(0.6); y=Inches(2.05); bw=Inches(2.25); bh=Inches(1.5); gap=Inches(0.18)
for i,(h,sub) in enumerate(steps):
    box(s,x,y,bw,bh,h,sub.replace("\n"," "),fill=LIGHT)
    if i<len(steps)-1: arrow(s,x+bw-Inches(0.01),y+bh/2-Inches(0.15))
    x+=bw+gap
_txt(s,Inches(0.6),Inches(3.9),Inches(12.13),Inches(0.4),[[("Refresh: ",13,True,NAVY),("on a schedule, or immediately when a source changes (change data capture).",13,False,GRAYTX)]])
notes(s,"Five steps, left to right. This is the pipeline behind every RAG answer. Refresh can run on a timer or trigger the moment a source document changes.")

# ============================================================ 14 WAREHOUSE RELATIONSHIP
s=add(); title(s,"Does this replace the data warehouse? No — it connects to it",14,"Data")
box(s,Inches(0.9),Inches(1.7),Inches(3.6),Inches(1.6),"Source systems","files, databases, CRM, call records",fill=LIGHT)
arrow(s,Inches(4.65),Inches(2.35))
box(s,Inches(5.15),Inches(1.7),Inches(3.6),Inches(1.6),"Lakehouse / warehouse","one place data lands and is governed",fill=RGBColor(0xE9,0xF3,0xF0),headc=TEAL)
arrow(s,Inches(8.9),Inches(2.35))
box(s,Inches(9.4),Inches(1.7),Inches(3.1),Inches(1.6),"Two consumers","search index  +  reporting dashboards",fill=LIGHT)
_txt(s,Inches(0.9),Inches(3.7),Inches(12.13),Inches(0.4),[[("The lakehouse can be the same platform used for general company reporting — LLM apps are just another consumer of it, and another contributor to it (via captured telemetry).",13,False,GRAYTX,True)]])
table(s,Inches(0.9),Inches(4.3),Inches(11.5),
      ["Flow","What moves"],
      [["Into the lakehouse","source documents, plus captured request/response telemetry"],
       ["Out of the lakehouse","the search index (for RAG) and curated training data (for fine-tuning)"],
       ["Governance","data classification and lineage tracked centrally, one place"]],
      cw=[Inches(3.0),Inches(8.5)],fs=12,rh=Inches(0.55))
notes(s,"Common question, so it gets its own slide. The lakehouse is not replaced — it becomes the shared landing place, feeding the search index on one side and reporting/training data on the other.")

# ============================================================ 15 MULTI-AGENT FRAMEWORKS
s=add(); title(s,"Multi-agent frameworks: comparing the options",15,"Multi-agent")
table(s,Inches(0.6),Inches(1.45),Inches(12.13),
      ["Framework","Style","Best fit"],
      [["Microsoft Agent Framework","Combines two earlier Microsoft toolkits; built for durable, long-running workflows","Azure-native builds; enterprise reliability needs"],
       ["LangGraph","Explicit graph / state-machine control","Fine-grained control over each step"],
       ["CrewAI","Agents given roles, like a small team","Fast prototyping"],
       ["OpenAI Agents SDK","Lightweight, minimal setup","Simple agent apps tied to one model provider"]],
      cw=[Inches(3.0),Inches(5.13),Inches(4.0)],fs=12,rh=Inches(0.65))
_txt(s,Inches(0.6),Inches(4.4),Inches(12.13),Inches(0.4),[[("Recommended default: the Microsoft Agent Framework, hosted on the platform's managed agent service.",13,True,NAVY,True)]])
box(s,Inches(0.6),Inches(5.0),Inches(12.13),Inches(1.4),"Agent-to-agent (A2A)","a standard that lets an agent built by one team call an agent built by another team, even on a different framework",fill=LIGHT,hs=13,ss=12)
notes(s,"Four real frameworks, compared honestly. The default recommendation is the Microsoft one because it's built into the same Azure platform and supports long-running, resumable work. A2A is the standard that lets agents from different teams cooperate later.")

# ============================================================ 16 AGENT TYPES & PATTERNS
s=add(); title(s,"Agent types, and how they're coordinated",16,"Multi-agent")
types=[("Router","decides what the request needs"),("Planner","breaks work into steps"),("Retrieval","searches documents"),
       ("Tool / action","calls outside systems"),("Critic","checks the draft answer"),("Guardrail","checks policy & safety"),
       ("Summarizer","wraps up the result"),("Human-proxy","asks a person to approve")]
cw=Inches(2.92); ch=Inches(0.85); x0=Inches(0.6); y0=Inches(1.55)
for i,(h,sub) in enumerate(types):
    r=i//4; c=i%4; x=x0+c*(cw+Inches(0.16)); y=y0+r*(ch+Inches(0.14))
    box(s,x,y,cw,ch,h,sub,fill=LIGHT,hs=12,ss=10)
_txt(s,Inches(0.6),Inches(3.65),Inches(12.13),Inches(0.35),[[("Coordination patterns:",13,True,NAVY)]])
pats=[("Sequential","A -> B -> C"),("Concurrent","A -> [B,C] -> merge"),("Group chat","A <-> B <-> C"),("Handoff","A => specialist"),("Planner-led","Plan -> workers")]
x=Inches(0.6); y=Inches(4.05)
for name,flow in pats:
    b=rrect(s,x,y,Inches(2.28),Inches(0.95),LIGHT); settext(b,[(name,12,True,NAVY),(flow,11,False,TEAL)]); x+=Inches(2.4)
_txt(s,Inches(0.6),Inches(5.25),Inches(12.13),Inches(0.6),
     [[("Agents are defined as code (YAML) in the same Git repository, tested by the same evaluation gate, and traced the same way as everything else.",13,True,NAVY,True)]])
notes(s,"Eight common agent roles, then five coordination patterns. The last line is the key point of this whole section: multi-agent does not get a separate Ops process — it plugs into the exact same Git, evaluation, and tracing setup as everything else.")

# ============================================================ 17 AZURE STACK TOUR
s=add(); title(s,"The Azure pieces, mapped to what they do",17,"Azure stack")
table(s,Inches(0.6),Inches(1.4),Inches(12.13),
      ["Azure service","What it is for"],
      [["Azure AI Foundry","the portal: model catalog, agents, evaluation, tracing"],
       ["Foundry Agent Service","runs and manages agents reliably (hosted, resumable)"],
       ["Azure OpenAI","the language and voice models themselves"],
       ["Model Router","auto-picks the cheapest model that meets a quality bar"],
       ["Azure AI Search","the search index behind RAG"],
       ["Azure AI Document Intelligence","reads and extracts text from files"],
       ["Content Safety","checks inputs and outputs for unsafe content and personal data"],
       ["Microsoft Fabric / OneLake","the lakehouse / data warehouse layer"],
       ["Azure Cosmos DB","fast storage for agent memory and state"],
       ["API Management","the gateway: one door in, usage limits, logging"],
       ["Entra ID / Key Vault","identity and secret storage"],
       ["Azure Monitor / App Insights","logs, traces, dashboards, alerts"]],
      cw=[Inches(3.4),Inches(8.73)],fs=11,rh=Inches(0.44))
notes(s,"A reference list. Don't read every row aloud — this is here so people can see the exact Azure service behind every component we've discussed, in one place.")

# ============================================================ 18 END TO END ARCHITECTURE
s=add(); title(s,"Putting it together: one end-to-end picture",18,"Architecture")
# GitHub column
box(s,Inches(0.5),Inches(1.5),Inches(2.5),Inches(2.4),"GitHub","prompts, agents,\nevals, infra as code",fill=LIGHT).text_frame
_txt(s,Inches(0.5),Inches(4.0),Inches(2.5),Inches(0.5),[[("CI/CD pipeline",11.5,True,TEAL)]],align=PP_ALIGN.CENTER)
arrow(s,Inches(3.05),Inches(2.6))
# Azure runtime column
box(s,Inches(3.55),Inches(1.5),Inches(3.1),Inches(2.4),"Azure runtime","gateway, orchestrator,\nagents, models, search,\nguardrails, memory",fill=RGBColor(0xE9,0xF3,0xF0),headc=TEAL)
arrow(s,Inches(6.7),Inches(2.6))
# Telemetry column
box(s,Inches(7.2),Inches(1.5),Inches(2.6),Inches(2.4),"Telemetry","traces, cost,\nquality scores",fill=LIGHT)
arrow(s,Inches(9.85),Inches(2.6))
# Lakehouse + dashboards
box(s,Inches(10.35),Inches(1.5),Inches(2.48),Inches(2.4),"Lakehouse","dashboards +\ntraining data",fill=LIGHT)
# feedback loop back
arrow(s,Inches(9.7),Inches(4.25),w=Inches(0.3),h=Inches(0.22),shape=MSO_SHAPE.LEFT_ARROW,color=RGBColor(0xBF,0xD6,0xD2))
connect(s,Inches(11.6),Inches(3.9),Inches(11.6),Inches(4.36),color=RGBColor(0xBF,0xD6,0xD2),w=2)
connect(s,Inches(11.6),Inches(4.36),Inches(0.6),Inches(4.36),color=RGBColor(0xBF,0xD6,0xD2),w=2)
connect(s,Inches(0.6),Inches(4.36),Inches(0.6),Inches(3.9),color=RGBColor(0xBF,0xD6,0xD2),w=2)
_txt(s,Inches(2.0),Inches(4.16),Inches(9.4),Inches(0.3),[[("Feedback flows back into golden datasets in GitHub — the loop closes.",12,True,TEAL,True)]],align=PP_ALIGN.CENTER)
table(s,Inches(0.6),Inches(4.85),Inches(12.13),
      ["Flow","What happens"],
      [["1. Change","a prompt/agent edit is reviewed and evaluated before merge"],
       ["2. Request","a user request flows through the gateway, orchestrator, and agents"],
       ["3. Telemetry","every step is traced, scored, and logged"],
       ["4. Feedback","real usage and ratings become new golden test cases"]],
      cw=[Inches(2.1),Inches(10.03)],fs=11,rh=Inches(0.38))
notes(s,"This is the single-picture summary of everything covered so far. Four flows: a change goes through review and testing; a live request runs through the system; everything is measured; and real usage feeds back into the test data. That closed loop is what LLMOps means end to end.")

# ============================================================ 19 PHASED ROADMAP DETAIL
s=add(); title(s,"The phased plan in detail",19,"Roadmap")
table(s,Inches(0.6),Inches(1.45),Inches(12.13),
      ["Level","Duration","What gets added","What you can now do"],
      [["0 — Baseline","2 weeks","Repo, Azure landing, basic tracing, one use case in a dev environment","Prove the wiring works end to end"],
       ["1 — Managed","4 weeks","CI evaluation gate, prompt registry, RAG search, guardrails, staged environments","Run the first use case safely in production"],
       ["2 — Production-grade","2-3 months","Full golden sets, canary releases, feedback dashboards, cost tracking, multi-agent orchestration","Scale one use case with confidence, add a second"],
       ["3 — Scaled","4+ months, ongoing","Self-service onboarding, model router, fine-tuning loop, cross-team agents","Onboard new use cases in weeks, not months"]],
      cw=[Inches(1.7),Inches(1.4),Inches(5.5),Inches(3.53)],fs=10.8,rh=Inches(0.85))
notes(s,"The same four levels, now with specifics. First production use case lands around week six. Each level is additive — we never redo earlier work.")

# ============================================================ 20 TEAM, TIMELINE, RISKS
s=add(); title(s,"Team, timeline, and risks to plan for",20,"Delivery")
box(s,Inches(0.6),Inches(1.5),Inches(5.8),Inches(2.4),"Core team",None,fill=LIGHT)
bullets(s,[("1 lead / architect",0),("2-3 engineers",0),("1 data engineer",0),("Part-time: security, delivery support",0)],
        x=Inches(0.85),y=Inches(2.0),w=Inches(5.3),h=Inches(1.8),size=13,gap=8)
box(s,Inches(6.7),Inches(1.5),Inches(6.03),Inches(2.4),"Assumptions",None,fill=LIGHT)
bullets(s,[("Cloud access & model quota already approved",0),("Data sources are reachable",0),("Experts give ~4 hours a week for test data",0)],
        x=Inches(6.95),y=Inches(2.0),w=Inches(5.5),h=Inches(1.8),size=13,gap=8)
table(s,Inches(0.6),Inches(4.15),Inches(12.13),
      ["Risk","Mitigation"],
      [["Cloud quota / capacity delays","Request early, in parallel with other setup"],
       ["Test-data quality (\"garbage in\")","Start small, review with experts, grow over time"],
       ["Cost surprises from multiple agent calls","Set budgets and alerts from day one"],
       ["New use cases pulled in before platform is ready","Agree the phased plan up front"]],
      cw=[Inches(4.5),Inches(7.63)],fs=11.5,rh=Inches(0.58))
notes(s,"Keep this brief. Small core team, clear assumptions, and the risks we already know about with a one-line fix for each.")

# ============================================================ 21 SUMMARY
s=add(); title(s,"Summary",21,"Wrap-up")
bullets(s,[
    ("LLMOps is twelve components, most of them reusable across every future use case.",0,True,NAVY),
    ("Multi-agent orchestration plugs into the same Git, evaluation, and tracing setup — no separate process.",0),
    ("Start small: one use case, a real (if minimal) loop, live in about six weeks.",0),
    ("Grow level by level; nothing built early gets thrown away later.",0),
    ("Every component maps to a specific, named Azure service.",0),
    ("Next step: agree the Level 0-1 scope and get the cloud environment set up.",0,True,NAVY),
],y=Inches(1.6),gap=13)
notes(s,"Recap. One clear ask at the end: agree the first six weeks of scope and get access to start.")

# ============================================================ 22 GLOSSARY
s=add(); title(s,"Terms in plain English",22,"Appendix")
gloss=[("LLM","Large Language Model — AI that reads and writes text."),
       ("LLMOps","Large Language Model Operations — the discipline of running LLM apps reliably."),
       ("RAG","Retrieval-Augmented Generation — answering using our own searched documents."),
       ("CI/CD","Continuous Integration / Continuous Delivery — automated testing and release."),
       ("Golden dataset","A saved set of test questions with expected answers, used to check quality."),
       ("Embeddings","Text turned into numbers so it can be searched by meaning."),
       ("OIDC","A way to log in to the cloud without storing a password or key."),
       ("Agent","An LLM that can take steps and use tools, not just chat."),
       ("Orchestrator","The part that decides which agent does what, and when."),
       ("A2A","Agent-to-Agent — a standard for agents to call each other."),
       ("MCP","Model Context Protocol — a standard way for agents to use tools."),
       ("Canary release","Releasing to a small slice of traffic first, to catch problems early."),
       ("Fine-tuning","Further training a model on our own reviewed examples."),
       ("Lakehouse","A combined data-storage and analytics platform (here: Microsoft Fabric)."),
       ("PII","Personally Identifiable Information — personal data like names or phone numbers."),
       ("Change data capture","Detecting and reacting to changes in a source system automatically.")]
colw=Inches(5.95); x0=Inches(0.6); y0=Inches(1.45); rowh=Inches(0.63)
for i,(term,defn) in enumerate(gloss):
    r=i//2; c=i%2; x=x0+c*(colw+Inches(0.25)); y=y0+r*rowh
    _txt(s,x,y,colw,rowh,[[(term+"  ",12,True,TEAL),(defn,11,False,GRAYTX)]],ls=1.0)
notes(s,"Appendix — leave this up during questions.")

out=os.path.join(os.path.dirname(os.path.dirname(os.path.abspath(__file__))),"presentation","LLMOps-Research-Deep-Dive.pptx")
prs.save(out)
print(f"Saved: {out}  ({len(prs.slides._sldIdLst)} slides)")
