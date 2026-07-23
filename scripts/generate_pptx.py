#!/usr/bin/env python3
"""Generate the Afni Enterprise LLMOps proposal PowerPoint deck.

Requires: python-pptx  (pip install python-pptx)
Output:   presentation/Afni-LLMOps-Proposal.pptx
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# ----------------------------------------------------------------------------
# Theme
# ----------------------------------------------------------------------------
NAVY   = RGBColor(0x12, 0x1F, 0x3D)   # deep navy - primary
INDIGO = RGBColor(0x1B, 0x3A, 0x6B)   # secondary
TEAL   = RGBColor(0x00, 0xA6, 0xA6)   # accent
CYAN   = RGBColor(0x2E, 0xC4, 0xD3)
AMBER  = RGBColor(0xF5, 0xA6, 0x23)   # highlight
LIGHT  = RGBColor(0xF4, 0xF6, 0xFA)   # light panel
GRAY   = RGBColor(0x5A, 0x64, 0x74)   # body gray
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
DARKTX = RGBColor(0x1C, 0x24, 0x33)
GREEN  = RGBColor(0x2E, 0x9E, 0x5B)

SW, SH = Inches(13.333), Inches(7.5)
FONT = "Segoe UI"

prs = Presentation()
prs.slide_width = SW
prs.slide_height = SH
BLANK = prs.slide_layouts[6]


def _set_fill(shape, color):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def rect(slide, x, y, w, h, color, line=None, shape=MSO_SHAPE.RECTANGLE):
    sp = slide.shapes.add_shape(shape, x, y, w, h)
    sp.fill.solid()
    sp.fill.fore_color.rgb = color
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line
        sp.line.width = Pt(1)
    sp.shadow.inherit = False
    return sp


def textbox(slide, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
            space_after=6, line_spacing=1.05):
    """runs: list of paragraphs; each paragraph = list of (text, size, bold, color, italic)."""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(space_after)
        p.space_before = Pt(0)
        p.line_spacing = line_spacing
        for (text, size, bold, color, *rest) in para:
            italic = rest[0] if rest else False
            r = p.add_run()
            r.text = text
            r.font.size = Pt(size)
            r.font.bold = bold
            r.font.italic = italic
            r.font.name = FONT
            r.font.color.rgb = color
    return tb


def bullets(slide, x, y, w, h, items, size=15, color=DARKTX, gap=8, lh=1.08):
    """items: list of (text, level, bold, color_override)."""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = 0
    tf.margin_top = 0
    for i, it in enumerate(items):
        text = it[0]
        level = it[1] if len(it) > 1 else 0
        bold = it[2] if len(it) > 2 else False
        col = it[3] if len(it) > 3 else color
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(gap)
        p.space_before = Pt(0)
        p.line_spacing = lh
        p.level = level
        bullet = "▸  " if level == 0 else "–  "
        r = p.add_run()
        r.text = bullet + text
        r.font.size = Pt(size - level * 1)
        r.font.bold = bold
        r.font.name = FONT
        r.font.color.rgb = col
    return tb


def footer(slide, page):
    rect(slide, 0, SH - Inches(0.32), SW, Inches(0.32), NAVY)
    textbox(slide, Inches(0.5), SH - Inches(0.33), Inches(9), Inches(0.3),
            [[("Evoke Technologies  ·  Enterprise LLMOps for Afni  ·  Confidential", 9, False, RGBColor(0xC7, 0xD0, 0xE0))]],
            anchor=MSO_ANCHOR.MIDDLE)
    textbox(slide, SW - Inches(1.3), SH - Inches(0.33), Inches(0.8), Inches(0.3),
            [[(str(page), 9, True, WHITE)]], align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)


def content_header(slide, kicker, title):
    rect(slide, 0, 0, SW, Inches(1.15), NAVY)
    rect(slide, 0, Inches(1.15), SW, Inches(0.06), TEAL)
    rect(slide, Inches(0.5), Inches(0.28), Inches(0.12), Inches(0.6), AMBER)
    textbox(slide, Inches(0.8), Inches(0.2), Inches(11.8), Inches(0.35),
            [[(kicker.upper(), 12, True, CYAN)]])
    textbox(slide, Inches(0.8), Inches(0.5), Inches(11.9), Inches(0.6),
            [[(title, 26, True, WHITE)]])


def new_slide():
    s = prs.slides.add_slide(BLANK)
    rect(s, 0, 0, SW, SH, WHITE)
    return s


def card(slide, x, y, w, h, title, body_items, accent=TEAL, title_size=14, body_size=12):
    rect(slide, x, y, w, h, LIGHT)
    rect(slide, x, y, w, Inches(0.09), accent)
    textbox(slide, x + Inches(0.18), y + Inches(0.16), w - Inches(0.36), Inches(0.5),
            [[(title, title_size, True, NAVY)]])
    bullets(slide, x + Inches(0.18), y + Inches(0.62), w - Inches(0.36), h - Inches(0.7),
            body_items, size=body_size, gap=5, lh=1.05)


def table(slide, x, y, w, headers, rows, col_widths=None, header_color=INDIGO,
          fsize=11, hsize=11, row_h=Inches(0.34)):
    nrows = len(rows) + 1
    ncols = len(headers)
    h = row_h * nrows
    gt = slide.shapes.add_table(nrows, ncols, x, y, w, h).table
    if col_widths:
        for i, cw in enumerate(col_widths):
            gt.columns[i].width = cw
    # header
    for j, htext in enumerate(headers):
        c = gt.cell(0, j)
        c.fill.solid(); c.fill.fore_color.rgb = header_color
        c.margin_left = Inches(0.08); c.margin_right = Inches(0.06)
        c.margin_top = Inches(0.02); c.margin_bottom = Inches(0.02)
        c.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = c.text_frame.paragraphs[0]
        r = p.add_run(); r.text = htext
        r.font.size = Pt(hsize); r.font.bold = True; r.font.color.rgb = WHITE; r.font.name = FONT
    # body
    for i, row in enumerate(rows):
        for j, val in enumerate(row):
            c = gt.cell(i + 1, j)
            c.fill.solid()
            c.fill.fore_color.rgb = WHITE if i % 2 == 0 else LIGHT
            c.margin_left = Inches(0.08); c.margin_right = Inches(0.06)
            c.margin_top = Inches(0.02); c.margin_bottom = Inches(0.02)
            c.vertical_anchor = MSO_ANCHOR.MIDDLE
            p = c.text_frame.paragraphs[0]
            r = p.add_run(); r.text = str(val)
            r.font.size = Pt(fsize); r.font.name = FONT; r.font.color.rgb = DARKTX
            if j == 0:
                r.font.bold = True; r.font.color.rgb = NAVY
    return gt


PAGE = [0]
def pg():
    PAGE[0] += 1
    return PAGE[0]

print("Building deck...")

# ============================================================ 1. TITLE
s = new_slide()
rect(s, 0, 0, SW, SH, NAVY)
rect(s, 0, 0, Inches(0.22), SH, TEAL)
rect(s, Inches(0.9), Inches(2.05), Inches(1.7), Inches(0.11), AMBER)
textbox(s, Inches(0.9), Inches(1.2), Inches(11), Inches(0.5),
        [[("ENTERPRISE PROPOSAL  ·  2026", 14, True, CYAN)]])
textbox(s, Inches(0.9), Inches(2.35), Inches(11.6), Inches(2.2),
        [[("Enterprise LLMOps Platform", 46, True, WHITE)],
         [("for Afni", 46, True, WHITE)]], line_spacing=1.0, space_after=2)
textbox(s, Inches(0.92), Inches(4.35), Inches(11), Inches(0.6),
        [[("Industrializing multi-agent Generative AI — Voice AI for contact centers "
           "and AI-driven HR recruitment on one governed platform.", 16, False, RGBColor(0xC7,0xD0,0xE0), True)]])
rect(s, Inches(0.9), Inches(5.55), Inches(11.5), Pt(1.2), INDIGO)
textbox(s, Inches(0.9), Inches(5.75), Inches(11), Inches(1.0),
        [[("Prepared by Evoke Technologies · GenAI Architecture Practice", 13, True, WHITE)],
         [("Shyam — Senior GenAI Architect, embedded at Afni", 12, False, RGBColor(0x9F,0xB0,0xCC))],
         [("Draft v1.0 · Confidential — for Afni & Evoke review", 11, False, RGBColor(0x9F,0xB0,0xCC))]],
        space_after=3)

# ============================================================ 2. AGENDA
s = new_slide(); content_header(s, "Orientation", "What this proposal covers")
footer(s, pg())
left = [("Why now — the BPO inflection point", 0, True),
        ("The vision & the platform-first strategy", 0, True),
        ("Platform reference architecture", 0, True),
        ("Multi-agent systems — one reusable pattern", 0, True),
        ("Flagship use case 1 — Voice AI", 0, True),
        ("Flagship use case 2 — HR recruitment", 0, True)]
right = [("LLMOps lifecycle & toolchain", 0, True),
         ("Responsible AI, security & compliance", 0, True),
         ("Observability & FinOps", 0, True),
         ("Operating model & Center of Excellence", 0, True),
         ("Roadmap: Crawl → Walk → Run", 0, True),
         ("Business case, risks & next steps", 0, True)]
card(s, Inches(0.6), Inches(1.5), Inches(6.0), Inches(5.3), "Foundations & strategy",
     left, accent=TEAL, body_size=14)
card(s, Inches(6.75), Inches(1.5), Inches(6.0), Inches(5.3), "Delivery & governance",
     right, accent=AMBER, body_size=14)

# ============================================================ 3. CONTEXT
s = new_slide(); content_header(s, "Afni business context", "Who Afni is — and why GenAI matters now")
footer(s, pg())
bullets(s, Inches(0.6), Inches(1.45), Inches(6.6), Inches(4.2), [
    ("Global BPO & customer-engagement leader, founded 1936, HQ Bloomington, IL.", 0),
    ("~3,400+ professionals across US, Mexico, the Philippines, and Afni@Home.", 0),
    ("Five service lines: Acquisition & Growth, Care & Retention, Collections, "
     "P&C Insurance (incl. subrogation), under a partnership / Gainshare model.", 0),
    ("Serves insurance, financial services, telecom, healthcare, fitness, media.", 0),
    ("Voice is the highest-volume, highest-cost channel — and hiring at scale is "
     "constant. Both are prime GenAI targets.", 0, True, NAVY),
], size=14, gap=11)
# stat panel
sx = Inches(7.5)
rect(s, sx, Inches(1.45), Inches(5.2), Inches(4.35), LIGHT)
rect(s, sx, Inches(1.45), Inches(5.2), Inches(0.5), INDIGO)
textbox(s, sx+Inches(0.2), Inches(1.5), Inches(4.8), Inches(0.4),
        [[("THE STRATEGIC PRESSURE ON BPOs", 12, True, WHITE)]], anchor=MSO_ANCHOR.MIDDLE)
stats = [("Margin & labor-cost pressure", "AI-native competitors reset price/quality expectations"),
         ("Agent attrition & ramp", "High turnover keeps hiring and training in permanent motion"),
         ("Client CX expectations", "24/7, instant, consistent, compliant interactions"),
         ("AI disruption is here", "Move first as an AI-enabled partner — or be repriced")]
yy = Inches(2.1)
for t, d in stats:
    textbox(s, sx+Inches(0.2), yy, Inches(4.85), Inches(0.4), [[(t, 13, True, TEAL)]])
    textbox(s, sx+Inches(0.2), yy+Inches(0.32), Inches(4.85), Inches(0.5), [[(d, 11.5, False, GRAY)]])
    yy += Inches(0.92)

# ============================================================ 4. VISION
s = new_slide(); content_header(s, "The vision", "One governed platform for fleets of AI agents")
footer(s, pg())
rect(s, Inches(0.6), Inches(1.45), Inches(12.1), Inches(1.35), INDIGO)
textbox(s, Inches(1.0), Inches(1.6), Inches(11.4), Inches(1.1),
        [[("“Give Afni one secure, governed Azure platform to build, evaluate, deploy, "
           "govern, and continuously improve fleets of cooperating AI agents — with the same "
           "operational rigor Afni already applies to running contact centers.”", 17, True, WHITE, True)]],
        anchor=MSO_ANCHOR.MIDDLE)
cards = [("Platform-first", "Build foundations once; every use case reuses them", TEAL),
         ("Multi-agent", "Orchestrator + specialists — the same pattern everywhere", CYAN),
         ("Governed by design", "Responsible AI, compliance & human-in-the-loop built in", AMBER),
         ("Measurable value", "Quality, cost, and business KPIs instrumented end to end", GREEN)]
cw = Inches(2.92); gap = Inches(0.14); x = Inches(0.6)
for t, d, c in cards:
    card(s, x, Inches(3.1), cw, Inches(3.5),
         t, [(d, 0)], accent=c, body_size=13)
    x += cw + gap

# ============================================================ 5. WHY LLMOps
s = new_slide(); content_header(s, "Why LLMOps", "Pilots are easy. Production at scale is the hard part.")
footer(s, pg())
card(s, Inches(0.6), Inches(1.5), Inches(6.0), Inches(5.25),
     "Without LLMOps — the pilot trap",
     [("Demos that never reach production ('POC purgatory')", 0),
      ("No repeatable evaluation — quality is a matter of opinion", 0),
      ("Prompt & model changes ship untested; regressions leak to customers", 0),
      ("Runaway token spend with no cost attribution", 0),
      ("Hallucinations & compliance gaps discovered by customers, not tests", 0),
      ("Every use case rebuilds the same plumbing from scratch", 0)],
     accent=RGBColor(0xC0,0x39,0x2B), body_size=13.5)
card(s, Inches(6.75), Inches(1.5), Inches(6.0), Inches(5.25),
     "With enterprise LLMOps",
     [("Versioned prompts & agents with automated evaluation gates", 0),
      ("Canary / blue-green deploys with instant rollback", 0),
      ("Full tracing: quality, groundedness, latency, drift, token cost", 0),
      ("Guardrails & human-in-the-loop enforced as policy, not hope", 0),
      ("Cost showback per use case and business unit", 0),
      ("A reusable platform — new use cases in weeks, not quarters", 0)],
     accent=GREEN, body_size=13.5)

# ============================================================ 6. SECTION DIVIDER - PLATFORM
def divider(number, title, subtitle):
    s = new_slide()
    rect(s, 0, 0, SW, SH, NAVY)
    rect(s, 0, Inches(3.0), SW, Inches(1.5), INDIGO)
    rect(s, Inches(0.9), Inches(3.0), Inches(0.14), Inches(1.5), AMBER)
    textbox(s, Inches(1.3), Inches(2.15), Inches(10), Inches(0.7),
            [[(number, 60, True, RGBColor(0x2A,0x40,0x6E))]])
    textbox(s, Inches(1.3), Inches(3.2), Inches(11), Inches(0.8),
            [[(title, 34, True, WHITE)]], anchor=MSO_ANCHOR.MIDDLE)
    textbox(s, Inches(1.32), Inches(4.6), Inches(10.5), Inches(0.6),
            [[(subtitle, 15, False, CYAN, True)]])
    return s

divider("01", "Platform Architecture", "The reference stack every Afni use case is built on")

# ============================================================ 7. ARCHITECTURE
s = new_slide(); content_header(s, "Reference architecture", "A layered platform on Azure AI Foundry")
footer(s, pg())
layers = [
    ("Experience & channels", "Voice / telephony (CCaaS + Azure Communication Services) · Web & chat · Agent-assist desktop · ATS / HR portals", CYAN),
    ("Orchestration & agents", "Azure AI Agent Service · Semantic Kernel / AutoGen (Microsoft Agent Framework) · orchestrator + specialist agents · tool calling", TEAL),
    ("Models & AI services", "Azure OpenAI GPT-4o / mini · gpt-realtime (speech-to-speech) · Azure AI Speech (STT/TTS) · embeddings · Content Safety", INDIGO),
    ("Knowledge & RAG", "Azure AI Search (hybrid + semantic ranker) · AI Document Intelligence · vector store · grounding & citations", RGBColor(0x3E,0x5C,0x99)),
    ("Data & integration", "Cosmos DB (agent state/memory) · Data Lake / Microsoft Fabric · Azure SQL · API Management · Functions · secure connectors to CRM / HRIS / billing", RGBColor(0x4A,0x6A,0xA8)),
    ("Platform, DevOps & LLMOps", "Prompt flow & registry · Azure AI evaluation SDK · GitHub Actions / Azure DevOps CI-CD · Container Apps / AKS", GRAY),
    ("Security & governance", "Microsoft Entra ID · Key Vault · Defender for Cloud · Microsoft Purview · private endpoints / VNet · Responsible AI", NAVY),
    ("Observability & FinOps", "Azure Monitor · Application Insights · OpenTelemetry (GenAI conventions) · token metering & cost showback", RGBColor(0x2E,0x7D,0x7D)),
]
y = Inches(1.42); lh = Inches(0.66); x = Inches(0.6); w = Inches(12.1)
for name, desc, col in layers:
    rect(s, x, y, w, lh - Inches(0.06), LIGHT)
    rect(s, x, y, Inches(3.0), lh - Inches(0.06), col)
    textbox(s, x+Inches(0.18), y, Inches(2.8), lh-Inches(0.06),
            [[(name, 12.5, True, WHITE)]], anchor=MSO_ANCHOR.MIDDLE)
    textbox(s, x+Inches(3.2), y, Inches(8.7), lh-Inches(0.06),
            [[(desc, 11, False, DARKTX)]], anchor=MSO_ANCHOR.MIDDLE)
    y += lh

# ============================================================ 8. SECTION - MULTI-AGENT
divider("02", "Multi-Agent Systems", "One orchestration pattern, reused across every use case")

# ============================================================ 9. MULTI-AGENT PATTERN
s = new_slide(); content_header(s, "The pattern", "Orchestrator routes work to specialist agents")
footer(s, pg())
# orchestrator box
rect(s, Inches(5.1), Inches(1.5), Inches(3.1), Inches(0.95), NAVY)
textbox(s, Inches(5.1), Inches(1.5), Inches(3.1), Inches(0.95),
        [[("ORCHESTRATOR", 13, True, WHITE)], [("supervisor · routing · policy", 10, False, CYAN)]],
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, space_after=2)
specialists = [
    ("Intent / Router", "classify & route", TEAL),
    ("Knowledge / RAG", "grounded answers", CYAN),
    ("Action / Tooling", "CRM · HRIS · billing", INDIGO),
    ("Compliance", "disclosures · PII · TCPA", AMBER),
    ("Sentiment", "emotion & escalation cues", GREEN),
    ("Escalation / Handoff", "warm transfer to human", RGBColor(0x8E,0x44,0xAD)),
    ("Summarize / QA", "disposition · scoring", RGBColor(0x2E,0x7D,0x7D)),
]
cols = 4; cw = Inches(2.9); gap = Inches(0.15)
startx = Inches(0.6); y0 = Inches(3.15)
for i, (t, d, c) in enumerate(specialists):
    r = i // cols; cc = i % cols
    x = startx + cc * (cw + gap)
    y = y0 + r * Inches(1.55)
    card(s, x, y, cw, Inches(1.35), t, [(d, 0)], accent=c, title_size=12.5, body_size=11)
textbox(s, Inches(0.6), Inches(6.35), Inches(12.1), Inches(0.55),
        [[("Patterns used: ", 12, True, NAVY),
          ("supervisor-orchestrator · sequential · concurrent · hand-off · group-chat · "
           "reflection/critic · human-in-the-loop. Deterministic guardrails wrap probabilistic agents.",
           12, False, GRAY)]])

# ============================================================ 10. ONE PATTERN TWO USE CASES
s = new_slide(); content_header(s, "Reuse in action", "The same agents serve voice and HR")
footer(s, pg())
table(s, Inches(0.6), Inches(1.55), Inches(12.1),
      ["Specialist agent", "Voice AI (contact center)", "HR recruitment"],
      [["Intent / Router", "Why is the customer calling?", "Where is the candidate in the funnel?"],
       ["Knowledge / RAG", "Policy, plan, billing answers", "Role, benefits, process FAQs"],
       ["Action / Tooling", "Update CRM, take payment, verify", "Read ATS, parse résumé, book interview"],
       ["Compliance", "TCPA, PCI pause/mask, disclosures", "EEOC, NYC LL144, notice & consent"],
       ["Sentiment", "Detect frustration → escalate", "Gauge candidate experience"],
       ["Escalation / Handoff", "Warm transfer to live agent", "Route to human recruiter"],
       ["Summarize / QA", "Call summary, disposition, QA score", "Structured interview notes (assist)"]],
      col_widths=[Inches(2.6), Inches(4.75), Inches(4.75)], fsize=12, row_h=Inches(0.6))
textbox(s, Inches(0.6), Inches(6.7), Inches(12), Inches(0.4),
        [[("Build once, deploy twice — and the third and fourth use cases inherit the same spine.",
           13, True, NAVY, True)]])

# ============================================================ 11. SECTION - VOICE
divider("03", "Flagship Use Case 1", "Voice AI for contact centers")

# ============================================================ 12. VOICE OVERVIEW
s = new_slide(); content_header(s, "Voice AI", "Three modes — from copilot to autonomous")
footer(s, pg())
modes = [
    ("1 · Agent-assist copilot", "Lowest risk, fastest value",
     ["Live transcription for human reps", "Real-time next-best-action & knowledge", "Sentiment & compliance nudges", "Auto summary & disposition"], TEAL),
    ("2 · Autonomous voice agent", "Scoped, containable call types",
     ["FAQs, verification, appointments", "Payment reminders (collections)", "Natural sub-second speech-to-speech", "Warm handoff when needed"], INDIGO),
    ("3 · Post-call analytics & QA", "100% coverage, not samples",
     ["Every call summarized & scored", "Compliance adherence checks", "Coaching insights for agents", "Trend & driver analytics"], AMBER),
]
x = Inches(0.6); cw = Inches(3.95); gap = Inches(0.13)
for t, sub, items, c in modes:
    card(s, x, Inches(1.5), cw, Inches(5.2), t,
         [(sub, 0, True, c)] + [(it, 0) for it in items], accent=c, body_size=12.5)
    x += cw + gap

# ============================================================ 13. VOICE CALL FLOW
s = new_slide(); content_header(s, "Voice AI", "End-to-end call flow")
footer(s, pg())
flow = ["Caller", "Telephony / CCaaS", "Speech (realtime STT)", "Orchestrator",
        "Specialist agents\n+ RAG + tools", "Systems of record\n(CRM · billing)", "Response (TTS)\nor human handoff"]
n = len(flow); bx = Inches(0.55); by = Inches(1.75)
bw = Inches(1.62); bh = Inches(1.15); gap = Inches(0.13)
colors = [GRAY, INDIGO, TEAL, NAVY, TEAL, INDIGO, GREEN]
x = bx
for i, (label, c) in enumerate(zip(flow, colors)):
    rect(s, x, by, bw, bh, c)
    textbox(s, x, by, bw, bh, [[(label, 10.5, True, WHITE)]],
            align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    if i < n - 1:
        textbox(s, x+bw-Inches(0.02), by, gap+Inches(0.05), bh,
                [[("▸", 16, True, AMBER)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    x += bw + gap
# guardrail band
rect(s, bx, Inches(3.15), Inches(12.2), Inches(0.55), LIGHT)
textbox(s, bx, Inches(3.15), Inches(12.2), Inches(0.55),
        [[("Cross-cutting: ", 12, True, NAVY),
          ("Content Safety · PII redaction · TCPA/PCI guardrails · full tracing on every turn (sub-second latency budget)", 12, False, GRAY)]],
        anchor=MSO_ANCHOR.MIDDLE)
# KPIs
table(s, bx, Inches(3.95), Inches(12.2),
      ["KPI", "Baseline (today)", "Target with Voice AI *"],
      [["Call containment / deflection", "Low / manual IVR", "20–40% of eligible call types"],
       ["Average Handle Time (AHT)", "Program baseline", "15–25% reduction"],
       ["QA coverage", "2–10% sampled", "100% automated"],
       ["Agent ramp time", "Weeks", "Materially shorter with copilot"],
       ["Compliance adherence", "Sampled review", "Monitored on every call"]],
      col_widths=[Inches(4.4), Inches(3.9), Inches(3.9)], fsize=11.5, row_h=Inches(0.42))
textbox(s, bx, Inches(6.95), Inches(12), Inches(0.3),
        [[("* Illustrative ranges — replaced with Afni actuals during discovery.", 10, False, GRAY, True)]])

# ============================================================ 14. SECTION - HR
divider("04", "Flagship Use Case 2", "AI-driven HR recruitment")

# ============================================================ 15. HR FUNNEL
s = new_slide(); content_header(s, "HR recruitment", "An agent at every step of high-volume hiring")
footer(s, pg())
steps = [
    ("JD generation", "consistent, inclusive postings"),
    ("Sourcing & screening", "parse + rank vs. criteria"),
    ("Conversational screen", "chat + optional voice pre-screen"),
    ("Scheduling", "calendar / ATS automation"),
    ("Interview scoring (assist)", "structured, human-decided"),
    ("Offer & onboarding", "faster, consistent"),
]
y = Inches(1.55); bh = Inches(0.7)
for i, (t, d) in enumerate(steps):
    rect(s, Inches(0.6), y, Inches(0.55), bh, INDIGO)
    textbox(s, Inches(0.6), y, Inches(0.55), bh, [[(str(i+1), 18, True, WHITE)]],
            align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    rect(s, Inches(1.2), y, Inches(6.4), bh, LIGHT)
    textbox(s, Inches(1.4), y, Inches(3.0), bh, [[(t, 13, True, NAVY)]], anchor=MSO_ANCHOR.MIDDLE)
    textbox(s, Inches(4.3), y, Inches(3.2), bh, [[(d, 11.5, False, GRAY)]], anchor=MSO_ANCHOR.MIDDLE)
    y += bh + Inches(0.1)
# fairness panel
fx = Inches(7.9)
rect(s, fx, Inches(1.55), Inches(4.8), Inches(4.75), NAVY)
rect(s, fx, Inches(1.55), Inches(4.8), Inches(0.5), AMBER)
textbox(s, fx+Inches(0.2), Inches(1.55), Inches(4.5), Inches(0.5),
        [[("PRINCIPLE: AI ASSISTS, HUMANS DECIDE", 12, True, NAVY)]], anchor=MSO_ANCHOR.MIDDLE)
bullets(s, fx+Inches(0.25), Inches(2.2), Inches(4.35), Inches(4.0), [
    ("No autonomous candidate rejection", 0, True, WHITE),
    ("EEOC + NYC Local Law 144 bias audits", 0, False, RGBColor(0xD5,0xDE,0xEE)),
    ("Illinois AI Video Interview Act", 0, False, RGBColor(0xD5,0xDE,0xEE)),
    ("EU AI Act — high-risk employment", 0, False, RGBColor(0xD5,0xDE,0xEE)),
    ("Candidate notice, consent & explainability", 0, False, RGBColor(0xD5,0xDE,0xEE)),
    ("Continuous adverse-impact monitoring", 0, False, RGBColor(0xD5,0xDE,0xEE)),
], size=12.5, gap=12)

# ============================================================ 16. HR VALUE
s = new_slide(); content_header(s, "HR recruitment", "Why it's the ideal internal proving ground")
footer(s, pg())
bullets(s, Inches(0.6), Inches(1.5), Inches(6.0), Inches(4.5), [
    ("Afni hires at high volume, continuously, across three countries.", 0, True, NAVY),
    ("Internal use case — Afni controls data, risk, and rollout.", 0),
    ("Reuses the voice platform for pre-screen — proves reusability.", 0),
    ("Fast, visible ROI: recruiter hours, time-to-fill, candidate experience.", 0),
    ("Builds organizational muscle in Responsible AI before customer-facing scale.", 0),
], size=14, gap=12)
table(s, Inches(6.9), Inches(1.5), Inches(5.8),
      ["KPI", "Target *"],
      [["Recruiter screening effort", "30–50% reduction"],
       ["Time-to-fill", "Materially shorter"],
       ["Cost-per-hire", "Reduced"],
       ["Candidate experience (NPS)", "Improved"],
       ["Offer-accept rate", "Improved"],
       ["90-day attrition", "Reduced via better matching"]],
      col_widths=[Inches(3.7), Inches(2.1)], fsize=12, row_h=Inches(0.5))
textbox(s, Inches(6.9), Inches(5.4), Inches(5.6), Inches(0.3),
        [[("* Illustrative — validated in discovery.", 10, False, GRAY, True)]])

# ============================================================ 17. SECTION - LLMOps
divider("05", "The LLMOps Backbone", "How agents are built, evaluated, shipped & watched")

# ============================================================ 18. LIFECYCLE
s = new_slide(); content_header(s, "LLMOps lifecycle", "A continuous, governed loop")
footer(s, pg())
stages = [
    ("Curate", "data & knowledge", CYAN),
    ("Engineer", "versioned prompts & agents", TEAL),
    ("Evaluate", "offline + online + red-team", INDIGO),
    ("Ship", "CI/CD · canary · rollback", RGBColor(0x3E,0x5C,0x99)),
    ("Serve", "gateway · quotas · caching", GRAY),
    ("Observe", "quality · cost · drift", GREEN),
]
x = Inches(0.55); cw = Inches(1.95); gap = Inches(0.12); y = Inches(1.7)
for i, (t, d, c) in enumerate(stages):
    rect(s, x, y, cw, Inches(1.35), c)
    textbox(s, x, y+Inches(0.18), cw, Inches(0.5), [[(t, 15, True, WHITE)]], align=PP_ALIGN.CENTER)
    textbox(s, x, y+Inches(0.68), cw, Inches(0.6), [[(d, 10.5, False, WHITE)]], align=PP_ALIGN.CENTER)
    if i < len(stages)-1:
        textbox(s, x+cw-Inches(0.02), y, gap+Inches(0.05), Inches(1.35),
                [[("▸", 15, True, AMBER)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    x += cw + gap
rect(s, Inches(0.55), Inches(3.3), Inches(12.2), Inches(0.5), NAVY)
textbox(s, Inches(0.55), Inches(3.3), Inches(12.2), Inches(0.5),
        [[("Wrapped end-to-end by Governance & Responsible AI · feedback (thumbs, QA, incidents) flows back into datasets", 12, True, WHITE)]],
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
card(s, Inches(0.55), Inches(4.1), Inches(6.0), Inches(2.6), "Evaluation — the quality gate",
     [("Golden datasets + LLM-as-judge + human review", 0),
      ("Groundedness / faithfulness scoring for RAG", 0),
      ("Regression gates block promotion on quality drop", 0),
      ("Red-teaming & safety evals before every release", 0),
      ("Online A/B & shadow testing in production", 0)], accent=TEAL, body_size=12)
card(s, Inches(6.75), Inches(4.1), Inches(6.0), Inches(2.6), "Toolchain on Azure",
     [("Prompt flow & prompt/model registry", 0),
      ("Azure AI evaluation SDK", 0),
      ("GitHub Actions / Azure DevOps pipelines", 0),
      ("App Insights + OpenTelemetry GenAI traces", 0),
      ("APIM gateway for metering, quotas & caching", 0)], accent=INDIGO, body_size=12)

# ============================================================ 19. SECTION - GOVERNANCE
divider("06", "Trust & Control", "Responsible AI, security, compliance, observability")

# ============================================================ 20. RESPONSIBLE AI
s = new_slide(); content_header(s, "Responsible AI & governance", "Trust engineered in, not bolted on")
footer(s, pg())
pillars = [("Fairness", TEAL), ("Reliability & safety", CYAN), ("Privacy & security", INDIGO),
           ("Inclusiveness", GREEN), ("Transparency", AMBER), ("Accountability", RGBColor(0x8E,0x44,0xAD))]
x = Inches(0.6); cw = Inches(1.94); gap = Inches(0.1)
for t, c in pillars:
    rect(s, x, Inches(1.5), cw, Inches(0.85), c)
    textbox(s, x, Inches(1.5), cw, Inches(0.85), [[(t, 12, True, WHITE)]],
            align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    x += cw + gap
card(s, Inches(0.6), Inches(2.6), Inches(6.0), Inches(4.05), "Governance mechanisms",
     [("AI use-case intake with risk-tiering (low / medium / high)", 0),
      ("Human-in-the-loop mandatory for consequential decisions", 0),
      ("Model & system cards for every deployed agent", 0),
      ("Content Safety: prompt shields, groundedness, PII, protected material", 0),
      ("Audit trails, AI incident response & red-teaming", 0),
      ("AI governance board with a regular operating cadence", 0)], accent=NAVY, body_size=12.5)
card(s, Inches(6.75), Inches(2.6), Inches(6.0), Inches(4.05), "Risk tiering drives controls",
     [("High risk (e.g., hiring decisions, collections): full HITL, bias audits, legal sign-off", 0, True, RGBColor(0xC0,0x39,0x2B)),
      ("Medium (customer voice answers): guardrails + sampled human QA + monitoring", 0, True, AMBER),
      ("Low (internal drafting, summaries): standard guardrails + spot checks", 0, True, GREEN),
      ("Controls scale with consequence — not one-size-fits-all.", 0, False, GRAY)], accent=AMBER, body_size=12.5)

# ============================================================ 21. SECURITY & COMPLIANCE
s = new_slide(); content_header(s, "Security & compliance", "Enterprise-grade by construction")
footer(s, pg())
bullets(s, Inches(0.6), Inches(1.5), Inches(5.7), Inches(4.8), [
    ("Identity: Microsoft Entra ID, RBAC, managed identities", 0),
    ("Secrets in Azure Key Vault; no keys in code", 0),
    ("Network isolation: VNet + private endpoints; no public data egress", 0),
    ("Encryption at rest & in transit; customer-managed keys optional", 0),
    ("Data residency honored across US / Mexico / Philippines", 0),
    ("Defender for Cloud posture; Purview lineage & DLP", 0),
    ("Prompt-injection & data-exfiltration defenses", 0),
], size=13.5, gap=11)
table(s, Inches(6.5), Inches(1.5), Inches(6.2),
      ["Framework", "Where it applies"],
      [["PCI-DSS", "Payment capture — pause/mask"],
       ["HIPAA", "Healthcare client programs"],
       ["TCPA", "Outbound voice & consent"],
       ["SOC 2", "Platform controls"],
       ["GDPR", "PII & data-subject rights"],
       ["EEOC / NYC LL144", "Hiring fairness & bias audit"]],
      col_widths=[Inches(2.3), Inches(3.9)], fsize=12, row_h=Inches(0.52))

# ============================================================ 22. OBSERVABILITY & FINOPS
s = new_slide(); content_header(s, "Observability & FinOps", "You cannot scale what you cannot see or cost")
footer(s, pg())
card(s, Inches(0.6), Inches(1.5), Inches(6.0), Inches(5.15), "Observe every interaction",
     [("Quality & groundedness scores in production", 0),
      ("Latency (p50/p95) against sub-second voice SLOs", 0),
      ("Errors, drift & safety events with alerting", 0),
      ("Token usage per request, agent & tool call", 0),
      ("Distributed tracing via OpenTelemetry GenAI conventions", 0),
      ("Dashboards for ops, engineering & governance", 0)], accent=TEAL, body_size=13)
card(s, Inches(6.75), Inches(1.5), Inches(6.0), Inches(5.15), "Control the spend (FinOps)",
     [("Token metering & quotas at the APIM gateway", 0),
      ("Cost showback per use case & business unit", 0),
      ("Model right-sizing: GPT-4o vs mini vs open-weight", 0),
      ("Semantic caching & prompt compression", 0),
      ("Budget guardrails and anomaly alerts", 0),
      ("Cost-per-resolved-call / cost-per-screen as first-class KPIs", 0)], accent=AMBER, body_size=13)

# ============================================================ 23. SECTION - DELIVERY
divider("07", "Delivery", "Operating model, roadmap, business case & risk")

# ============================================================ 24. OPERATING MODEL
s = new_slide(); content_header(s, "Operating model", "A GenAI Center of Excellence (hub-and-spoke)")
footer(s, pg())
rect(s, Inches(4.55), Inches(1.5), Inches(4.2), Inches(0.9), NAVY)
textbox(s, Inches(4.55), Inches(1.5), Inches(4.2), Inches(0.9),
        [[("GenAI Center of Excellence", 14, True, WHITE)],
         [("platform · standards · governance · enablement", 10, False, CYAN)]],
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, space_after=2)
roles = [("Exec sponsor", TEAL), ("AI product owner", CYAN), ("GenAI architect (lead)", INDIGO),
         ("Prompt / agent engineers", GREEN), ("LLMOps / MLOps engineers", RGBColor(0x3E,0x5C,0x99)),
         ("Data engineers", GRAY), ("RAI / governance officer", AMBER), ("Security engineer", RGBColor(0x8E,0x44,0xAD))]
x = Inches(0.6); cw = Inches(2.95); gap = Inches(0.12); y = Inches(2.9)
for i, (t, c) in enumerate(roles):
    r = i // 4; cc = i % 4
    xx = x + cc * (cw + gap); yy = y + r * Inches(0.95)
    rect(s, xx, yy, cw, Inches(0.8), LIGHT)
    rect(s, xx, yy, Inches(0.1), Inches(0.8), c)
    textbox(s, xx+Inches(0.25), yy, cw-Inches(0.3), Inches(0.8), [[(t, 12.5, True, NAVY)]], anchor=MSO_ANCHOR.MIDDLE)
rect(s, Inches(0.6), Inches(4.95), Inches(12.15), Inches(1.6), INDIGO)
textbox(s, Inches(0.85), Inches(5.1), Inches(11.7), Inches(1.35),
        [[("Federated model: ", 14, True, WHITE),
          ("the CoE owns the platform, guardrails and standards; contact-center and HR "
           "'spokes' own their use cases and outcomes. SMEs from Operations, HR, Compliance "
           "and Security embed into delivery pods. A RACI governs every lifecycle activity.", 14, False, RGBColor(0xD5,0xDE,0xEE))]],
        anchor=MSO_ANCHOR.MIDDLE)

# ============================================================ 25. ROADMAP
s = new_slide(); content_header(s, "Roadmap", "Crawl → Walk → Run over ~9–12 months")
footer(s, pg())
phases = [
    ("Phase 0", "Weeks 0–4", "Foundations & Discovery",
     ["Landing zone & security baseline", "Use-case intake & success metrics", "Data access & guardrail policy"], GRAY),
    ("Phase 1 · Crawl", "Months 1–3", "Platform MVP + pilots",
     ["Agent-assist copilot (1 program)", "HR screening pilot", "Offline eval + observability baseline"], TEAL),
    ("Phase 2 · Walk", "Months 4–7", "Scale & autonomy",
     ["Autonomous voice for scoped calls", "HR voice pre-screen + scheduling", "Online eval, FinOps, CoE stood up"], INDIGO),
    ("Phase 3 · Run", "Months 8–12", "Enterprise scale",
     ["Multiple programs & geographies", "Add subrogation & QA analytics", "Full governance, DR, improvement flywheel"], GREEN),
]
x = Inches(0.55); cw = Inches(2.98); gap = Inches(0.12)
for name, when, title, items, c in phases:
    rect(s, x, Inches(1.55), cw, Inches(0.95), c)
    textbox(s, x, Inches(1.6), cw, Inches(0.5), [[(name, 14, True, WHITE)]], align=PP_ALIGN.CENTER)
    textbox(s, x, Inches(2.05), cw, Inches(0.4), [[(when, 11, True, WHITE)]], align=PP_ALIGN.CENTER)
    rect(s, x, Inches(2.55), cw, Inches(0.55), LIGHT)
    textbox(s, x, Inches(2.55), cw, Inches(0.55), [[(title, 12, True, NAVY)]],
            align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    bullets(s, x+Inches(0.15), Inches(3.25), cw-Inches(0.25), Inches(3.2),
            [(it, 0) for it in items], size=11.5, gap=8)
    x += cw + gap
textbox(s, Inches(0.55), Inches(6.55), Inches(12), Inches(0.4),
        [[("Each phase has explicit exit criteria; value is delivered from Phase 1, not deferred to the end.",
           12.5, True, NAVY, True)]])

# ============================================================ 26. BUSINESS CASE
s = new_slide(); content_header(s, "Business case", "Where the value comes from (illustrative)")
footer(s, pg())
table(s, Inches(0.6), Inches(1.55), Inches(12.1),
      ["Value lever", "Use case", "Illustrative impact *"],
      [["Call containment / deflection", "Voice AI", "20–40% of eligible calls automated"],
       ["AHT reduction (agent-assist)", "Voice AI", "15–25% shorter handle time"],
       ["QA coverage", "Voice AI", "From ~5% sampled to 100%"],
       ["Recruiter screening effort", "HR", "30–50% fewer manual hours"],
       ["Time-to-fill & cost-per-hire", "HR", "Both materially reduced"],
       ["Attrition (better matching)", "HR", "Lower 90-day attrition"]],
      col_widths=[Inches(4.4), Inches(2.2), Inches(5.5)], fsize=12, row_h=Inches(0.5))
rect(s, Inches(0.6), Inches(5.55), Inches(12.1), Inches(0.95), INDIGO)
textbox(s, Inches(0.85), Inches(5.6), Inches(11.6), Inches(0.85),
        [[("Illustrative payback: 9–15 months. ", 14, True, WHITE),
          ("All figures are placeholders to be replaced with Afni's actual volumes, rates and cost "
           "structure during Phase 0 discovery.", 13, False, RGBColor(0xD5,0xDE,0xEE), True)]],
        anchor=MSO_ANCHOR.MIDDLE)

# ============================================================ 27. RISKS
s = new_slide(); content_header(s, "Risks & mitigations", "Named early, managed deliberately")
footer(s, pg())
table(s, Inches(0.6), Inches(1.55), Inches(12.1),
      ["Risk", "Mitigation"],
      [["Hallucination / wrong answers", "RAG grounding, groundedness eval gates, HITL, guardrails"],
       ["Compliance breach (TCPA/PCI/EEOC)", "Compliance agent, risk-tiering, bias audits, legal sign-off"],
       ["Runaway token cost", "APIM metering, model right-sizing, caching, budget alerts"],
       ["Quality regression on change", "Golden-set regression gates, canary + rollback"],
       ["Low adoption / change fatigue", "Copilot-first, agent involvement, enablement, clear ROI"],
       ["Vendor / model lock-in", "Gateway abstraction, model catalog, portable orchestration"]],
      col_widths=[Inches(4.6), Inches(7.5)], fsize=11.5, row_h=Inches(0.62))

# ============================================================ 28. NEXT STEPS / CLOSE
s = new_slide()
rect(s, 0, 0, SW, SH, NAVY)
rect(s, 0, 0, Inches(0.22), SH, TEAL)
textbox(s, Inches(0.9), Inches(0.8), Inches(11), Inches(0.5),
        [[("RECOMMENDATION & NEXT STEPS", 14, True, CYAN)]])
textbox(s, Inches(0.9), Inches(1.35), Inches(11.5), Inches(0.9),
        [[("Approve a 4-week Phase 0 to lock foundations", 30, True, WHITE)]])
steps2 = [
    ("1", "Confirm the two flagship use cases and success metrics with Ops & HR leaders"),
    ("2", "Stand up the Azure landing zone, security baseline and guardrail policy"),
    ("3", "Run use-case intake & risk-tiering; secure data access"),
    ("4", "Launch the agent-assist copilot and HR-screening pilots in Phase 1"),
]
y = Inches(2.6)
for n, t in steps2:
    rect(s, Inches(0.9), y, Inches(0.6), Inches(0.6), AMBER)
    textbox(s, Inches(0.9), y, Inches(0.6), Inches(0.6), [[(n, 20, True, NAVY)]],
            align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    textbox(s, Inches(1.7), y, Inches(10.8), Inches(0.6), [[(t, 15, False, WHITE)]], anchor=MSO_ANCHOR.MIDDLE)
    y += Inches(0.82)
rect(s, Inches(0.9), Inches(6.15), Inches(11.5), Pt(1.2), INDIGO)
textbox(s, Inches(0.9), Inches(6.35), Inches(11.5), Inches(0.8),
        [[("Evoke Technologies · GenAI Architecture Practice   |   Thank you — questions & discussion welcome", 14, True, WHITE)]])

# ---------------------------------------------------------------------------
out = os.path.join(os.path.dirname(os.path.dirname(os.path.abspath(__file__))),
                   "presentation", "Afni-LLMOps-Proposal.pptx")
prs.save(out)
print(f"Saved: {out}  ({len(prs.slides.__iter__.__self__._sldIdLst)} slides)")





