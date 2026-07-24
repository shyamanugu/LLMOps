#!/usr/bin/env python3
"""Editable, diagram-led technical deck: enterprise-grade LLMOps framework on Azure,
built around reusable components and multi-agent orchestration.
Native shapes only (no images) so it stays editable in PowerPoint. Speaker notes included.
Requires: python-pptx.  Output: presentation/LLMOps-Framework-Technical.pptx
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR

NAVY=RGBColor(0x1F,0x3A,0x5F); BLUE=RGBColor(0x2F,0x5C,0x9E); TEAL=RGBColor(0x2A,0x9D,0x8F)
GRAYTX=RGBColor(0x3C,0x46,0x54); LIGHT=RGBColor(0xEE,0xF1,0xF5); MUTE=RGBColor(0x6B,0x74,0x80)
WHITE=RGBColor(0xFF,0xFF,0xFF); LINEC=RGBColor(0xC9,0xD2,0xDD); AMBER=RGBColor(0xC8,0x7B,0x1E)
SKY=RGBColor(0x3E,0x6F,0xB0); SLATE=RGBColor(0x51,0x6B,0x8A); DTEAL=RGBColor(0x2E,0x7D,0x7D)
FONT="Calibri"; SW,SH=Inches(13.333),Inches(7.5)

prs=Presentation(); prs.slide_width=SW; prs.slide_height=SH; BLANK=prs.slide_layouts[6]

def add():
    s=prs.slides.add_slide(BLANK)
    r=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,0,0,SW,SH); r.fill.solid(); r.fill.fore_color.rgb=WHITE; r.line.fill.background(); r.shadow.inherit=False
    return s
def _txt(s,x,y,w,h,runs,align=PP_ALIGN.LEFT,anchor=MSO_ANCHOR.TOP,sa=6,ls=1.05):
    tf=s.shapes.add_textbox(x,y,w,h).text_frame; tf.word_wrap=True; tf.vertical_anchor=anchor
    tf.margin_left=0; tf.margin_right=0; tf.margin_top=0; tf.margin_bottom=0
    for i,para in enumerate(runs):
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph()
        p.alignment=align; p.space_after=Pt(sa); p.space_before=Pt(0); p.line_spacing=ls
        for (t,sz,b,c,*rest) in para:
            it=rest[0] if rest else False
            r=p.add_run(); r.text=t; r.font.size=Pt(sz); r.font.bold=b; r.font.italic=it; r.font.name=FONT; r.font.color.rgb=c
    return tf
def fillrect(s,x,y,w,h,color):
    sp=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,x,y,w,h); sp.fill.solid(); sp.fill.fore_color.rgb=color
    sp.line.fill.background(); sp.shadow.inherit=False; return sp
def rrect(s,x,y,w,h,fill,line=LINEC,lw=1.0):
    sp=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,x,y,w,h); sp.fill.solid(); sp.fill.fore_color.rgb=fill
    if line is None: sp.line.fill.background()
    else: sp.line.color.rgb=line; sp.line.width=Pt(lw)
    sp.shadow.inherit=False
    try: sp.adjustments[0]=0.08
    except Exception: pass
    return sp
def settext(sp,lines,anchor=MSO_ANCHOR.MIDDLE):
    tf=sp.text_frame; tf.word_wrap=True; tf.vertical_anchor=anchor
    tf.margin_left=Inches(0.08); tf.margin_right=Inches(0.08); tf.margin_top=Inches(0.03); tf.margin_bottom=Inches(0.03)
    for i,(t,sz,b,c) in enumerate(lines):
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph()
        p.alignment=PP_ALIGN.CENTER; p.space_after=Pt(1); p.line_spacing=1.0
        r=p.add_run(); r.text=t; r.font.size=Pt(sz); r.font.bold=b; r.font.name=FONT; r.font.color.rgb=c
def box(s,x,y,w,h,head,sub=None,fill=LIGHT,headc=NAVY,hs=12.5,ss=10,subc=GRAYTX):
    sp=rrect(s,x,y,w,h,fill); lines=[(head,hs,True,headc)]
    if sub: lines.append((sub,ss,False,subc))
    settext(sp,lines); return sp
def arrow(s,x,y,w=Inches(0.34),h=Inches(0.3),color=TEAL,shape=MSO_SHAPE.RIGHT_ARROW):
    sp=s.shapes.add_shape(shape,x,y,w,h); sp.fill.solid(); sp.fill.fore_color.rgb=color; sp.line.fill.background(); sp.shadow.inherit=False; return sp
def connect(s,x1,y1,x2,y2,color=LINEC,w=1.4):
    ln=s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,x1,y1,x2,y2); ln.line.color.rgb=color; ln.line.width=Pt(w); return ln
def title(s,text,num,kicker=None):
    if kicker: _txt(s,Inches(0.7),Inches(0.34),Inches(11.9),Inches(0.3),[[(kicker.upper(),11,True,TEAL)]])
    _txt(s,Inches(0.7),Inches(0.58),Inches(11.9),Inches(0.7),[[(text,25,True,NAVY)]])
    fillrect(s,Inches(0.72),Inches(1.22),Inches(1.5),Pt(3),TEAL)
    _txt(s,SW-Inches(1.1),SH-Inches(0.44),Inches(0.7),Inches(0.3),[[(str(num),10,False,MUTE)]],align=PP_ALIGN.RIGHT)
    _txt(s,Inches(0.7),SH-Inches(0.44),Inches(9),Inches(0.3),[[("Enterprise LLMOps on Azure — technical overview",9,False,MUTE)]])
def bullets(s,items,x=Inches(0.75),y=Inches(1.55),w=Inches(11.8),h=Inches(5.2),size=16,gap=10,color=GRAYTX):
    tf=s.shapes.add_textbox(x,y,w,h).text_frame; tf.word_wrap=True; tf.margin_left=0; tf.margin_top=0
    for i,it in enumerate(items):
        text=it[0]; lvl=it[1] if len(it)>1 else 0; bold=it[2] if len(it)>2 else False; col=it[3] if len(it)>3 else color
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph()
        p.space_after=Pt(gap); p.space_before=Pt(0); p.line_spacing=1.06; p.level=lvl
        r=p.add_run(); r.text=("•  " if lvl==0 else "–  ")+text
        r.font.size=Pt(size-lvl); r.font.bold=bold; r.font.name=FONT; r.font.color.rgb=col
    return tf
def band(s,y,label,services,lblcolor,h=0.6):
    x=Inches(0.6); w=Inches(12.13); hh=Inches(h)
    rrect(s,x,y,w,hh,LIGHT)
    lab=rrect(s,x,y,Inches(3.0),hh,lblcolor,line=None); settext(lab,[(label,12.5,True,WHITE)])
    _txt(s,x+Inches(3.2),y,w-Inches(3.4),hh,[[(services,11,False,GRAYTX)]],anchor=MSO_ANCHOR.MIDDLE)
def table(s,x,y,w,headers,rows,cw,fs=11.5,hs=11.5,rh=Inches(0.5)):
    nr=len(rows)+1; t=s.shapes.add_table(nr,len(headers),x,y,w,rh*nr).table
    for i,cc in enumerate(cw): t.columns[i].width=cc
    for j,hh in enumerate(headers):
        c=t.cell(0,j); c.fill.solid(); c.fill.fore_color.rgb=NAVY; c.vertical_anchor=MSO_ANCHOR.MIDDLE
        c.margin_left=Inches(0.1); c.margin_top=Inches(0.02); c.margin_bottom=Inches(0.02)
        r=c.text_frame.paragraphs[0].add_run(); r.text=hh; r.font.size=Pt(hs); r.font.bold=True; r.font.color.rgb=WHITE; r.font.name=FONT
    for i,row in enumerate(rows):
        for j,v in enumerate(row):
            c=t.cell(i+1,j); c.fill.solid(); c.fill.fore_color.rgb=WHITE if i%2==0 else LIGHT; c.vertical_anchor=MSO_ANCHOR.MIDDLE
            c.margin_left=Inches(0.1); c.margin_top=Inches(0.02); c.margin_bottom=Inches(0.02)
            r=c.text_frame.paragraphs[0].add_run(); r.text=v; r.font.size=Pt(fs); r.font.name=FONT; r.font.color.rgb=GRAYTX
            if j==0: r.font.bold=True; r.font.color.rgb=NAVY
    return t
def catcol(s,x,header,color,items,top=1.55,cw=Inches(2.28),ih=Inches(0.56),pitch=Inches(0.62)):
    hb=rrect(s,x,Inches(top),cw,Inches(0.5),color,line=None); settext(hb,[(header,11.5,True,WHITE)])
    yy=Inches(top+0.62)
    for name,svc in items:
        b=rrect(s,x,yy,cw,ih,LIGHT); settext(b,[(name,10.5,True,NAVY),(svc,8.5,False,MUTE)])
        yy=yy+pitch
def notes(s,text): s.notes_slide.notes_text_frame.text=text

print("Building Azure LLMOps multi-agent technical deck...")

# ---- 1 TITLE ----
s=add()
fillrect(s,0,0,Inches(0.28),SH,NAVY)
_txt(s,Inches(0.9),Inches(1.8),Inches(11.5),Inches(1.4),[[("Enterprise LLMOps Framework on Azure",31,True,NAVY)]])
fillrect(s,Inches(0.94),Inches(2.7),Inches(2.0),Pt(3),TEAL)
_txt(s,Inches(0.92),Inches(3.0),Inches(11.4),Inches(1.6),
     [[("A reusable foundation for building, testing, shipping, and running multi-agent large language model (LLM) applications — on Microsoft Azure.",17,False,GRAYTX)],
      [("LLMOps = Large Language Model Operations. Multi-agent = several specialised AI agents working together, coordinated by an orchestrator.",13,False,MUTE,True)]],sa=10)
_txt(s,Inches(0.92),Inches(6.35),Inches(11),Inches(0.5),[[("Technical overview  ·  presenter working copy",12,True,NAVY)]])
notes(s,"Set expectations. This is the technical view: the reusable building blocks of an enterprise LLMOps platform, how multi-agent orchestration works on Azure, how we build the base in about a month, and what it costs to run. I'll define every term as we go.")

# ---- 2 WHAT LLMOPS IS ----
s=add(); title(s,"What LLMOps means, and why multi-agent",2,"Basics")
bullets(s,[
    ("LLM (Large Language Model): the AI model that reads and writes text, like the model behind a chat assistant.",0),
    ("LLMOps (Large Language Model Operations): the way we build, test, release, and run LLM apps safely and the same way every time. It is DevOps adapted for language models.",0),
    ("Multi-agent: instead of one model doing everything, we use several focused agents — one to understand the request, one to search knowledge, one to call systems, one to check compliance — coordinated by an orchestrator.",0,True,NAVY),
    ("Why it matters: real tasks have many steps. Splitting the work makes it more accurate, easier to test, and easier to control. An 'agent' is an LLM that can take steps and use tools, not just chat.",0),
],y=Inches(1.55),gap=13)
by=Inches(5.05); bw=Inches(2.75); bh=Inches(1.25)
box(s,Inches(0.9),by,bw,bh,"One big model","does everything, hard to test",fill=LIGHT)
arrow(s,Inches(3.75),by+bh/2-Inches(0.15))
box(s,Inches(4.2),by,Inches(8.0),bh,"Orchestrator + focused agents","understand · search · act · check — each testable, each governed",fill=RGBColor(0xE9,0xF3,0xF0),headc=TEAL)
notes(s,"One line: LLMOps is DevOps for LLM apps. Then the key point of this deck: we do this with multiple focused agents coordinated by an orchestrator, which is more accurate and easier to govern than one giant prompt.")

# ---- 3 LAYERED ARCHITECTURE (hero) ----
s=add(); title(s,"The framework, as layers",3,"Architecture")
y=1.5
band(s,Inches(y),"Apps & channels","Web app, chat, voice, and internal tools where people ask questions",SLATE); y+=0.66
band(s,Inches(y),"API gateway","Azure API Management — one secure entry point, usage limits, keys, logging",SKY); y+=0.66
band(s,Inches(y),"Orchestration & agents","Azure AI Foundry Agent Service · Microsoft Agent Framework — coordinate agents, tools, memory",BLUE); y+=0.66
band(s,Inches(y),"Models","Azure OpenAI (GPT models) · Model Router · text embeddings",NAVY); y+=0.66
band(s,Inches(y),"Knowledge / RAG","Azure AI Search · Azure AI Document Intelligence — answer from our own documents",TEAL); y+=0.66
band(s,Inches(y),"Data","Microsoft Fabric / OneLake · Azure Cosmos DB · Blob Storage",DTEAL); y+=0.7
cc=rrect(s,Inches(0.6),Inches(y),Inches(12.13),Inches(0.6),RGBColor(0x22,0x2E,0x3E),line=None)
_txt(s,Inches(0.78),Inches(y),Inches(11.8),Inches(0.6),
     [[("Across everything:  ",11.5,True,RGBColor(0x8F,0xD3,0xCC)),
       ("Security & governance (Entra ID · Key Vault · Content Safety · Purview · Defender)   ·   "
        "Observability (Azure Monitor · App Insights)   ·   CI/CD (Azure DevOps / GitHub Actions)",11,False,WHITE)]],anchor=MSO_ANCHOR.MIDDLE)
notes(s,"The map for the whole talk. Top to bottom: a request comes in, passes the gateway, the orchestration layer coordinates the agents, which use the models and the knowledge search sitting on our data. Security, monitoring and CI/CD apply across all of it. Every box names the Azure service.")

# ---- 4 REUSABLE COMPONENTS CATALOG ----
s=add(); title(s,"The reusable building blocks (build once, reuse everywhere)",4,"Component catalog")
xs=[Inches(0.6),Inches(3.0),Inches(5.4),Inches(7.8),Inches(10.2)]
catcol(s,xs[0],"Agents & orchestration",NAVY,[
    ("Orchestration engine","Foundry Agent Service"),("Agent registry","versioned agents"),
    ("Orchestration patterns","Agent Framework"),("Tool registry (MCP)","reusable connectors"),
    ("Agent-to-agent (A2A)","cross-agent calls"),("Memory & state","Azure Cosmos DB"),
    ("Durable workflows","pause / resume")])
catcol(s,xs[1],"Knowledge & data",TEAL,[
    ("RAG search","Azure AI Search"),("Document ingestion","Document Intelligence"),
    ("Embeddings & vectorize","Azure OpenAI"),("Data platform","Fabric / OneLake"),
    ("Caching layer","cuts model cost")])
catcol(s,xs[2],"Models & prompts",BLUE,[
    ("Model gateway","API Management"),("Model catalog + Router","cost / quality"),
    ("Prompt library","versioned prompts")])
catcol(s,xs[3],"Quality & safety",SLATE,[
    ("Guardrails","Content Safety"),("Evaluation service","Foundry evaluations"),
    ("Human-in-the-loop","approvals"),("Feedback flywheel","real use → tests")])
catcol(s,xs[4],"Run & govern",DTEAL,[
    ("Observability","Monitor · OpenTelemetry"),("FinOps / cost metering","budgets · showback"),
    ("CI/CD pipeline","DevOps / GitHub"),("Security & identity","Entra ID · Key Vault"),
    ("Environments / IaC","dev · test · prod")])
notes(s,"This is the completeness slide. Every one of these is a shared component we build once and every project reuses. Read down each column. Call out the agent components on the left — registry, tool registry via MCP, agent-to-agent via A2A, memory, durable workflows — that is what makes it a real multi-agent platform, not a single chatbot.")

# ---- 5 MULTI-AGENT ORCHESTRATION ARCHITECTURE ----
s=add(); title(s,"Enterprise multi-agent orchestration",5,"The core")
_txt(s,Inches(0.7),Inches(1.32),Inches(11.9),Inches(0.34),
     [[("An orchestrator coordinates focused agents. It uses shared memory and tools, and every step is guarded and logged.",13,False,GRAYTX,True)]])
# orchestrator + memory + tools
box(s,Inches(5.35),Inches(1.85),Inches(2.65),Inches(0.95),"Orchestrator","Foundry Agent Service",fill=NAVY,headc=WHITE,ss=10,subc=RGBColor(0xAE,0xC6,0xDE))
box(s,Inches(0.7),Inches(1.85),Inches(2.7),Inches(0.95),"Memory","session · user · procedural (Cosmos DB)",fill=LIGHT)
box(s,Inches(9.95),Inches(1.85),Inches(2.75),Inches(0.95),"Tools & interop","MCP tools + A2A agents",fill=LIGHT)
connect(s,Inches(3.4),Inches(2.32),Inches(5.35),Inches(2.32))
connect(s,Inches(9.95),Inches(2.32),Inches(8.0),Inches(2.32))
# specialist agents row
specs=[("Router","understand & route"),("Knowledge / RAG","find documents"),("Action / Tools","call systems (MCP)"),
       ("Compliance","policy & PII checks"),("QA / Summarize","score & wrap up")]
sw=Inches(2.3); gap=Inches(0.14); x=Inches(0.7); y=Inches(3.55)
for i,(h,sub) in enumerate(specs):
    box(s,x,y,sw,Inches(0.95),h,sub,fill=LIGHT)
    connect(s,Inches(6.67),Inches(2.8),x+sw/2,y)
    x+=sw+gap
# bottom band
bb=rrect(s,Inches(0.7),Inches(4.9),Inches(12.0),Inches(0.95),RGBColor(0x22,0x2E,0x3E),line=None)
_txt(s,Inches(0.95),Inches(4.9),Inches(11.6),Inches(0.95),
     [[("On every step:  ",12,True,RGBColor(0x8F,0xD3,0xCC)),
       ("durable execution (pause / resume / retry)   ·   guardrails (Content Safety + PII)   ·   "
        "human approval for sensitive actions   ·   full tracing (each step logged and evaluated)",11.5,False,WHITE)]],anchor=MSO_ANCHOR.MIDDLE)
_txt(s,Inches(0.7),Inches(6.05),Inches(12),Inches(0.7),
     [[("On Azure: ",12.5,True,NAVY),("the Agent Service hosts the orchestrator and keeps state; the Microsoft Agent Framework "
       "(Semantic Kernel) builds the agents; MCP connects tools; A2A lets agents from other teams cooperate.",12.5,False,GRAYTX)]])
notes(s,"This is the heart of the deck. The orchestrator is the coordinator. It pulls from shared memory, calls tools through MCP, and can hand off to other agents through A2A. The specialist agents each do one job. The dark band is the enterprise part: every step is durable, guarded, approvable, and traced. MCP = Model Context Protocol (standard way to connect tools). A2A = Agent-to-Agent (standard way for agents to talk).")

# ---- 6 ORCHESTRATION PATTERNS ----
s=add(); title(s,"How agents work together — orchestration patterns",6,"Multi-agent")
pats=[("Sequential","A -> B -> C","Steps run one after another, each builds on the last."),
      ("Concurrent","A -> [B, C] -> merge","Run agents in parallel, then combine the results."),
      ("Group chat","A <-> B <-> C","Agents discuss with each other to reach a better answer."),
      ("Handoff","A => specialist B","One agent transfers the task to a specialist."),
      ("Planner","Planner -> workers","A lead agent plans the work and assigns it to others.")]
cw=Inches(3.85); ch=Inches(1.75); gx=Inches(0.22); gy=Inches(0.25)
pos=[(Inches(0.7),Inches(1.7)),(Inches(4.77),Inches(1.7)),(Inches(8.84),Inches(1.7)),
     (Inches(2.7),Inches(3.7)),(Inches(6.77),Inches(3.7))]
for (name,flow,desc),(x,y) in zip(pats,pos):
    sp=rrect(s,x,y,cw,ch,LIGHT)
    _txt(s,x+Inches(0.2),y+Inches(0.15),cw-Inches(0.4),Inches(0.4),[[(name,14,True,NAVY)]])
    ch2=rrect(s,x+Inches(0.2),y+Inches(0.6),cw-Inches(0.4),Inches(0.5),WHITE,line=LINEC); settext(ch2,[(flow,12.5,True,TEAL)])
    _txt(s,x+Inches(0.2),y+Inches(1.18),cw-Inches(0.4),Inches(0.5),[[(desc,11.5,False,GRAYTX)]])
_txt(s,Inches(0.7),Inches(5.75),Inches(12),Inches(0.7),
     [[("All five are supported by the Microsoft Agent Framework. We pick the pattern that fits the task — and reuse it across projects.",12.5,False,GRAYTX,True)]])
notes(s,"Five ways agents cooperate. Keep it simple: sequential is a line, concurrent is parallel, group chat is a discussion, handoff is a transfer to a specialist, planner is a lead that assigns work. We choose per task. All are built into the Agent Framework, so we do not write this coordination from scratch.")

# ---- 7 WHAT MAKES IT ENTERPRISE-GRADE ----
s=add(); title(s,"What makes the orchestration enterprise-grade",7,"Multi-agent")
props=[("Durable execution","Checkpoints, pause/resume, retries — long jobs survive failures."),
       ("Memory","Session, user, and procedural memory so agents keep context."),
       ("Agent registry","Agents are versioned, discoverable, and governed like code."),
       ("Tool registry (MCP)","Reusable, least-privilege connectors to systems and data."),
       ("Agent-to-agent (A2A)","Agents from different teams can cooperate over a standard."),
       ("Guardrails on every step","Safety, policy, and personal-data checks around each agent."),
       ("Human-in-the-loop","People approve sensitive or irreversible actions."),
       ("Full tracing & evaluation","Every step is logged, measured, and can be replayed.")]
cw=Inches(5.9); ch=Inches(1.15); gx=Inches(0.23); gy=Inches(0.2); x0=Inches(0.7); y0=Inches(1.6)
for i,(h,sub) in enumerate(props):
    r=i//2; c=i%2; x=x0+c*(cw+gx); y=y0+r*(ch+gy)
    sp=rrect(s,x,y,cw,ch,LIGHT); fillrect(s,x,y,Inches(0.12),ch,TEAL)
    _txt(s,x+Inches(0.3),y+Inches(0.16),cw-Inches(0.45),Inches(0.4),[[(h,13.5,True,NAVY)]])
    _txt(s,x+Inches(0.3),y+Inches(0.58),cw-Inches(0.45),Inches(0.5),[[(sub,11.5,False,GRAYTX)]])
notes(s,"This slide answers 'why is this enterprise-grade and not a demo'. Eight properties. The ones people miss: durable execution (jobs recover from failure), the agent registry (agents managed like code), and tracing (every step replayable for audit). These come from Azure AI Foundry Agent Service and the Agent Framework, not custom code.")

# ---- 8 RAG ----
s=add(); title(s,"Knowledge / RAG on Azure",8,"Component deep-dive")
_txt(s,Inches(0.75),Inches(1.35),Inches(11.8),Inches(0.5),
     [[("RAG (Retrieval-Augmented Generation): before answering, we search our own documents and give the model the relevant pieces, so it answers from our data instead of guessing.",13,False,GRAYTX,True)]])
steps=[("Documents","policies, PDFs, records"),("Document Intelligence","read & split text"),
       ("Embeddings + AI Search","store as searchable index"),("Retrieve","find relevant pieces"),
       ("Model answers","grounded, with sources")]
x=Inches(0.7); y=Inches(2.55); bw=Inches(2.15); bh=Inches(1.4); gap=Inches(0.28)
for i,(h,sub) in enumerate(steps):
    box(s,x,y,bw,bh,h,sub,fill=LIGHT)
    if i<len(steps)-1: arrow(s,x+bw-Inches(0.01),y+bh/2-Inches(0.15))
    x+=bw+gap
_txt(s,Inches(0.75),Inches(4.45),Inches(11.8),Inches(1.6),
     [[("Azure services: ",13,True,NAVY),("Azure AI Document Intelligence reads and splits files; Azure OpenAI creates embeddings "
       "(text turned into numbers); Azure AI Search stores them and finds matches by meaning; the model answers using only those matches.",13,False,GRAYTX)],
      [("Why it matters: this keeps answers accurate and traceable to a source, instead of the model making things up.",13,False,GRAYTX,True)]],sa=10)
notes(s,"Walk the five boxes. The message: RAG grounds answers in our own data and shows sources. Embeddings just means text stored as numbers so we can search by meaning.")

# ---- 9 SECURITY LAYERS ----
s=add(); title(s,"Keeping it safe — layers of defense",9,"Security")
layers=[("Identity & network","Entra ID sign-in, private network, no public exposure"),
        ("Gateway limits","Azure API Management — keys, rate limits, quotas"),
        ("Input checks","Azure AI Content Safety — block unsafe or malicious prompts"),
        ("Least access","tools and data limited to only what is needed; human approval for sensitive actions"),
        ("Output checks","filter unsafe content, hide personal data (PII) before it leaves"),
        ("Monitor & respond","Microsoft Defender for Cloud, alerts, and an audit trail")]
y=Inches(1.5); h=Inches(0.72)
for h_,sub in layers:
    rrect(s,Inches(2.2),y,Inches(9.0),h-Inches(0.08),LIGHT)
    _txt(s,Inches(2.45),y,Inches(2.9),h-Inches(0.08),[[(h_,13,True,NAVY)]],anchor=MSO_ANCHOR.MIDDLE)
    _txt(s,Inches(5.5),y,Inches(5.5),h-Inches(0.08),[[(sub,11.5,False,GRAYTX)]],anchor=MSO_ANCHOR.MIDDLE)
    fillrect(s,Inches(1.9),y,Inches(0.14),h-Inches(0.08),TEAL); y+=h
_txt(s,Inches(1.9),y+Inches(0.04),Inches(9.3),Inches(0.5),
     [[("Rule of thumb: treat everything going in and coming out of the model as untrusted, and check both.",12.5,True,NAVY,True)]])
notes(s,"Walk down the layers from outside in. The one line to land: treat all model input and output as untrusted and check both. PII = Personally Identifiable Information.")

# ---- 10 EVALUATION & OBSERVABILITY ----
s=add(); title(s,"Making sure answers stay good, and watching them run",10,"Quality & monitoring")
box(s,Inches(0.7),Inches(1.55),Inches(5.9),Inches(4.55),"Evaluation (before release)",None,fill=LIGHT)
bullets(s,[("Test answers with example questions and expected results.",0),
    ("Score for: correct, on-topic, and answered from our data.",0),
    ("Every release is checked against a saved test set, so quality does not slip.",0),
    ("Azure: Azure AI Foundry evaluations run these checks in the pipeline.",0,True,NAVY)],
    x=Inches(0.95),y=Inches(2.1),w=Inches(5.45),h=Inches(4.0),size=13.5,gap=12)
box(s,Inches(6.9),Inches(1.55),Inches(5.83),Inches(4.55),"Observability (in production)",None,fill=LIGHT)
bullets(s,[("Watch quality, speed, errors, and cost of every request.",0),
    ("Traces show each step the agents took, for troubleshooting.",0),
    ("Alerts fire when something drifts or breaks.",0),
    ("Azure: Azure Monitor + Application Insights, using OpenTelemetry (an open tracing standard).",0,True,NAVY)],
    x=Inches(7.15),y=Inches(2.1),w=Inches(5.4),h=Inches(4.0),size=13.5,gap=12)
notes(s,"Two halves: test before release, watch in production. We do not hope the model is right, we measure it, and we block a release if quality drops. In multi-agent systems the traces are essential — they show which agent did what.")

# ---- 11 CI/CD ----
s=add(); title(s,"How changes are tested and shipped (CI/CD)",11,"Delivery")
_txt(s,Inches(0.75),Inches(1.3),Inches(11.8),Inches(0.45),
     [[("CI/CD = Continuous Integration / Continuous Delivery: an automated pipeline that tests every change and releases it safely.",13,False,GRAYTX,True)]])
steps=[("Change","edit agent/prompt"),("Build","package as code"),("Test gates","quality · safety · cost"),
       ("Release","start small (canary)"),("Production","watch & improve")]
x=Inches(0.7); y=Inches(2.35); bw=Inches(2.15); bh=Inches(1.25); gap=Inches(0.28)
for i,(h,sub) in enumerate(steps):
    fill=RGBColor(0xF6,0xE9,0xE0) if i==2 else LIGHT; hc=AMBER if i==2 else NAVY
    box(s,x,y,bw,bh,h,sub,fill=fill,headc=hc)
    if i<len(steps)-1: arrow(s,x+bw-Inches(0.01),y+bh/2-Inches(0.15))
    x+=bw+gap
_txt(s,Inches(0.7),Inches(3.95),Inches(11.9),Inches(0.4),[[("The test gates block a release if it fails any of:",13,True,AMBER)]])
g=[("Quality","correct & on-topic"),("Groundedness","from our data, with sources"),
   ("Safety","no unsafe content; red-team"),("Cost & speed","within budget & time")]
gx=Inches(0.7)
for h,sub in g:
    box(s,gx,Inches(4.4),Inches(2.92),Inches(1.0),h,sub,fill=LIGHT); gx+=Inches(3.05)
_txt(s,Inches(0.75),Inches(5.65),Inches(11.8),Inches(0.6),
     [[("Azure: ",13,True,NAVY),("Azure DevOps or GitHub Actions runs the pipeline; releases go out gradually and roll back automatically if a problem shows up.",13,False,GRAYTX)]])
notes(s,"The heart of LLMOps. Nothing ships without passing the gates. Point at the amber Test gates box, then the four gates. On Azure this runs in Azure DevOps or GitHub Actions with gradual release and automatic rollback.")

# ---- 12 LIFECYCLE ----
s=add(); title(s,"The LLMOps lifecycle",12,"How it all flows")
steps=["Prepare\ndata & prompts","Build\nagents","Test\n(evaluate)","Release\n(CI/CD)","Run & monitor","Improve"]
x=Inches(0.75); y=Inches(2.35); bw=Inches(1.85); bh=Inches(1.35); gap=Inches(0.14)
for i,st in enumerate(steps):
    head=st.split("\n")[0]; sub=st.split("\n")[1] if "\n" in st else ""
    box(s,x,y,bw,bh,head,sub,fill=LIGHT)
    if i<len(steps)-1: arrow(s,x+bw-Inches(0.02),y+bh/2-Inches(0.15),w=Inches(0.24))
    x+=bw+gap
arrow(s,Inches(0.9),Inches(4.1),w=Inches(10.8),h=Inches(0.22),color=RGBColor(0xBF,0xD6,0xD2),shape=MSO_SHAPE.LEFT_ARROW)
_txt(s,Inches(0.9),Inches(4.45),Inches(11.6),Inches(1.4),
     [[("The loop matters most: ",14,True,NAVY),("real use produces feedback — wrong answers, new questions, cost data — which flows back to improve prompts, data, agents, and tests.",14,False,GRAYTX)],
      [("This is what turns a one-time build into an operation that keeps getting better.",14,False,GRAYTX,True)]],sa=10)
notes(s,"Left to right, then it loops back. Emphasize the loop: feedback from real use improves the system. That loop is why it's 'Ops' and not just 'build once'.")

# ---- 13 BUILD PLAN ----
s=add(); title(s,"Building the base framework on Azure — about one month",13,"Plan")
table(s,Inches(0.6),Inches(1.5),Inches(12.13),
      ["Week","Focus","What we set up on Azure"],
      [["Week 1","Set up & secure","Subscription & resource groups, Entra ID access, Key Vault, private network, API Management gateway, Azure OpenAI, logging."],
       ["Week 2","Agents & knowledge","Foundry Agent Service orchestration, first specialist agents, tool registry (MCP), memory (Cosmos DB), Azure AI Search (RAG)."],
       ["Week 3","Quality & pipeline","Guardrails (Content Safety), Foundry evaluations, Azure DevOps/GitHub CI/CD, dev / test / production."],
       ["Week 4","Monitor & harden","Azure Monitor + App Insights dashboards, cost tracking, security review, docs, and a small multi-agent sample to prove it works."]],
      cw=[Inches(1.2),Inches(2.5),Inches(8.43)],fs=11.5,rh=Inches(0.98))
_txt(s,Inches(0.6),Inches(5.95),Inches(12.13),Inches(0.9),
     [[("Honest scope: ",13,True,NAVY),("one month gives a working, reusable base — orchestration, one set of agents, knowledge, guardrails, evaluation, and CI/CD. "
       "It is enough to onboard the first real project safely, and we add components over time.",13,False,GRAYTX)]])
notes(s,"Concrete Azure build plan. Week 2 is where the multi-agent core goes in: the Agent Service, first agents, tool registry via MCP, and memory. Be honest: one month is a solid base, not everything. By week 4 a small multi-agent sample runs end to end.")

# ---- 14 ONBOARDING ----
s=add(); title(s,"Using the framework for a new project",14,"Onboarding")
steps=[("Define need","and success measures"),("Pick pattern","reuse building blocks"),
       ("Add / connect agents","specific tools & data"),("Test quality","fix until it passes"),
       ("Release small","watch closely"),("Improve","from real feedback")]
x=Inches(0.7); y=Inches(2.3); bw=Inches(1.85); bh=Inches(1.35); gap=Inches(0.14)
for i,(h,sub) in enumerate(steps):
    box(s,x,y,bw,bh,h,sub,fill=LIGHT)
    if i<len(steps)-1: arrow(s,x+bw-Inches(0.02),y+bh/2-Inches(0.15),w=Inches(0.24))
    x+=bw+gap
_txt(s,Inches(0.7),Inches(4.35),Inches(11.9),Inches(1.5),
     [[("Because the base already exists, most of this is configuration — reuse existing agents, connectors, and prompts, then add only what is new.",15,True,NAVY)],
      [("That is why the first project takes the longest, and every one after it is faster and cheaper.",14,False,GRAYTX,True)]],sa=10)
notes(s,"This is what the framework is for. Most steps are configuration, not building. We reuse existing agents and connectors and add only the new bits. Keep it generic.")

# ---- 15 RUNNING COSTS ----
s=add(); title(s,"What it costs to keep running on Azure",15,"Running cost")
_txt(s,Inches(0.6),Inches(1.4),Inches(12.13),Inches(0.4),
     [[("These are ongoing costs, not a one-time build. Most go up as the system is used more.",12.5,False,GRAYTX,True)]])
table(s,Inches(0.6),Inches(1.85),Inches(12.13),
      ["Cost area","Azure service","What drives it","Size"],
      [["Model usage","Azure OpenAI","Per token; more agents & steps = more calls","Largest, variable"],
       ["Compute / hosting","Container Apps / Functions","Always-on agent services","Medium, steady"],
       ["Search index","Azure AI Search","Storage + number of lookups","Medium"],
       ["Data & memory","Cosmos DB / Fabric / Blob","Volume of data and state kept","Low–Medium"],
       ["Monitoring","Azure Monitor / App Insights","Volume of logs and traces","Low–Medium"],
       ["People","—","Small team to maintain & support","Biggest overall"]],
      cw=[Inches(2.3),Inches(3.0),Inches(4.83),Inches(2.0)],fs=10.5,rh=Inches(0.56))
_txt(s,Inches(0.6),Inches(6.2),Inches(12.13),Inches(0.7),
     [[("Two big levers: ",12.5,True,NAVY),("use the Model Router to pick a smaller, cheaper model when a task allows it, and cache repeated answers. "
       "Note: multi-agent means several model calls per request, so caching and right-sizing matter even more.",12.5,False,GRAYTX)]])
notes(s,"Slow down here — they asked about running costs. Two messages: costs are ongoing, and the two biggest are model usage (variable) and people (fixed). Important nuance: multi-agent means several model calls per request, so cost control (routing + caching) matters more. I can give a real monthly figure once we know expected volume.")

# ---- 16 TEAM & RISKS ----
s=add(); title(s,"Who runs it, and what to watch",16,"People & risk")
box(s,Inches(0.7),Inches(1.55),Inches(5.9),Inches(4.55),"Team & roles",None,fill=LIGHT)
bullets(s,[("Platform lead / architect — owns design & standards.",0),
    ("Engineers — build the framework, agents, and pipeline.",0),
    ("Data engineer — prepares and connects data.",0),
    ("Security & governance — access, privacy, compliance.",0),
    ("Support / operations — keeps it running.",0)],
    x=Inches(0.95),y=Inches(2.1),w=Inches(5.45),h=Inches(4.0),size=13,gap=11)
box(s,Inches(6.9),Inches(1.55),Inches(5.83),Inches(4.55),"Risks to watch",None,fill=LIGHT)
bullets(s,[("Wrong or made-up answers -> grounding + testing.",0),
    ("Cost surprises (many agent calls) -> limits, budgets, caching.",0),
    ("Too much automation -> keep humans in the loop.",0),
    ("Locked to one model -> the gateway lets us switch.",0),
    ("Security gaps -> treat all input/output as untrusted.",0)],
    x=Inches(7.15),y=Inches(2.1),w=Inches(5.4),h=Inches(4.0),size=13,gap=11)
notes(s,"Two short columns. A small focused team can build the base in a month and support many projects. On risks, the cost point is bigger for multi-agent because each request makes several calls — so limits and caching matter.")

# ---- 17 SUMMARY ----
s=add(); title(s,"Summary and next steps",17,"Wrap-up")
bullets(s,[("LLMOps is how we run LLM apps reliably, not just build a demo.",0),
    ("The value is a set of reusable components — build once, reuse for every project.",0,True,NAVY),
    ("Multi-agent orchestration is the core: an orchestrator coordinating focused agents, made enterprise-grade by durable execution, memory, guardrails, and tracing.",0),
    ("Every component maps to a specific Azure service, so the build is concrete.",0),
    ("Build the base in about one month; then onboard projects quickly on top.",0,True,NAVY),
    ("Plan for ongoing running costs — mainly model usage and people.",0),
    ("Next step: agree scope for the one-month base build and set up the Azure environment.",0,True,NAVY)],
    y=Inches(1.6),gap=11)
notes(s,"Recap and make one clear ask: agreement on scope for the one-month Azure base build, and access to set up the environment. Stress the two headlines — reusable components and multi-agent orchestration.")

# ---- 18 GLOSSARY ----
s=add(); title(s,"Terms in plain English (quick reference)",18,"Appendix")
gloss=[("LLM","Large Language Model — AI that reads and writes text."),
       ("LLMOps","Large Language Model Operations — building, testing, shipping, running LLM apps."),
       ("Agent","An LLM that can use tools and take steps, not just chat."),
       ("Orchestrator","The coordinator that decides which agent does what, and when."),
       ("MCP","Model Context Protocol — a standard way for agents to use tools/connectors."),
       ("A2A","Agent-to-Agent — a standard way for agents to talk to each other."),
       ("Agent registry","A managed, versioned catalog of agents, governed like code."),
       ("Durable execution","Long jobs that can pause, resume, and recover from failures."),
       ("RAG","Retrieval-Augmented Generation — the model answers using our own documents."),
       ("Embeddings","Text turned into numbers so we can search by meaning."),
       ("CI/CD","Continuous Integration / Continuous Delivery — automated testing and release."),
       ("Guardrails","Automatic safety checks on what goes in and comes out."),
       ("PII","Personally Identifiable Information — personal data like names, phone numbers."),
       ("Token","The small unit of text a model reads/writes; billing is per token."),
       ("Canary release","Releasing to a small slice first to catch problems early."),
       ("OpenTelemetry","An open standard for collecting traces and metrics.")]
colw=Inches(5.95); x0=Inches(0.7); y0=Inches(1.5); rowh=Inches(0.66)
for i,(term,defn) in enumerate(gloss):
    r=i//2; c=i%2; x=x0+c*(colw+Inches(0.25)); y=y0+r*rowh
    _txt(s,x,y,colw,rowh,[[(term+"  ",12,True,TEAL),(defn,11,False,GRAYTX)]],ls=1.0)
notes(s,"Appendix. Leave it up during questions for the acronyms — especially MCP, A2A, and agent registry, which are new to most people.")

out=os.path.join(os.path.dirname(os.path.dirname(os.path.abspath(__file__))),"presentation","LLMOps-Framework-Technical.pptx")
prs.save(out)
print(f"Saved: {out}  ({len(prs.slides._sldIdLst)} slides)")
