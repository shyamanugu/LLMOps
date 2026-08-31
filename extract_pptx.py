import sys
sys.stdout.reconfigure(encoding='utf-8')

from pptx import Presentation

pptx_path = "presentation/LLMOps-Implementation-v2.pptx"
prs = Presentation(pptx_path)

print("=" * 80)
print("LLMOps-Implementation-v2.pptx - Content Summary")
print("=" * 80)
print(f"Total Slides: {len(prs.slides)}\n")

for i, slide in enumerate(prs.slides):
    print(f"\nSLIDE {i + 1}")
    print("-" * 80)
    
    # Extract all text shapes
    text_found = False
    for shape in slide.shapes:
        if hasattr(shape, "text") and shape.text.strip():
            text_found = True
            # Clean up the text
            text = shape.text.strip()
            # Print with line breaks preserved
            for line in text.split('\n'):
                if line.strip():
                    print(f"  {line.strip()}")
    
    if not text_found:
        print("  (No text content)")
